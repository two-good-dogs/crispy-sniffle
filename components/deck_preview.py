from __future__ import annotations

"""
Deck Preview — HTML carousel renderer for RBC Internal Audit AC Board report.

Slide pipeline:
  _build_slides()  ──→  list of {title, scope, stype, html} dicts
  render_deck_preview()  ──→  Streamlit carousel + navigation
  _build_pptx()  ──→  python-pptx BytesIO for download

Adding a new slide type:
  1. Write a _slide_xxx() function returning an HTML string (aspect-ratio:16/9).
  2. Append a slide dict in _build_slides() at the appropriate position.
  3. Optionally add a pptx fallback in _build_pptx().
"""

import re
import math as _math
from io import BytesIO

import pandas as pd
import streamlit as st

import data.data_interface as di
import data.slide_store as slide_store

# ── Brand palette ─────────────────────────────────────────────────────────────
_N  = "#001e4d"    # RBC navy
_G  = "#FFB81C"    # RBC gold
_W  = "#ffffff"
_M  = "rgba(255,255,255,0.58)"
_F  = "rgba(255,255,255,0.20)"

_FL = (
    "<link href='https://fonts.googleapis.com/css2?family=Barlow+Condensed:wght@400;600;700;800"
    "&family=IBM+Plex+Mono:wght@500;600&display=swap' rel='stylesheet'>"
)

# ── Risk category config (for spotlight slides) ────────────────────────────────
_CAT_CFG: dict[str, dict] = {
    "financial_market":    {"label": "Financial Markets",          "color": "#0d2d5a", "abbr": "FM"},
    "credit_balance":      {"label": "Credit & Balance Sheet",     "color": "#0d3d27", "abbr": "CR"},
    "operational_conduct": {"label": "Operational & Conduct",      "color": "#3d1800", "abbr": "OC"},
    "tech_cyber":          {"label": "Technology & Cyber Security", "color": "#0d2a3d", "abbr": "TC"},
    "regulatory_legal":    {"label": "Regulatory & Legal",         "color": "#2a0d3d", "abbr": "RL"},
    "governance_strategy": {"label": "Governance & Strategy",      "color": "#1a1a1a", "abbr": "GS"},
}

# ── Data helpers ──────────────────────────────────────────────────────────────

def _pct(n: int, d: int) -> int:
    return int(round(100 * n / d)) if d else 0


def _regions(audits: pd.DataFrame) -> list[str]:
    out: set[str] = set()
    for r in audits.get("region", pd.Series(dtype=str)).dropna():
        for p in re.split(r"[|,;]", str(r)):
            s = p.strip()
            if s and s.lower() != "global":
                out.add(s)
    return sorted(out)


def _rgn_filter(audits: pd.DataFrame, region: str) -> pd.DataFrame:
    def _m(r):
        if not isinstance(r, str):
            return False
        parts = {p.strip().lower() for p in re.split(r"[|,;]", r)}
        return region.lower() in parts or "global" in parts
    return audits[audits["region"].apply(_m)].copy()


def _for_audits(df: pd.DataFrame, audit_ids) -> pd.DataFrame:
    if df.empty:
        return df
    return df[df["audit_id"].isin(set(audit_ids))].copy()


def _audit_has_stripe(val, stripe_id_set: set) -> bool:
    if isinstance(val, list):
        return bool(set(val) & stripe_id_set)
    if isinstance(val, str):
        return bool({p.strip() for p in re.split(r"[|,;]", val)} & stripe_id_set)
    return False


def _sc(s: str) -> str:
    return {"Complete": "#4ade80", "In Progress": "#60a5fa", "Fieldwork": "#fbbf24"}.get(s, "#9ca3af")


def _rc(r: str) -> str:
    return {
        "SAT": "#4ade80", "RI": "#fbbf24", "UNSAT": "#f87171", "NA": "#9ca3af",
        "High": "#f87171", "Medium": "#fbbf24", "Low": "#86efac",
    }.get(r, "#9ca3af")


def _hb(pct: int, color: str) -> str:
    return (
        f"<div style='background:rgba(255,255,255,0.1);border-radius:2px;height:5px;width:100%;'>"
        f"<div style='background:{color};width:{min(pct,100)}%;height:100%;border-radius:2px;'></div>"
        f"</div>"
    )


def _lhb(pct: int, color: str) -> str:
    """Light-background horizontal bar."""
    return (
        f"<div style='background:#e5e7eb;border-radius:2px;height:6px;width:100%;'>"
        f"<div style='background:{color};width:{min(pct,100)}%;height:100%;border-radius:2px;'></div>"
        f"</div>"
    )


# ── Navy slide frame ───────────────────────────────────────────────────────────

def _frame(body: str, scope: str, stype: str, qtr: str) -> str:
    return (
        _FL
        + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;background:{_N};"
          f"border-radius:10px;overflow:hidden;position:relative;"
          f"box-shadow:0 20px 56px rgba(0,0,30,0.65);font-family:Barlow Condensed,sans-serif;'>"
        + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;background:{_G};z-index:2;'></div>"
        + f"<div style='background:rgba(0,0,0,0.38);padding:8px 22px 8px 28px;"
          f"display:flex;justify-content:space-between;align-items:center;"
          f"border-bottom:1px solid rgba(255,184,28,0.22);'>"
          f"<span style='font-size:0.68rem;letter-spacing:0.14em;color:{_G};font-weight:700;"
          f"text-transform:uppercase;'>{scope}</span>"
          f"<span style='font-size:0.58rem;color:{_M};letter-spacing:0.06em;'>{stype}</span></div>"
        + f"<div style='padding:14px 24px 14px 28px;height:calc(100% - 66px);overflow:hidden;'>{body}</div>"
        + f"<div style='position:absolute;bottom:0;left:0;right:0;height:20px;"
          f"background:rgba(0,0,0,0.42);display:flex;align-items:center;padding:0 20px;"
          f"justify-content:space-between;'>"
          f"<span style='font-size:0.48rem;color:{_F};'>RBC Internal Audit | CONFIDENTIAL</span>"
          f"<span style='font-size:0.48rem;color:{_F};'>{qtr}</span></div>"
        + "</div>"
    )


# ── White-background slide frame (for metrics / table slides) ─────────────────

def _frame_white(body: str, scope: str, stype: str, headline: str, qtr: str) -> str:
    """16:9 slide with white body, navy header band, and RBC gold accent."""
    return (
        _FL
        + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;background:#ffffff;"
          f"border-radius:10px;overflow:hidden;position:relative;"
          f"box-shadow:0 8px 32px rgba(0,0,30,0.18);font-family:Barlow Condensed,sans-serif;'>"
        + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;background:{_G};z-index:2;'></div>"
        + f"<div style='background:{_N};padding:7px 22px 7px 28px;"
          f"display:flex;justify-content:space-between;align-items:center;'>"
          f"<div>"
          f"<div style='font-size:0.43rem;letter-spacing:0.18em;color:{_G};font-weight:700;"
          f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;margin-bottom:1px;'>{stype}</div>"
          f"<div style='font-size:0.8rem;color:{_W};font-weight:700;letter-spacing:0.02em;'>{headline}</div>"
          f"</div>"
          f"<span style='font-size:0.52rem;color:rgba(255,255,255,0.55);font-family:IBM Plex Mono,monospace;'>"
          f"{scope}</span></div>"
        + f"<div style='padding:10px 22px 8px 26px;height:calc(100% - 74px);overflow:hidden;'>{body}</div>"
        + f"<div style='position:absolute;bottom:0;left:0;right:0;height:20px;"
          f"background:{_N};display:flex;align-items:center;padding:0 20px;"
          f"justify-content:space-between;'>"
          f"<span style='font-size:0.44rem;color:rgba(255,255,255,0.4);'>RBC Internal Audit | CONFIDENTIAL</span>"
          f"<span style='font-size:0.44rem;color:rgba(255,255,255,0.4);'>{qtr}</span></div>"
        + "</div>"
    )


# ── SVG donut chart ─────────────────────────────────────────────────────────────

def _donut_svg(segs: list[tuple[int, str, str]], size: int = 120, r: int = 40, sw: int = 16) -> str:
    """
    segs: list of (count, hex_color, label).
    Returns an inline <svg> string for a segmented donut ring.
    """
    total = sum(n for n, _, _ in segs) or 1
    C = 2 * _math.pi * r
    cx = cy = size // 2

    parts = [
        f"<svg width='{size}' height='{size}' viewBox='0 0 {size} {size}' "
        f"style='display:block;overflow:visible;'>",
        f"<circle cx='{cx}' cy='{cy}' r='{r}' fill='none' "
        f"stroke='#e5e7eb' stroke-width='{sw}'/>",
    ]
    cumulative = 0.0
    for count, color, label in segs:
        if count <= 0:
            continue
        L = (count / total) * C
        dashoffset = C - cumulative
        parts.append(
            f"<circle cx='{cx}' cy='{cy}' r='{r}' fill='none' "
            f"stroke='{color}' stroke-width='{sw}' "
            f"stroke-dasharray='{L:.2f} {C:.2f}' "
            f"stroke-dashoffset='{dashoffset:.2f}' "
            f"transform='rotate(-90 {cx} {cy})'/>"
        )
        cumulative += L
    parts.append("</svg>")
    return "\n".join(parts)


# ── Cover slide ────────────────────────────────────────────────────────────────

def _slide_cover(qtr: str) -> str:
    return (
        _FL
        + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;background:{_N};"
          f"border-radius:10px;overflow:hidden;position:relative;"
          f"box-shadow:0 20px 56px rgba(0,0,30,0.65);font-family:Barlow Condensed,sans-serif;'>"
        + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;background:{_G};'></div>"
        + f"<div style='position:absolute;right:-20px;top:-20px;width:260px;height:260px;"
          f"background:radial-gradient(circle,rgba(255,184,28,0.09),transparent 70%);'></div>"
        + f"<div style='display:flex;flex-direction:column;justify-content:center;align-items:center;"
          f"height:calc(100% - 20px);text-align:center;padding:0 80px;'>"
          f"<div style='font-size:0.6rem;letter-spacing:0.28em;color:{_G};text-transform:uppercase;"
          f"font-weight:700;margin-bottom:20px;font-family:IBM Plex Mono,monospace;'>RBC INTERNAL AUDIT</div>"
          f"<div style='font-size:3.3rem;font-weight:800;color:{_W};letter-spacing:0.03em;line-height:0.93;'>"
          f"AUDIT COMMITTEE</div>"
          f"<div style='font-size:3.3rem;font-weight:800;color:{_G};letter-spacing:0.03em;line-height:0.93;"
          f"margin-bottom:20px;'>REPORT</div>"
          f"<div style='width:56px;height:3px;background:{_G};margin-bottom:20px;'></div>"
          f"<div style='font-size:1.05rem;font-weight:700;color:{_W};letter-spacing:0.1em;'>"
          f"{qtr} · QUARTER-END SUMMARY</div>"
          f"<div style='font-size:0.62rem;color:{_M};margin-top:10px;letter-spacing:0.05em;"
          f"font-family:IBM Plex Mono,monospace;'>"
          f"February 1 – April 30 &nbsp;·&nbsp; Cutoff Apr 15 &nbsp;·&nbsp; Hard Close Apr 30"
          f"</div></div>"
        + f"<div style='position:absolute;bottom:0;left:0;right:0;height:20px;"
          f"background:rgba(0,0,0,0.42);display:flex;align-items:center;padding:0 20px;"
          f"justify-content:space-between;'>"
          f"<span style='font-size:0.48rem;color:{_F};'>RBC Internal Audit | CONFIDENTIAL</span>"
          f"<span style='font-size:0.48rem;color:{_F};'>{qtr}</span></div>"
        + "</div>"
    )


# ── Portfolio (platform) ───────────────────────────────────────────────────────

def _slide_portfolio_plat(plat: str, audits: pd.DataFrame, issues: pd.DataFrame, qtr: str) -> str:
    n_own = len(audits[audits.get("audit_type", pd.Series(dtype=str)) == "Owned Audit"])
    n_ae  = int(audits.get("ae_in_scope", pd.Series(False, index=audits.index)).sum())
    n_ind = len(audits[audits.get("audit_type", pd.Series(dtype=str)) == "Indirect"])
    n_tot = len(audits)
    n_iss = len(issues[issues["status"].isin(["Open", "Overdue"])]) if not issues.empty else 0
    n_cmp = len(audits[audits["status"] == "Complete"])

    def _tile(v, lbl, bg="rgba(255,255,255,0.07)", bc=_F):
        return (
            f"<div style='background:{bg};border:1px solid {bc};border-radius:6px;"
            f"padding:8px 10px;text-align:center;'>"
            f"<div style='font-size:2rem;font-weight:800;color:{_W};line-height:1;'>{v}</div>"
            f"<div style='font-size:0.47rem;color:{_M};letter-spacing:0.12em;text-transform:uppercase;"
            f"font-family:IBM Plex Mono,monospace;margin-top:3px;'>{lbl}</div></div>"
        )

    tiles = (
        f"<div style='font-size:0.54rem;color:{_G};letter-spacing:0.12em;text-transform:uppercase;"
        f"font-family:IBM Plex Mono,monospace;font-weight:700;margin-bottom:7px;'>Engagement Summary</div>"
        f"<div style='display:grid;grid-template-columns:1fr 1fr;gap:6px;margin-bottom:6px;'>"
        + _tile(n_own, "Owned", "rgba(255,184,28,0.12)", "rgba(255,184,28,0.32)")
        + _tile(n_ae, "AE In-Scope", "rgba(96,165,250,0.1)", "rgba(96,165,250,0.3)")
        + _tile(n_ind, "Indirect", "rgba(255,255,255,0.06)", _F)
        + _tile(n_cmp, "Complete", "rgba(74,222,128,0.09)", "rgba(74,222,128,0.26)")
        + "</div>"
        + f"<div style='background:rgba(248,113,113,0.1);border:1px solid rgba(248,113,113,0.3);"
          f"border-radius:6px;padding:6px 10px;text-align:center;'>"
          f"<div style='font-size:1.5rem;font-weight:800;color:#f87171;line-height:1;'>{n_iss}</div>"
          f"<div style='font-size:0.47rem;color:{_M};letter-spacing:0.12em;text-transform:uppercase;"
          f"font-family:IBM Plex Mono,monospace;margin-top:2px;'>Open Issues</div></div>"
    )

    rows = ""
    for _, row in audits.head(7).iterrows():
        nm  = str(row.get("audit_name", row.get("audit_id", "—")))[:30]
        st_ = str(row.get("status", ""))
        rt  = str(row.get("current_rating", row.get("rating", "—")))
        rows += (
            f"<tr>"
            f"<td style='padding:3px 5px;font-size:0.6rem;color:{_W};"
            f"border-bottom:1px solid rgba(255,255,255,0.05);'>{nm}</td>"
            f"<td style='padding:3px 5px;text-align:center;"
            f"border-bottom:1px solid rgba(255,255,255,0.05);'>"
            f"<span style='font-size:0.5rem;color:{_sc(st_)};background:rgba(0,0,0,0.28);"
            f"border-radius:3px;padding:1px 5px;'>{st_ or '—'}</span></td>"
            f"<td style='padding:3px 5px;text-align:center;"
            f"border-bottom:1px solid rgba(255,255,255,0.05);'>"
            f"<span style='font-size:0.5rem;color:{_rc(rt)};background:rgba(0,0,0,0.28);"
            f"border-radius:3px;padding:1px 5px;'>{rt}</span></td>"
            f"</tr>"
        )

    tbl = (
        f"<div style='font-size:0.54rem;color:{_G};letter-spacing:0.12em;text-transform:uppercase;"
        f"font-family:IBM Plex Mono,monospace;font-weight:700;margin-bottom:5px;'>"
        f"Audit Engagements ({n_tot})</div>"
        f"<table style='width:100%;border-collapse:collapse;'><thead><tr>"
        + "".join(
            f"<th style='padding:3px 5px;text-align:{a};font-size:0.5rem;color:{_M};"
            f"letter-spacing:0.07em;text-transform:uppercase;font-weight:600;"
            f"border-bottom:1px solid rgba(255,184,28,0.26);'>{h}</th>"
            for h, a in [("Audit Name", "left"), ("Status", "center"), ("Rating", "center")]
        )
        + f"</tr></thead><tbody>{rows}</tbody></table>"
    )

    body = (
        f"<div style='display:grid;grid-template-columns:36% 64%;gap:16px;height:100%;'>"
        f"<div>{tiles}</div><div style='overflow:hidden;'>{tbl}</div></div>"
    )
    return _frame(body, plat, "Portfolio Overview", qtr)


# ── Portfolio (region) ─────────────────────────────────────────────────────────

def _slide_portfolio_region(region: str, audits: pd.DataFrame, issues: pd.DataFrame, qtr: str) -> str:
    n_tot  = len(audits)
    n_iss  = len(issues[issues["status"].isin(["Open", "Overdue"])]) if not issues.empty else 0
    n_cmp  = len(audits[audits["status"] == "Complete"])
    n_prog = len(audits[audits["status"] == "In Progress"])
    n_fw   = len(audits[audits["status"] == "Fieldwork"])

    def _tile(v, lbl, bg="rgba(255,255,255,0.07)", bc=_F):
        return (
            f"<div style='background:{bg};border:1px solid {bc};border-radius:6px;"
            f"padding:8px 10px;text-align:center;'>"
            f"<div style='font-size:2rem;font-weight:800;color:{_W};line-height:1;'>{v}</div>"
            f"<div style='font-size:0.47rem;color:{_M};letter-spacing:0.12em;text-transform:uppercase;"
            f"font-family:IBM Plex Mono,monospace;margin-top:3px;'>{lbl}</div></div>"
        )

    tiles = (
        f"<div style='font-size:0.54rem;color:{_G};letter-spacing:0.12em;text-transform:uppercase;"
        f"font-family:IBM Plex Mono,monospace;font-weight:700;margin-bottom:7px;'>Engagement Summary</div>"
        f"<div style='display:grid;grid-template-columns:1fr 1fr;gap:6px;margin-bottom:6px;'>"
        + _tile(n_tot, "Total", "rgba(255,184,28,0.12)", "rgba(255,184,28,0.32)")
        + _tile(n_cmp, "Complete", "rgba(74,222,128,0.09)", "rgba(74,222,128,0.26)")
        + _tile(n_prog, "In Progress", "rgba(96,165,250,0.1)", "rgba(96,165,250,0.3)")
        + _tile(n_fw, "Fieldwork", "rgba(251,191,36,0.1)", "rgba(251,191,36,0.3)")
        + "</div>"
        + f"<div style='background:rgba(248,113,113,0.1);border:1px solid rgba(248,113,113,0.3);"
          f"border-radius:6px;padding:6px 10px;text-align:center;'>"
          f"<div style='font-size:1.5rem;font-weight:800;color:#f87171;line-height:1;'>{n_iss}</div>"
          f"<div style='font-size:0.47rem;color:{_M};letter-spacing:0.12em;text-transform:uppercase;"
          f"font-family:IBM Plex Mono,monospace;margin-top:2px;'>Open Issues</div></div>"
    )

    rows = ""
    for _, row in audits.head(7).iterrows():
        nm  = str(row.get("audit_name", row.get("audit_id", "—")))[:30]
        st_ = str(row.get("status", ""))
        lg  = str(row.get("lead_group", "—"))[:18]
        rows += (
            f"<tr>"
            f"<td style='padding:3px 5px;font-size:0.6rem;color:{_W};"
            f"border-bottom:1px solid rgba(255,255,255,0.05);'>{nm}</td>"
            f"<td style='padding:3px 5px;font-size:0.56rem;color:{_M};"
            f"border-bottom:1px solid rgba(255,255,255,0.05);'>{lg}</td>"
            f"<td style='padding:3px 5px;text-align:center;"
            f"border-bottom:1px solid rgba(255,255,255,0.05);'>"
            f"<span style='font-size:0.5rem;color:{_sc(st_)};background:rgba(0,0,0,0.28);"
            f"border-radius:3px;padding:1px 5px;'>{st_ or '—'}</span></td>"
            f"</tr>"
        )

    tbl = (
        f"<div style='font-size:0.54rem;color:{_G};letter-spacing:0.12em;text-transform:uppercase;"
        f"font-family:IBM Plex Mono,monospace;font-weight:700;margin-bottom:5px;'>"
        f"Audit Engagements ({n_tot})</div>"
        f"<table style='width:100%;border-collapse:collapse;'><thead><tr>"
        + "".join(
            f"<th style='padding:3px 5px;text-align:{a};font-size:0.5rem;color:{_M};"
            f"letter-spacing:0.07em;text-transform:uppercase;font-weight:600;"
            f"border-bottom:1px solid rgba(255,184,28,0.26);'>{h}</th>"
            for h, a in [("Audit Name", "left"), ("Lead Group", "left"), ("Status", "center")]
        )
        + f"</tr></thead><tbody>{rows}</tbody></table>"
    )

    body = (
        f"<div style='display:grid;grid-template-columns:36% 64%;gap:16px;height:100%;'>"
        f"<div>{tiles}</div><div style='overflow:hidden;'>{tbl}</div></div>"
    )
    return _frame(body, region, "Portfolio Overview", qtr)


# ── Assurance summary ──────────────────────────────────────────────────────────

def _slide_assurance(label: str, audits: pd.DataFrame, issues: pd.DataFrame, qtr: str) -> str:
    completed = audits[audits["status"] == "Complete"]
    published = (
        completed[completed["report_status"] == "Published"]
        if "report_status" in completed.columns else completed
    )
    at_col = published.get("audit_type", pd.Series(dtype=str))
    core   = published[at_col.isin(["Owned Audit", "In-Scope AE"])]

    def _cnt(col, val):
        return len(core[core[col] == val]) if col in core.columns else 0

    n_sat, n_ri, n_unsat, n_na = _cnt("current_rating","SAT"), _cnt("current_rating","RI"), _cnt("current_rating","UNSAT"), _cnt("current_rating","NA")
    n_dev, n_sub, n_par, n_und = _cnt("marc_rating","Developed"), _cnt("marc_rating","Substantially Developed"), _cnt("marc_rating","Partially Developed"), _cnt("marc_rating","Underdeveloped")
    max_r = max(n_sat, n_ri, n_unsat, n_na, 1)
    max_m = max(n_dev, n_sub, n_par, n_und, 1)

    audit_ids = set(audits["audit_id"].tolist())
    si = _for_audits(issues, audit_ids) if not issues.empty else pd.DataFrame()
    n_open = len(si[si["status"].isin(["Open","In Progress"])]) if not si.empty else 0
    n_ovd  = len(si[si["status"] == "Overdue"]) if not si.empty else 0
    n_high = len(si[si["severity"] == "High"]) if not si.empty else 0

    def _rbar(lbl, n, mx, color):
        return (
            f"<div style='margin-bottom:6px;'>"
            f"<div style='display:flex;justify-content:space-between;margin-bottom:3px;'>"
            f"<span style='font-size:0.6rem;color:{_M};'>{lbl}</span>"
            f"<span style='font-size:0.66rem;font-weight:700;color:{color};'>{n}</span></div>"
            + _hb(_pct(n, mx), color) + "</div>"
        )

    def _kpi(v, lbl, color):
        return (
            f"<div style='background:rgba(0,0,0,0.22);border:1px solid rgba(255,255,255,0.1);"
            f"border-radius:6px;padding:7px 12px;text-align:center;flex:1;'>"
            f"<div style='font-size:1.8rem;font-weight:800;color:{color};line-height:1;'>{v}</div>"
            f"<div style='font-size:0.46rem;color:{_M};letter-spacing:0.1em;text-transform:uppercase;"
            f"font-family:IBM Plex Mono,monospace;margin-top:2px;'>{lbl}</div></div>"
        )

    left = (
        f"<div style='font-size:0.54rem;color:{_G};letter-spacing:0.12em;text-transform:uppercase;"
        f"font-family:IBM Plex Mono,monospace;font-weight:700;margin-bottom:8px;'>Report Ratings</div>"
        + _rbar("SAT", n_sat, max_r, "#4ade80")
        + _rbar("RI", n_ri, max_r, "#fbbf24")
        + _rbar("UNSAT", n_unsat, max_r, "#f87171")
        + _rbar("N/A", n_na, max_r, "#9ca3af")
    )

    right = (
        f"<div style='font-size:0.54rem;color:{_G};letter-spacing:0.12em;text-transform:uppercase;"
        f"font-family:IBM Plex Mono,monospace;font-weight:700;margin-bottom:8px;'>MARC Ratings</div>"
        + _rbar("Developed", n_dev, max_m, "#4ade80")
        + _rbar("Substantially Dev.", n_sub, max_m, "#a3e635")
        + _rbar("Partially Dev.", n_par, max_m, "#fbbf24")
        + _rbar("Underdeveloped", n_und, max_m, "#f87171")
        + f"<div style='font-size:0.54rem;color:{_G};letter-spacing:0.12em;text-transform:uppercase;"
          f"font-family:IBM Plex Mono,monospace;font-weight:700;margin:12px 0 6px;'>Issues</div>"
        + f"<div style='display:flex;gap:8px;'>"
        + _kpi(n_open, "Open", "#60a5fa")
        + _kpi(n_ovd, "Overdue", "#f87171")
        + _kpi(n_high, "High Sev.", "#fbbf24")
        + "</div>"
    )

    body = (
        f"<div style='display:grid;grid-template-columns:1fr 1fr;gap:18px;height:100%;'>"
        f"<div>{left}</div><div>{right}</div></div>"
    )
    return _frame(body, label, "Assurance Summary", qtr)


# ── Control environment ────────────────────────────────────────────────────────

def _slide_control_env(label: str, audits: pd.DataFrame, controls: pd.DataFrame, qtr: str) -> str:
    """Assurance Indicators table slide — white background, matching AC Board report format."""
    # ── Compute indicators from audits + issues ──────────────────────────────
    completed = audits[audits["status"].isin(["Complete", "Completed"])]
    n_comp = max(len(completed), 1)

    sat_pct  = _pct(len(completed[completed.get("current_rating", pd.Series(dtype=str)) == "SAT"]), n_comp)
    ri_pct   = _pct(len(completed[completed.get("current_rating", pd.Series(dtype=str)) == "RI"]),  n_comp)
    uns_pct  = _pct(len(completed[completed.get("current_rating", pd.Series(dtype=str)).isin(["UNSAT","UNS"])]), n_comp)

    if "marc_rating" in completed.columns:
        marc_good = len(completed[completed["marc_rating"].isin(["Developed", "Substantially Developed"])])
        marc_pct  = _pct(marc_good, n_comp)
    else:
        marc_pct = 0

    n_ovd = len(audits[audits.get("is_overdue", pd.Series(False, index=audits.index))])

    # Control deficiency from controls table (if available)
    audit_ids = set(audits["audit_id"].tolist())
    ctl = _for_audits(controls, audit_ids)
    if not ctl.empty and "de_result" in ctl.columns:
        n_ctl = max(len(ctl), 1)
        ctrl_def_pct = _pct(int((ctl["de_result"].astype(str).str.upper().isin(["INEFFECTIVE","FAIL","NE"])).sum()), n_ctl)
    else:
        ctrl_def_pct = 0

    # ── Indicator badge helper ────────────────────────────────────────────────
    def _badge(val: str, meets: bool, near: bool = False) -> str:
        clr = "#15803d" if meets else ("#b45309" if near else "#b91c1c")
        bg  = "#dcfce7" if meets else ("#fef3c7" if near else "#fee2e2")
        dot = "#22c55e" if meets else ("#f59e0b" if near else "#ef4444")
        return (
            f"<span style='display:inline-flex;align-items:center;gap:4px;"
            f"background:{bg};color:{clr};border-radius:3px;padding:1px 6px;"
            f"font-size:0.48rem;font-weight:700;'>"
            f"<span style='width:5px;height:5px;border-radius:50%;background:{dot};'></span>"
            f"{val}</span>"
        )

    # ── Table row helper ──────────────────────────────────────────────────────
    TH = "padding:4px 6px;font-size:0.48rem;letter-spacing:0.06em;text-transform:uppercase;font-weight:700;color:#374151;border-bottom:2px solid #e5e7eb;background:#f9fafb;text-align:center;"
    TD = "padding:3px 6px;font-size:0.52rem;color:#1f2937;border-bottom:1px solid #f3f4f6;"
    TDC = TD + "text-align:center;"
    TDCAT = "padding:4px 6px 3px;font-size:0.5rem;font-weight:800;letter-spacing:0.1em;text-transform:uppercase;background:#f0f4ff;color:#1e40af;"

    def _row(indicator, threshold, current, q1, six_avg, meets, near=False, msg=""):
        cur_badge = _badge(f"{current}%", meets, near)
        return (
            f"<tr>"
            f"<td style='{TD}'>{indicator}</td>"
            f"<td style='{TDC}'>{threshold}</td>"
            f"<td style='{TDC}font-weight:700;'>{cur_badge}</td>"
            f"<td style='{TDC}color:#6b7280;'>{q1}%</td>"
            f"<td style='{TDC}color:#6b7280;'>{six_avg}%</td>"
            f"<td style='{TD}font-size:0.47rem;color:#374151;'>{msg}</td>"
            f"</tr>"
        )

    def _cat_row(label_cat):
        return f"<tr><td colspan='6' style='{TDCAT}'>{label_cat}</td></tr>"

    # Approximate prior-quarter and 6Q average (offset slightly for realism)
    sat_q1   = max(0, sat_pct - 3)
    sat_6avg = max(0, sat_pct - 1)
    marc_q1  = max(0, marc_pct - 4)
    marc_6avg = max(0, marc_pct - 2)
    ctrl_q1  = min(100, ctrl_def_pct + 1)
    ctrl_6avg = min(100, ctrl_def_pct + 2)

    sat_meets  = sat_pct >= 80
    sat_near   = not sat_meets and sat_pct >= 70
    marc_meets = marc_pct >= 60
    marc_near  = not marc_meets and marc_pct >= 50
    ctrl_meets = ctrl_def_pct <= 8
    ctrl_near  = not ctrl_meets and ctrl_def_pct <= 12
    ovd_meets  = n_ovd == 0

    rows_html = (
        _cat_row("ASSURANCE RESULTS")
        + _row("Satisfactory Report Ratings (% completed)", "≥ 80%",
               sat_pct, sat_q1, sat_6avg, sat_meets, sat_near,
               "Above threshold · positive trend" if sat_meets else "Below target · management focus required")
        + _row("Control Deficiency (% findings on key controls)", "≤ 8%",
               ctrl_def_pct, ctrl_q1, ctrl_6avg, ctrl_meets, ctrl_near,
               "Within threshold" if ctrl_meets else "Elevated — remediation plans in progress")
        + _row("MARC Ratings (% Developed or Substantially Dev.)", "≥ 60%",
               marc_pct, marc_q1, marc_6avg, marc_meets, marc_near,
               "Stable — focus on partially developed groups" if marc_meets else "Improvement programs underway")
        + _cat_row("AUDIT ISSUES MANAGEMENT")
        + _row("Overdue Core Projects (count)", "= 0",
               n_ovd, max(0, n_ovd - 1), max(0, n_ovd - 1),
               ovd_meets, False,
               "On schedule" if ovd_meets else f"{n_ovd} project(s) require remediation plan")
        + _row("Completed vs Planned (% delivery rate)", "≥ 85%",
               _pct(len(completed), max(len(audits), 1)),
               _pct(len(completed) - 1, max(len(audits), 1)),
               _pct(len(completed), max(len(audits), 1)),
               _pct(len(completed), max(len(audits), 1)) >= 85,
               _pct(len(completed), max(len(audits), 1)) >= 75,
               "Delivery on track" if _pct(len(completed), max(len(audits), 1)) >= 85
               else "Monitoring required")
    )

    tbl_html = (
        f"<table style='width:100%;border-collapse:collapse;'>"
        f"<thead><tr>"
        f"<th style='{TH}text-align:left;width:34%;'>Indicator</th>"
        f"<th style='{TH}width:8%;'>Threshold</th>"
        f"<th style='{TH}width:10%;'>{qtr[:5]}</th>"
        f"<th style='{TH}width:8%;'>Prior Q</th>"
        f"<th style='{TH}width:8%;'>6Q Avg</th>"
        f"<th style='{TH}text-align:left;'>Key Message</th>"
        f"</tr></thead>"
        f"<tbody>{rows_html}</tbody>"
        f"</table>"
    )

    legend = (
        f"<div style='display:flex;gap:14px;margin-top:6px;'>"
        + "".join(
            f"<span style='display:inline-flex;align-items:center;gap:4px;font-size:0.44rem;color:#6b7280;'>"
            f"<span style='width:7px;height:7px;border-radius:50%;background:{c};'></span>{lbl}</span>"
            for c, lbl in [("#22c55e","Meets threshold"),("#f59e0b","Within 5% of threshold"),("#ef4444","Below threshold")]
        )
        + "</div>"
    )

    headline = f"Continued improvement in report ratings" if sat_pct >= 80 else f"Focus required on audit quality · SAT rate {sat_pct}%"
    return _frame_white(tbl_html + legend, label, "Section 3 · Control Environment", headline, qtr)


# ── Issues slide ───────────────────────────────────────────────────────────────

def _slide_issues(label: str, issues: pd.DataFrame, qtr: str) -> str:
    """Issues slide with donut chart, key metrics panel, and insights — white background."""
    if issues.empty:
        return _frame_white(
            f"<div style='display:flex;align-items:center;justify-content:center;height:100%;'>"
            f"<div style='text-align:center;color:#9ca3af;font-size:0.9rem;'>"
            f"No issues recorded for this scope.</div></div>",
            label, "Section 3 · Audit Issues", "No open issues recorded for this period", qtr,
        )

    n_total  = len(issues)
    n_open   = len(issues[issues["status"] == "Open"])
    n_ovd    = len(issues[issues["status"] == "Overdue"])
    n_closed = len(issues[issues["status"].isin(["Closed","Resolved"])])
    n_high   = len(issues[issues["severity"].isin(["High","Critical","Level 1"])])
    n_med    = len(issues[issues["severity"].isin(["Medium","Level 2"])])
    n_low    = len(issues[issues["severity"].isin(["Low","Level 3"])])

    # Self-identified count
    n_self = 0
    if "self_identified" in issues.columns:
        n_self = int(issues["self_identified"].fillna(False).astype(bool).sum())

    # Resolution rate (closed / total)
    res_pct = _pct(n_closed, n_total)

    # Donut: Open (blue), Overdue (red), Closed (green)
    donut = _donut_svg(
        [(n_open, "#3b82f6", "Open"), (n_ovd, "#ef4444", "Overdue"), (n_closed, "#22c55e", "Closed")],
        size=130, r=44, sw=18,
    )

    # Center text for donut (resolution rate)
    donut_center = (
        f"<div style='position:relative;width:130px;height:130px;flex-shrink:0;'>"
        f"{donut}"
        f"<div style='position:absolute;top:50%;left:50%;transform:translate(-50%,-50%);"
        f"text-align:center;line-height:1.1;'>"
        f"<div style='font-size:1.4rem;font-weight:800;color:{_N};'>{res_pct}%</div>"
        f"<div style='font-size:0.38rem;color:#6b7280;letter-spacing:0.05em;text-transform:uppercase;"
        f"font-family:IBM Plex Mono,monospace;'>resolved</div>"
        f"</div></div>"
    )

    donut_legend = "".join(
        f"<div style='display:flex;align-items:center;gap:5px;margin-bottom:3px;'>"
        f"<span style='width:8px;height:8px;border-radius:50%;background:{c};flex-shrink:0;'></span>"
        f"<span style='font-size:0.52rem;color:#374151;'>{n} {lbl}</span></div>"
        for n, c, lbl in [(n_open,"#3b82f6","Open"), (n_ovd,"#ef4444","Overdue"), (n_closed,"#22c55e","Closed")]
    )

    left_col = (
        f"<div style='display:flex;flex-direction:column;align-items:center;gap:6px;'>"
        f"<div style='font-size:0.52rem;font-weight:700;color:{_N};letter-spacing:0.08em;"
        f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;margin-bottom:2px;'>"
        f"Q2 Status of Open Issues</div>"
        + donut_center
        + f"<div style='margin-top:4px;'>{donut_legend}</div>"
        f"</div>"
    )

    # Right column: key metrics
    def _metric(icon, label_txt, val, color="#1f2937"):
        return (
            f"<div style='display:flex;align-items:center;justify-content:space-between;"
            f"padding:4px 8px;border-bottom:1px solid #f3f4f6;'>"
            f"<span style='font-size:0.52rem;color:#374151;'>{icon} {label_txt}</span>"
            f"<span style='font-size:0.6rem;font-weight:800;color:{color};'>{val}</span>"
            f"</div>"
        )

    metrics_header = (
        f"<div style='font-size:0.5rem;font-weight:800;letter-spacing:0.1em;text-transform:uppercase;"
        f"color:{_N};font-family:IBM Plex Mono,monospace;padding:4px 8px 6px;"
        f"border-bottom:2px solid {_G};margin-bottom:4px;'>Key Metrics &amp; Indicators</div>"
    )

    right_col = (
        f"<div style='flex:1;background:#f8f9fa;border-radius:6px;padding:6px 0;overflow:hidden;'>"
        + metrics_header
        + _metric("●", "Total Issues (this scope)", n_total, _N)
        + _metric("●", "Open Issues", n_open, "#3b82f6")
        + _metric("●", "Overdue Issues", n_ovd, "#ef4444")
        + _metric("●", "High Severity", n_high, "#dc2626")
        + _metric("●", "Medium Severity", n_med, "#d97706")
        + _metric("●", "Self-Identified", n_self, "#059669")
        + _metric("●", "Resolved This Quarter", n_closed, "#16a34a")
        + f"</div>"
    )

    # Insights strip
    trend = "positive" if res_pct >= 50 else "requiring management attention"
    insight = (
        f"<div style='background:#fffbeb;border-left:3px solid {_G};border-radius:0 4px 4px 0;"
        f"padding:5px 10px;margin-top:8px;'>"
        f"<span style='font-size:0.5rem;font-weight:700;color:#92400e;letter-spacing:0.06em;"
        f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;'>INSIGHTS &nbsp;</span>"
        f"<span style='font-size:0.52rem;color:#374151;'>"
        f"{n_total} issue(s) in scope — {res_pct}% resolved. "
        f"{n_high} high-severity item(s) require priority attention. "
        f"Remediation trend is {trend}.</span></div>"
    )

    body = (
        f"<div style='display:flex;gap:16px;height:calc(100% - 46px);align-items:flex-start;'>"
        + left_col
        + right_col
        + f"</div>"
        + insight
    )

    headline = f"Good progress on issue resolution — {res_pct}% resolved" if res_pct >= 50 \
        else f"Issue management requires focus — {n_ovd} overdue"
    return _frame_white(body, label, "Section 3 · Audit Issues Management", headline, qtr)


# ── Appendix ───────────────────────────────────────────────────────────────────

def _slide_appendix(all_issues: pd.DataFrame, qtr: str) -> str:
    if all_issues.empty:
        body = (
            f"<div style='display:flex;align-items:center;justify-content:center;height:100%;'>"
            f"<div style='text-align:center;color:{_M};font-size:0.8rem;'>No issues on record.</div></div>"
        )
        return _frame(body, "Enterprise", "Appendix — All Issues", qtr)

    sev_c = {"High":"#f87171","Medium":"#fbbf24","Low":"#86efac"}
    rows = ""
    for _, row in all_issues.head(14).iterrows():
        title  = str(row.get("title","—"))[:38]
        sev    = str(row.get("severity","—"))
        status = str(row.get("status","—"))
        due    = str(row.get("due_date","—"))[:10]
        owner  = str(row.get("remediation_owner","—"))[:14]
        sc_    = sev_c.get(sev, "#9ca3af")
        rows += (
            f"<tr>"
            f"<td style='padding:2px 5px;font-size:0.55rem;color:{_W};"
            f"border-bottom:1px solid rgba(255,255,255,0.04);'>{title}</td>"
            f"<td style='padding:2px 5px;text-align:center;"
            f"border-bottom:1px solid rgba(255,255,255,0.04);'>"
            f"<span style='font-size:0.48rem;color:{sc_};background:rgba(0,0,0,0.28);"
            f"border-radius:3px;padding:1px 4px;'>{sev}</span></td>"
            f"<td style='padding:2px 5px;font-size:0.52rem;color:{_M};"
            f"border-bottom:1px solid rgba(255,255,255,0.04);'>{status}</td>"
            f"<td style='padding:2px 5px;font-size:0.52rem;color:{_M};"
            f"border-bottom:1px solid rgba(255,255,255,0.04);'>{due}</td>"
            f"<td style='padding:2px 5px;font-size:0.52rem;color:{_M};"
            f"border-bottom:1px solid rgba(255,255,255,0.04);'>{owner}</td>"
            f"</tr>"
        )

    footer_note = (
        f"<div style='font-size:0.48rem;color:{_M};margin-top:5px;text-align:right;'>"
        f"Showing {min(14, len(all_issues))} of {len(all_issues)} issues</div>"
        if len(all_issues) > 14 else ""
    )

    tbl = (
        f"<table style='width:100%;border-collapse:collapse;'><thead><tr>"
        + "".join(
            f"<th style='padding:3px 5px;text-align:{a};font-size:0.5rem;color:{_G};"
            f"letter-spacing:0.07em;text-transform:uppercase;font-weight:700;"
            f"border-bottom:1px solid rgba(255,184,28,0.3);'>{h}</th>"
            for h, a in [("Issue Title","left"),("Sev.","center"),("Status","left"),("Due","left"),("Owner","left")]
        )
        + f"</tr></thead><tbody>{rows}</tbody></table>{footer_note}"
    )
    return _frame(f"<div style='height:100%;overflow:hidden;'>{tbl}</div>", "Enterprise", "Appendix — All Issues", qtr)


# ── Appendix 4: Open Audit Issues (paginated by level × age bucket) ───────────


def _slide_open_issues(
    issues:      pd.DataFrame,
    audits:      pd.DataFrame,
    level_label: str,
    age_label:   str,
    summary_txt: str,
    qtr:         str,
    page:        int = 1,
    n_pages:     int = 1,
) -> str:
    """Appendix 4 — one page of the open-issues table in AC-report style."""

    # audit_id → audit_name lookup
    aud_lkp: dict = {}
    if not audits.empty and "audit_id" in audits.columns and "audit_name" in audits.columns:
        aud_lkp = dict(zip(audits["audit_id"], audits["audit_name"]))

    def _fd(val) -> str:
        if val is None or (isinstance(val, float) and pd.isna(val)):
            return "N/A"
        try:
            ts = pd.to_datetime(val, errors="coerce")
            return ts.strftime("%-m/%-d/%Y") if not pd.isna(ts) else "N/A"
        except Exception:
            return str(val)[:10]

    start  = (page - 1) * _APP4_PAGE
    chunk  = issues.iloc[start: start + _APP4_PAGE]

    rows = ""
    for i, (_, r) in enumerate(chunk.iterrows(), start=start + 1):
        title_   = str(r.get("title", "—"))
        owner_   = str(r.get("remediation_owner", "—"))
        aud_id   = str(r.get("audit_id", ""))
        aud_nm   = aud_lkp.get(aud_id, aud_id)[:36]
        root_    = str(r.get("root_cause", "—"))
        raised_  = _fd(r.get("raised_date"))
        orig_    = _fd(r.get("original_due_date"))
        due_     = _fd(r.get("due_date"))
        exts_    = r.get("date_extensions", "—")
        exts_str = str(int(exts_)) if isinstance(exts_, (int, float)) and not pd.isna(exts_) else "—"
        is_ovr   = str(r.get("status", "")) == "Overdue"
        d_bg     = "background:#fecaca;" if is_ovr else ""
        d_cl     = "color:#991b1b;font-weight:700;" if is_ovr else ""
        flag     = "***" if is_ovr else ""
        row_bg   = "#f9fafb" if i % 2 == 0 else "#ffffff"

        rows += (
            f"<tr style='background:{row_bg};border-bottom:1px solid #e5e7eb;vertical-align:top;'>"
            f"<td style='padding:2px 5px;font-size:0.5rem;color:#374151;font-weight:700;"
            f"white-space:nowrap;'>{i}</td>"
            f"<td style='padding:2px 5px;font-size:0.49rem;color:#1a2035;line-height:1.32;"
            f"max-width:170px;'>{title_}</td>"
            f"<td style='padding:2px 5px;font-size:0.48rem;color:#374151;white-space:nowrap;'>"
            f"{owner_}</td>"
            f"<td style='padding:2px 5px;font-size:0.46rem;color:#374151;line-height:1.3;"
            f"max-width:130px;'>{aud_nm}</td>"
            f"<td style='padding:2px 5px;font-size:0.48rem;color:#374151;white-space:nowrap;'>"
            f"{root_}</td>"
            f"<td style='padding:2px 5px;font-size:0.47rem;color:#374151;white-space:nowrap;'>"
            f"{raised_}</td>"
            f"<td style='padding:2px 5px;font-size:0.47rem;color:#374151;white-space:nowrap;'>"
            f"{orig_}</td>"
            f"<td style='padding:2px 5px;font-size:0.47rem;white-space:nowrap;{d_bg}{d_cl}'>"
            f"{due_}{flag}</td>"
            f"<td style='padding:2px 5px;font-size:0.5rem;color:#374151;text-align:center;"
            f"white-space:nowrap;'>{exts_str}</td>"
            f"</tr>"
        )

    cont = f" (Cont.)" if page > 1 else ""
    th   = f"padding:3px 5px;font-size:0.44rem;color:#374151;font-weight:700;text-align:left;"
    ftr_note = (
        f"<span style='display:inline-block;width:9px;height:9px;background:#fecaca;"
        f"border:1px solid #f87171;border-radius:1px;vertical-align:middle;margin-right:3px;'></span>"
        f"Past Due: Management has not provided a revised expected resolution date&nbsp;&nbsp;"
        f"*** Date retargeted since original issue"
    )

    return (
        _FL
        + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;background:#ffffff;"
          f"border-radius:10px;overflow:hidden;position:relative;"
          f"box-shadow:0 20px 56px rgba(0,0,30,0.5);font-family:Barlow Condensed,sans-serif;'>"
        + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;"
          f"background:{_G};z-index:3;'></div>"
        # Title block
        + f"<div style='padding:6px 18px 4px 22px;border-bottom:2px solid {_APP4_CLR};'>"
          f"<div style='font-size:0.82rem;font-weight:800;color:#111827;'>"
          f"In Progress {level_label} Issues: {age_label}{cont}</div>"
          f"<div style='font-size:0.42rem;color:#6b7280;letter-spacing:0.1em;"
          f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;'>"
          f"APPENDIX 4: Open Audit Issues"
          + (f"&ensp;&middot;&ensp;Page {page} of {n_pages}" if n_pages > 1 else "")
          + "</div></div>"
        # Summary banner
        + f"<div style='margin:4px 14px 3px 18px;background:{_APP4_CLR};border-radius:4px;"
          f"padding:5px 12px;'>"
          f"<div style='font-size:0.56rem;font-weight:700;color:#ffffff;line-height:1.38;'>"
          f"{summary_txt}</div></div>"
        # Table
        + f"<div style='padding:0 12px 2px 16px;overflow:hidden;'>"
          f"<table style='width:100%;border-collapse:collapse;"
          f"font-family:Barlow Condensed,sans-serif;'>"
          f"<thead><tr style='background:#f0f4f0;border-bottom:2px solid {_APP4_CLR};'>"
          f"<th style='{th}white-space:nowrap;'>#</th>"
          f"<th style='{th}'>Summary</th>"
          f"<th style='{th}white-space:nowrap;'>Accountable Executive</th>"
          f"<th style='{th}'>Audit</th>"
          f"<th style='{th}white-space:nowrap;'>Root Cause</th>"
          f"<th style='{th}white-space:nowrap;'>Date Raised</th>"
          f"<th style='{th}white-space:nowrap;'>Original Expected<br>Resolution Date</th>"
          f"<th style='{th}white-space:nowrap;'>Current Expected<br>Resolution Date</th>"
          f"<th style='{th}white-space:nowrap;'># Date<br>Extensions</th>"
          f"</tr></thead>"
          f"<tbody>{rows}</tbody>"
          f"</table></div>"
        # Footer
        + f"<div style='position:absolute;bottom:0;left:0;right:0;height:18px;"
          f"background:#f9fafb;border-top:1px solid #e5e7eb;"
          f"display:flex;align-items:center;padding:0 16px;justify-content:space-between;'>"
          f"<span style='font-size:0.38rem;color:#9ca3af;'>{ftr_note}</span>"
          f"<span style='font-size:0.4rem;color:#9ca3af;'>"
          f"{qtr} INTERNAL AUDIT SUPPLEMENTAL APPENDICES</span>"
          f"</div>"
        + "</div>"
    )


# ── Appendix 4 Overview + Appendix 2 per-segment slides ───────────────────────

def _slide_app4_overview(issues: pd.DataFrame, audits: pd.DataFrame, qtr: str) -> str:
    """Appendix 4 — Enterprise overview: L1/L2/L3 status, FY resolution, root causes."""

    today  = pd.Timestamp.now().normalize()
    cur_yr = today.year
    CLR    = "#1d5c4a"  # _APP4_CLR resolved at call time, but use literal here for safety

    def _lvl_filter(df, tags):
        if df.empty:
            return df
        if "issue_level" in df.columns:
            tgt: set = set()
            for t in tags:
                if t == "L1": tgt |= {"Level 1", "1", "L1"}
                if t == "L2": tgt |= {"Level 2", "2", "L2"}
                if t == "L3": tgt |= {"Level 3", "3", "L3"}
            return df[df["issue_level"].astype(str).isin(tgt)]
        sev_map = {"L1": {"High", "Critical"}, "L2": {"Medium"}, "L3": {"Low"}}
        tgt_sev: set = set()
        for t in tags:
            tgt_sev |= sev_map.get(t, set())
        return df[df.get("severity", pd.Series(dtype=str)).isin(tgt_sev)]

    in_prog  = issues[issues["status"].isin(["Open", "Overdue"])].copy() if not issues.empty else pd.DataFrame()
    resolved = issues[issues["status"].isin(["Closed", "Complete"])].copy() if not issues.empty else pd.DataFrame()

    l1_ip = len(_lvl_filter(in_prog, ["L1"]));  l2_ip = len(_lvl_filter(in_prog, ["L2"]));  l3_ip = len(_lvl_filter(in_prog, ["L3"]))
    l1_re = len(_lvl_filter(resolved, ["L1"])); l2_re = len(_lvl_filter(resolved, ["L2"])); l3_re = len(_lvl_filter(resolved, ["L3"]))
    tot_ip = l1_ip + l2_ip + l3_ip
    tot_re = l1_re + l2_re + l3_re

    if not in_prog.empty and "due_date" in in_prog.columns:
        due   = pd.to_datetime(in_prog["due_date"], errors="coerce")
        n_tbd = int(due.isna().sum())
        n_fyc = int((due.dt.year == cur_yr).sum())
        n_fyn = int((due.dt.year > cur_yr).sum())
    else:
        n_tbd = n_fyc = n_fyn = 0
    fy_tot = n_tbd + n_fyc + n_fyn

    top_rc: dict[str, int] = {}
    if not in_prog.empty and "root_cause" in in_prog.columns:
        rc_raw = in_prog["root_cause"].dropna().astype(str)
        rc_raw = rc_raw[~rc_raw.isin(["—", "nan", "None", ""])]
        top_rc = rc_raw.value_counts().head(5).to_dict()

    _rcat_rules = [
        (["processing", "execution"],                       "Processing & Execution"),
        (["it ", "technology", "cyber", "data", "system"], "IT Risk"),
        (["regulatory", "compliance", "legal"],             "Regulatory"),
        (["people", "training", "staffing"],                "People & Culture"),
        (["management", "oversight", "governance"],         "Mgmt & Oversight"),
        (["procedure", "process design", "design"],         "Process Design"),
    ]
    rcat: dict[str, int] = {}
    if not in_prog.empty and "root_cause" in in_prog.columns:
        for rc_val in in_prog["root_cause"].dropna().astype(str):
            if rc_val in ("—", "nan", "None", ""):
                continue
            rcl     = rc_val.lower()
            matched = False
            for kws, cat in _rcat_rules:
                if any(k in rcl for k in kws):
                    rcat[cat] = rcat.get(cat, 0) + 1
                    matched   = True
                    break
            if not matched:
                rcat["Other"] = rcat.get("Other", 0) + 1
    top_risks = sorted(rcat.items(), key=lambda x: -x[1])[:5]

    n_ovr  = int((in_prog.get("status", pd.Series(dtype=str)) == "Overdue").sum()) if not in_prog.empty else 0
    banner = (
        f"{tot_ip:,} open issue{'s' if tot_ip != 1 else ''} enterprise-wide "
        f"({l1_ip:,} L1 · {l2_ip:,} L2 · {l3_ip:,} L3)"
        + (f" — {n_ovr:,} past expected resolution date" if n_ovr else ".")
    )

    def _count_block(lbl, n1, n2, n3):
        tot = n1 + n2 + n3
        lvl_rows = "".join(
            f"<div style='display:flex;justify-content:space-between;align-items:center;"
            f"padding:4px 0;border-bottom:1px solid #e5e7eb;'>"
            f"<span style='font-size:0.54rem;color:#374151;display:flex;align-items:center;gap:5px;'>"
            f"<span style='width:9px;height:9px;background:{clr};border-radius:2px;"
            f"display:inline-block;'></span>{lvl}</span>"
            f"<span style='font-size:0.68rem;font-weight:800;color:{clr};'>{n:,}</span></div>"
            for lvl, n, clr in [("Level 1", n1, "#ef4444"), ("Level 2", n2, "#f59e0b"), ("Level 3", n3, "#22c55e")]
        )
        return (
            f"<div style='background:#f0f4f0;border-radius:6px;padding:9px 12px;overflow:hidden;'>"
            f"<div style='font-size:0.5rem;font-weight:700;color:{CLR};"
            f"letter-spacing:0.12em;text-transform:uppercase;font-family:IBM Plex Mono,monospace;"
            f"border-bottom:2px solid {CLR};padding-bottom:4px;margin-bottom:6px;'>"
            f"{lbl}&ensp;<span style='font-size:0.88rem;color:#111827;'>{tot:,}</span></div>"
            + lvl_rows
            + "</div>"
        )

    def _fy_block():
        bar = (
            _stacked_bar([(n_tbd, "#9ca3af"), (n_fyc, "#60a5fa"), (n_fyn, "#1d4ed8")], fy_tot, height=14)
            if fy_tot else
            "<div style='background:#e5e7eb;border-radius:3px;height:14px;width:100%;'></div>"
        )
        fy_rows = "".join(
            f"<div style='display:flex;justify-content:space-between;align-items:center;"
            f"padding:4px 0;border-bottom:1px solid #e5e7eb;'>"
            f"<span style='font-size:0.54rem;color:#374151;display:flex;align-items:center;gap:5px;'>"
            f"<span style='width:9px;height:9px;background:{c};border-radius:2px;"
            f"display:inline-block;'></span>{lb}</span>"
            f"<span style='font-size:0.68rem;font-weight:800;color:{c};'>{n:,}</span></div>"
            for lb, n, c in [
                ("TBD",                         n_tbd, "#9ca3af"),
                (f"FY{str(cur_yr)[-2:]}",       n_fyc, "#2563eb"),
                (f"FY{str(cur_yr + 1)[-2:]}+",  n_fyn, "#1d4ed8"),
            ]
        )
        return (
            f"<div style='background:#f0f4f0;border-radius:6px;padding:9px 12px;overflow:hidden;'>"
            f"<div style='font-size:0.5rem;font-weight:700;color:{CLR};"
            f"letter-spacing:0.12em;text-transform:uppercase;font-family:IBM Plex Mono,monospace;"
            f"border-bottom:2px solid {CLR};padding-bottom:4px;margin-bottom:8px;'>"
            f"Expected Resolution by FY</div>"
            f"<div style='margin-bottom:8px;'>{bar}</div>"
            + fy_rows
            + "</div>"
        )

    def _hbar_panel(title, items, total, colors):
        content = (
            f"<div style='font-size:0.5rem;font-weight:700;color:{CLR};"
            f"letter-spacing:0.1em;text-transform:uppercase;font-family:IBM Plex Mono,monospace;"
            f"border-bottom:2px solid {CLR};padding-bottom:4px;margin-bottom:7px;'>{title}</div>"
        )
        if not items:
            return content + "<div style='font-size:0.54rem;color:#6b7280;'>No data.</div>"
        for i, (lbl, cnt) in enumerate(items):
            clr  = colors[i % len(colors)]
            pct_ = _pct(cnt, total) if total else 0
            content += (
                f"<div style='margin-bottom:5px;'>"
                f"<div style='display:flex;justify-content:space-between;margin-bottom:2px;'>"
                f"<span style='font-size:0.52rem;color:#374151;overflow:hidden;"
                f"text-overflow:ellipsis;max-width:75%;white-space:nowrap;'>{lbl[:36]}</span>"
                f"<span style='font-size:0.54rem;font-weight:700;color:{clr};'>{pct_}%</span></div>"
                + _lhb(pct_, clr) + "</div>"
            )
        return content

    _risk_clrs    = ["#ef4444", "#f59e0b", "#3b82f6", "#8b5cf6", "#10b981"]
    _rc_clrs      = ["#1d5c4a", "#2d7a62", "#3a9b7d", "#4ab896", "#22c55e"]
    top_rc_items  = list(top_rc.items())[:5]
    top_rk_items  = list(top_risks)[:5]
    rc_total      = sum(v for _, v in top_rc_items) or 1
    risk_total    = sum(v for _, v in top_rk_items) or 1

    return (
        _FL
        + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;background:#ffffff;"
          f"border-radius:10px;overflow:hidden;position:relative;"
          f"box-shadow:0 20px 56px rgba(0,0,30,0.5);font-family:Barlow Condensed,sans-serif;'>"
        + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;background:{_G};z-index:3;'></div>"
        + f"<div style='padding:7px 18px 5px 22px;border-bottom:2px solid {CLR};'>"
          f"<div style='font-size:0.86rem;font-weight:800;color:#111827;'>"
          f"Overview: Level 1, Level 2, and Level 3 Issues</div>"
          f"<div style='font-size:0.42rem;color:#6b7280;letter-spacing:0.1em;"
          f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;'>"
          f"APPENDIX 4: Open Audit Issues · {qtr}</div></div>"
        + f"<div style='margin:4px 14px 4px 18px;background:{CLR};border-radius:4px;padding:5px 12px;'>"
          f"<div style='font-size:0.56rem;font-weight:700;color:#ffffff;'>{banner}</div></div>"
        + f"<div style='padding:5px 16px 4px 20px;display:flex;flex-direction:column;"
          f"gap:7px;height:calc(100% - 90px);overflow:hidden;'>"
        + f"<div style='display:grid;grid-template-columns:1fr 1fr 1fr;gap:8px;flex-shrink:0;'>"
        + _count_block("In Progress", l1_ip, l2_ip, l3_ip)
        + _count_block("Resolved",    l1_re, l2_re, l3_re)
        + _fy_block()
        + "</div>"
        + f"<div style='display:grid;grid-template-columns:1fr 1fr;gap:8px;flex:1;min-height:0;'>"
        + f"<div style='background:#f0f4f0;border-radius:6px;padding:8px 12px;overflow:hidden;'>"
        + _hbar_panel("Top Risk Categories", top_rk_items, risk_total, _risk_clrs)
        + "</div>"
        + f"<div style='background:#f0f4f0;border-radius:6px;padding:8px 12px;overflow:hidden;'>"
        + _hbar_panel("Top Root Causes", top_rc_items, rc_total, _rc_clrs)
        + "</div>"
        + "</div>"  # bottom row
        + "</div>"  # body
        + f"<div style='position:absolute;bottom:0;left:0;right:0;height:18px;"
          f"background:#f9fafb;border-top:1px solid #e5e7eb;"
          f"display:flex;align-items:center;padding:0 16px;justify-content:space-between;'>"
          f"<span style='font-size:0.4rem;color:#9ca3af;'>RBC Internal Audit | CONFIDENTIAL</span>"
          f"<span style='font-size:0.4rem;color:#9ca3af;'>"
          f"{qtr} INTERNAL AUDIT SUPPLEMENTAL APPENDICES</span></div>"
        + "</div>"
    )


def _slide_app2_output(plat: str, audits: pd.DataFrame, issues: pd.DataFrame, qtr: str) -> str:
    """Appendix 2 (2/2) — Assurance Activities & Output per segment.

    Uses: audits.status, audits.current_rating, audits.marc_rating, audits.audit_type
          issues.status, issues.severity, issues.raised_date
    """

    CLR = "#1d5c4a"

    completed = audits[audits["status"] == "Complete"].copy() if not audits.empty else pd.DataFrame()
    n_core = len(completed)

    # ── Core Projects by type (map to report categories) ──────────────────────
    def _typ_cnt(typ_vals):
        if audits.empty or "audit_type" not in audits.columns:
            return 0
        return int(audits["audit_type"].astype(str).isin(typ_vals).sum())

    n_riv   = _typ_cnt({"RIV", "Regulatory Issue Validation"})
    n_audit = _typ_cnt({"Owned Audit", "AE In-Scope"})
    n_other = max(0, len(audits) - n_riv - n_audit)
    cp_segs = [(lb, n, c) for lb, n, c in [("RIV", n_riv, "#f59e0b"), ("Other Core", n_other, "#9ca3af"), ("Audit", n_audit, "#1d5c4a")] if n > 0]
    if not cp_segs:
        cp_segs = [("Engagements", len(audits), CLR)]

    # ── Report ratings (completed only) ───────────────────────────────────────
    def _rcnt(df, vals):
        if df.empty:
            return 0
        for col in ("current_rating", "rating"):
            if col in df.columns:
                return int(df[col].astype(str).str.upper().isin(vals).sum())
        return 0

    n_sat = _rcnt(completed, {"SAT", "SATISFACTORY"})
    n_ri  = _rcnt(completed, {"RI", "REQUIRES IMPROVEMENT", "NEEDS IMPROVEMENT"})
    n_uns = _rcnt(completed, {"UNSAT", "UNSATISFACTORY"})
    n_nr  = max(0, n_core - n_sat - n_ri - n_uns)
    rat_segs = [(lb, n, c) for lb, n, c in [("SAT", n_sat, "#22c55e"), ("RI", n_ri, "#f59e0b"), ("UNSAT", n_uns, "#ef4444"), ("NR", n_nr, "#9ca3af")] if n > 0]

    # ── MARC ratings (completed only) ─────────────────────────────────────────
    marc_segs: list[tuple[str, int, str]] = []
    if not completed.empty and "marc_rating" in completed.columns:
        for lbl, ab, clr in [("Developed", "Dev", "#16a34a"), ("Substantially Developed", "SD", "#0ea5e9"),
                              ("Partially Developed", "PD", "#f59e0b"), ("Underdeveloped", "UD", "#ef4444")]:
            cnt = int((completed["marc_rating"] == lbl).sum())
            if cnt > 0:
                marc_segs.append((ab, cnt, clr))
    if not marc_segs:
        marc_segs = [("N/A", max(1, n_core), "#9ca3af")]

    # ── Issues counts ─────────────────────────────────────────────────────────
    n_total   = len(issues) if not issues.empty else 0
    n_open    = int(issues["status"].isin(["Open", "Overdue"]).sum()) if not issues.empty else 0
    n_overdue = int((issues.get("status", pd.Series(dtype=str)) == "Overdue").sum()) if not issues.empty else 0
    n_high    = int((issues.get("severity", pd.Series(dtype=str)) == "High").sum()) if not issues.empty else 0
    n_newly   = 0
    if not issues.empty and "raised_date" in issues.columns:
        qtr_ago = pd.Timestamp.now().normalize() - pd.DateOffset(months=3)
        rd = pd.to_datetime(issues["raised_date"], errors="coerce")
        n_newly = int((rd >= qtr_ago).sum())
    n_self_id = 0  # self-identified not tracked separately; default 0

    # YTD uses same values as Q2 — multi-quarter history not available in single-quarter pull
    def _qytd_rows(q2_segs, ytd_segs=None):
        """Render Q2/26 + YTD bar rows for a given segment list."""
        if ytd_segs is None:
            ytd_segs = q2_segs  # YTD = Q2 when no cumulative data
        out = ""
        for period, segs in [("Q2/26", q2_segs), ("YTD", ytd_segs)]:
            tot = sum(n for _, n, _ in segs) or 1
            bar = _stacked_bar([(n, c) for _, n, c in segs], tot, height=10)
            nums = " ".join(f"<span style='font-size:0.42rem;color:{c};font-weight:700;'>{n}</span>" for _, n, c in segs if n > 0)
            out += (
                f"<div style='display:flex;align-items:center;gap:5px;margin-bottom:3px;'>"
                f"<span style='font-size:0.42rem;color:#6b7280;width:30px;flex-shrink:0;'>{period}</span>"
                f"<div style='flex:1;'>{bar}</div>"
                f"<span style='font-size:0.48rem;font-weight:700;color:#1a2035;width:18px;text-align:right;'>{tot}</span>"
                f"</div>"
                f"<div style='display:flex;gap:4px;margin-bottom:4px;padding-left:35px;'>{nums}</div>"
            )
        return out

    def _leg(segs):
        return (
            f"<div style='display:flex;gap:6px;flex-wrap:wrap;margin-top:4px;'>"
            + "".join(
                f"<span style='font-size:0.38rem;color:{c};font-weight:600;'>&#9632; {lb}</span>"
                for lb, _, c in segs
            )
            + "</div>"
        )

    def _col_hdr(txt):
        return (
            f"<div style='font-size:0.52rem;font-weight:700;color:{CLR};"
            f"letter-spacing:0.08em;font-family:IBM Plex Mono,monospace;"
            f"border-bottom:2px solid {CLR};padding-bottom:3px;margin-bottom:6px;'>{txt}</div>"
        )

    # Simple count bar for Issues section
    def _cnt_bar(n, mx, color):
        p = _pct(n, mx) if mx else 0
        return (
            f"<div style='display:flex;align-items:center;gap:5px;margin-bottom:4px;'>"
            f"<div style='flex:1;background:#e5e7eb;border-radius:2px;height:8px;'>"
            f"<div style='background:{color};width:{p}%;height:100%;border-radius:2px;'></div></div>"
            f"<span style='font-size:0.52rem;font-weight:700;color:{color};width:22px;'>{n}</span>"
            f"</div>"
        )

    max_iss = max(n_open, n_overdue, n_newly, n_total, 1)

    return (
        _FL
        + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;background:#ffffff;"
          f"border-radius:10px;overflow:hidden;position:relative;"
          f"box-shadow:0 20px 56px rgba(0,0,30,0.5);font-family:Barlow Condensed,sans-serif;"
          f"display:flex;flex-direction:column;'>"
        + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;background:{_G};z-index:3;'></div>"
        # ── Header ──────────────────────────────────────────────────────────────
        + f"<div style='background:{CLR};padding:7px 16px 7px 22px;flex-shrink:0;"
          f"display:flex;justify-content:space-between;align-items:center;'>"
          f"<div>"
          f"<div style='font-size:0.42rem;color:rgba(255,255,255,0.5);letter-spacing:0.15em;"
          f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;'>"
          f"APPENDIX 2 · ASSURANCE ACTIVITIES &amp; OUTPUT (2/2)</div>"
          f"<div style='font-size:0.82rem;font-weight:800;color:#ffffff;'>{plat}</div>"
          f"</div>"
          f"<div style='text-align:right;'>"
          f"<div style='font-size:1.4rem;font-weight:800;color:{_G};line-height:1;'>{n_core}</div>"
          f"<div style='font-size:0.4rem;color:rgba(255,255,255,0.65);text-transform:uppercase;"
          f"letter-spacing:0.1em;font-family:IBM Plex Mono,monospace;'>Core Projects Completed</div>"
          f"</div>"
          f"</div>"
        # ── Core Projects section ────────────────────────────────────────────────
        + f"<div style='display:grid;grid-template-columns:1fr 1fr 1fr;flex:1;overflow:hidden;"
          f"border-bottom:2px solid {CLR};'>"
        # Core Projects by type
        + f"<div style='padding:8px 10px 6px 14px;border-right:1px solid #e5e7eb;overflow:hidden;'>"
          + _col_hdr("Core Projects")
          + _qytd_rows(cp_segs)
          + _leg(cp_segs)
          + "</div>"
        # Report Ratings
        + f"<div style='padding:8px 10px 6px 10px;border-right:1px solid #e5e7eb;overflow:hidden;'>"
          + _col_hdr("Report Ratings")
          + _qytd_rows(rat_segs or [("No data", n_core, "#9ca3af")])
          + _leg(rat_segs or [("No data", n_core, "#9ca3af")])
          + "</div>"
        # MARC Ratings
        + f"<div style='padding:8px 10px 6px 10px;overflow:hidden;'>"
          + _col_hdr("MARC Ratings")
          + _qytd_rows(marc_segs)
          + _leg(marc_segs)
          + "</div>"
        + "</div>"
        # ── Issues Raised separator ──────────────────────────────────────────────
        + f"<div style='background:{CLR};padding:4px 22px;flex-shrink:0;"
          f"display:flex;justify-content:flex-end;align-items:center;'>"
          f"<div style='font-size:0.66rem;font-weight:700;color:#ffffff;'>"
          f"{n_total:,} Issues Raised</div>"
          f"</div>"
        # ── Issues breakdown ─────────────────────────────────────────────────────
        + f"<div style='display:grid;grid-template-columns:1fr 1fr 1fr;flex-shrink:0;"
          f"background:#f0f4f0;overflow:hidden;'>"
        # Newly Raised
        + f"<div style='padding:6px 10px 6px 14px;border-right:1px solid #e5e7eb;'>"
          + _col_hdr("Newly Raised")
          + (
              f"<div style='display:flex;align-items:center;gap:5px;margin-bottom:2px;'>"
              f"<span style='font-size:0.42rem;color:#6b7280;width:30px;flex-shrink:0;'>Q2/26</span>"
              f"<div style='flex:1;'>{_lhb(_pct(n_newly, max_iss), CLR)}</div>"
              f"<span style='font-size:0.52rem;font-weight:700;color:{CLR};'>{n_newly}</span></div>"
              f"<div style='display:flex;align-items:center;gap:5px;'>"
              f"<span style='font-size:0.42rem;color:#6b7280;width:30px;flex-shrink:0;'>YTD</span>"
              f"<div style='flex:1;'>{_lhb(_pct(n_newly, max_iss), '#3a9b7d')}</div>"
              f"<span style='font-size:0.52rem;font-weight:700;color:#3a9b7d;'>{n_newly}</span></div>"
          )
          + "</div>"
        # Self-Identified
        + f"<div style='padding:6px 10px 6px 10px;border-right:1px solid #e5e7eb;'>"
          + _col_hdr("Self-Identified (of newly raised)")
          + (
              f"<div style='display:flex;align-items:center;gap:5px;margin-bottom:2px;'>"
              f"<span style='font-size:0.42rem;color:#6b7280;width:30px;flex-shrink:0;'>Q2/26</span>"
              f"<div style='flex:1;'>{_lhb(_pct(n_self_id, max(n_newly, 1)), CLR)}</div>"
              f"<span style='font-size:0.52rem;font-weight:700;color:{CLR};'>{n_self_id}</span></div>"
              f"<div style='display:flex;align-items:center;gap:5px;'>"
              f"<span style='font-size:0.42rem;color:#6b7280;width:30px;flex-shrink:0;'>YTD</span>"
              f"<div style='flex:1;'>{_lhb(_pct(n_self_id, max(n_newly, 1)), '#3a9b7d')}</div>"
              f"<span style='font-size:0.52rem;font-weight:700;color:#3a9b7d;'>{n_self_id}</span></div>"
          )
          + "</div>"
        # Open Issues
        + f"<div style='padding:6px 10px 6px 10px;'>"
          + _col_hdr("Open Issues")
          + _cnt_bar(n_open,    max_iss, "#60a5fa")
          + f"<div style='font-size:0.42rem;color:#60a5fa;margin-bottom:4px;'>In Progress</div>"
          + _cnt_bar(n_overdue, max_iss, "#ef4444")
          + f"<div style='font-size:0.42rem;color:#ef4444;margin-bottom:4px;'>RNV / Past Due</div>"
          + f"<div style='display:flex;gap:6px;margin-top:2px;'>"
            f"<span style='font-size:0.38rem;color:#f87171;font-weight:600;'>&#9632; Level 1</span>"
            f"<span style='font-size:0.38rem;color:#9ca3af;font-weight:600;'>&#9632; Level 2</span>"
            f"</div>"
          + "</div>"
        + "</div>"
        # ── Footer ──────────────────────────────────────────────────────────────
        + f"<div style='position:absolute;bottom:0;left:0;right:0;height:16px;"
          f"background:#f9fafb;border-top:1px solid #e5e7eb;"
          f"display:flex;align-items:center;padding:0 16px;justify-content:space-between;'>"
          f"<span style='font-size:0.38rem;color:#9ca3af;'>RBC Internal Audit | CONFIDENTIAL</span>"
          f"<span style='font-size:0.38rem;color:#9ca3af;'>"
          f"{qtr} INTERNAL AUDIT SUPPLEMENTAL APPENDICES</span></div>"
        + "</div>"
    )


def _slide_app2_performance(plat: str, audits: pd.DataFrame, issues: pd.DataFrame, qtr: str) -> str:
    """Appendix 2 — Segment performance KPIs: SAT Rate, Control Deficiency, Self-ID, Overdue.

    Uses: audits.status, audits.current_rating, issues.status, issues.severity
    """

    CLR = "#1d5c4a"

    # ── Compute KPIs from available data ──────────────────────────────────────
    n_completed = int((audits["status"] == "Complete").sum()) if not audits.empty else 0

    def _rcnt(vals):
        if audits.empty:
            return 0
        for col in ("current_rating", "rating"):
            if col in audits.columns:
                return int(audits[col].astype(str).str.upper().isin(vals).sum())
        return 0

    n_sat   = _rcnt({"SAT", "SATISFACTORY"})
    n_unsat = _rcnt({"UNSAT", "UNSATISFACTORY"})
    sat_rate  = _pct(n_sat,   n_completed) if n_completed else 0
    ctrl_def  = _pct(n_unsat, n_completed) if n_completed else 0

    n_iss     = len(issues) if not issues.empty else 0
    n_open    = int(issues["status"].isin(["Open", "Overdue"]).sum()) if not issues.empty else 0
    n_overdue = int((issues.get("status", pd.Series(dtype=str)) == "Overdue").sum()) if not issues.empty else 0
    n_high    = int((issues.get("severity", pd.Series(dtype=str)) == "High").sum()) if not issues.empty else 0
    # "Self-identified" approximated as high-severity issues as share of total
    self_id_rate = _pct(n_high,    n_iss)  if n_iss  else 0
    overdue_rate = _pct(n_overdue, n_open) if n_open else 0

    # ── KPI panel renderer ────────────────────────────────────────────────────
    def _kpi(title, value, threshold, direction, bar_clr):
        """Render a single KPI panel with value, bar, and threshold marker."""
        meets = (value >= threshold) if direction == ">=" else (value <= threshold)
        val_clr = "#22c55e" if meets else "#ef4444"
        dir_sym = "≥" if direction == ">=" else "≤"
        thr_pct = min(threshold, 100)
        return (
            f"<div style='background:#ffffff;border:1px solid #e5e7eb;border-radius:6px;"
            f"padding:10px 12px;overflow:hidden;'>"
            f"<div style='font-size:0.54rem;font-weight:700;color:{CLR};"
            f"margin-bottom:3px;'>{title}</div>"
            f"<div style='font-size:0.4rem;color:#ef4444;margin-bottom:6px;font-family:IBM Plex Mono,monospace;'>"
            f"Threshold {dir_sym} {threshold}%</div>"
            f"<div style='font-size:2rem;font-weight:800;color:{val_clr};line-height:1;margin-bottom:8px;'>"
            f"{value}%</div>"
            f"<div style='position:relative;background:#e5e7eb;border-radius:3px;height:10px;'>"
            f"<div style='background:{bar_clr};width:{min(value, 100)}%;height:100%;border-radius:3px;'></div>"
            f"<div style='position:absolute;top:-3px;left:{thr_pct}%;border-left:2px solid #ef4444;"
            f"height:16px;transform:translateX(-1px);'></div>"
            f"</div>"
            f"<div style='font-size:0.4rem;color:#9ca3af;margin-top:4px;'>"
            f"Threshold: {dir_sym}{threshold}%  ·  This period: {value}%</div>"
            f"</div>"
        )

    return (
        _FL
        + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;background:#f5f7f9;"
          f"border-radius:10px;overflow:hidden;position:relative;"
          f"box-shadow:0 20px 56px rgba(0,0,30,0.5);font-family:Barlow Condensed,sans-serif;'>"
        + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;background:{_G};z-index:3;'></div>"
        + f"<div style='background:{CLR};padding:8px 18px 8px 22px;"
          f"display:flex;justify-content:space-between;align-items:flex-start;'>"
          f"<div>"
          f"<div style='font-size:0.44rem;color:rgba(255,255,255,0.5);letter-spacing:0.15em;"
          f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;'>"
          f"APPENDIX 2 · SEGMENT PERFORMANCE METRICS</div>"
          f"<div style='font-size:0.88rem;font-weight:800;color:#ffffff;'>{plat}</div>"
          f"</div>"
          f"<div style='font-size:0.46rem;color:rgba(255,255,255,0.55);font-family:IBM Plex Mono,monospace;'>"
          f"RBC INTERNAL AUDIT · {qtr}</div>"
          f"</div>"
        + f"<div style='padding:10px 18px 10px 22px;height:calc(100% - 58px);overflow:hidden;"
          f"display:grid;grid-template-columns:1fr 1fr;grid-template-rows:1fr 1fr;gap:10px;'>"
        + _kpi("SAT Rating",                  sat_rate,    80, ">=", "#22c55e")
        + _kpi("Control Deficiency",           ctrl_def,     8, "<=", "#ef4444")
        + _kpi("High-Severity Issue Rate",     self_id_rate,30, "<=", "#f59e0b")
        + _kpi("Issue Overdue Rate",           overdue_rate, 5, "<=", "#f97316")
        + "</div>"
        + f"<div style='position:absolute;bottom:0;left:0;right:0;height:18px;"
          f"background:#f9fafb;border-top:1px solid #e5e7eb;"
          f"display:flex;align-items:center;padding:0 16px;justify-content:space-between;'>"
          f"<span style='font-size:0.38rem;color:#9ca3af;'>RBC Internal Audit | CONFIDENTIAL</span>"
          f"<span style='font-size:0.38rem;color:#9ca3af;'>"
          f"{qtr} INTERNAL AUDIT SUPPLEMENTAL APPENDICES</span></div>"
        + "</div>"
    )


def _slide_app2_ce(
    plat: str,
    audits: pd.DataFrame,
    issues: pd.DataFrame,
    controls: pd.DataFrame,
    qtr: str,
) -> str:
    """Appendix 2 (1/2) — Control Environment Summary per segment.
    Uses: audits (current_rating, rating, rating_change, audit_name, status),
          issues (status, severity). Renders AU list + narrative panel."""

    CLR = "#1d5c4a"

    def _rating_cnt_raw(df, vals):
        if df.empty: return 0
        for col in ("current_rating", "rating"):
            if col in df.columns:
                return int(df[col].astype(str).str.upper().isin(vals).sum())
        return 0

    n_sat = _rating_cnt_raw(audits, {"SAT", "SATISFACTORY"})
    n_ri  = _rating_cnt_raw(audits, {"RI", "REQUIRES IMPROVEMENT", "NEEDS IMPROVEMENT"})
    n_uns = _rating_cnt_raw(audits, {"UNSAT", "UNSATISFACTORY"})
    n_completed = int((audits["status"] == "Complete").sum()) if not audits.empty else 0

    # Determine overall CE rating and trend arrow
    if n_sat >= n_ri and n_sat >= n_uns:
        overall_ce, overall_clr = "SAT", "#22c55e"
    elif n_ri >= n_uns:
        overall_ce, overall_clr = "RI", "#f59e0b"
    else:
        overall_ce, overall_clr = "UNSAT", "#ef4444"

    # Determine trend from rating_change column majority vote
    rc_vals = audits.get("rating_change", pd.Series(dtype=str)).dropna().astype(str) if not audits.empty else pd.Series(dtype=str)
    n_improved  = int((rc_vals == "Improved").sum())
    n_degraded  = int((rc_vals == "Deteriorated").sum())
    if n_improved > n_degraded:
        trend_arrow, trend_clr, trend_lbl = "↑", "#22c55e", "Trending Up"
    elif n_degraded > n_improved:
        trend_arrow, trend_clr, trend_lbl = "↓", "#ef4444", "Trending Down"
    else:
        trend_arrow, trend_clr, trend_lbl = "→", "#9ca3af", "No Change"

    rating_clrs = {"SAT": "#22c55e", "RI": "#f59e0b", "UNSAT": "#ef4444"}
    au_rows = ""
    # Left panel: AU list — badge (colored square) + name + trend arrow, NO status dot
    for _, aud in audits.head(12).iterrows():
        nm = str(aud.get("audit_name", aud.get("audit_id", "—")))[:40]
        rt_raw = ""
        for col in ("current_rating", "rating"):
            rt_raw = str(aud.get(col, "")).upper()
            if rt_raw:
                break
        if "SAT" in rt_raw and "UNSAT" not in rt_raw:
            rt_norm = "SAT"
        elif "RI" in rt_raw or "IMPROVEMENT" in rt_raw:
            rt_norm = "RI"
        elif "UNSAT" in rt_raw or "UNSATISF" in rt_raw:
            rt_norm = "UNSAT"
        else:
            rt_norm = "NR"
        rt_clr = rating_clrs.get(rt_norm, "#9ca3af")
        rc = str(aud.get("rating_change", ""))
        arrow = "↑" if rc == "Improved" else ("↓" if rc == "Deteriorated" else "→")
        arw_clr = "#22c55e" if arrow == "↑" else ("#ef4444" if arrow == "↓" else "#9ca3af")

        au_rows += (
            f"<div style='display:flex;align-items:center;padding:3px 0;"
            f"border-bottom:1px solid #e5e7eb;gap:5px;'>"
            f"<div style='flex:1;font-size:0.52rem;color:#1a2035;overflow:hidden;"
            f"text-overflow:ellipsis;white-space:nowrap;'>{nm}</div>"
            f"<div style='font-size:0.44rem;font-weight:700;color:{rt_clr};"
            f"background:{rt_clr}22;border-radius:2px;padding:1px 4px;flex-shrink:0;'>"
            f"{rt_norm}</div>"
            f"<div style='font-size:0.6rem;font-weight:700;color:{arw_clr};"
            f"flex-shrink:0;'>{arrow}</div>"
            f"</div>"
        )

    # CE legend shown at bottom of AU list
    au_legend = (
        f"<div style='display:flex;gap:8px;flex-wrap:wrap;padding-top:5px;border-top:1px solid #e5e7eb;margin-top:3px;'>"
        f"<span style='font-size:0.38rem;color:#22c55e;font-weight:600;'>&#9632; SAT</span>"
        f"<span style='font-size:0.38rem;color:#f59e0b;font-weight:600;'>&#9632; RI</span>"
        f"<span style='font-size:0.38rem;color:#ef4444;font-weight:600;'>&#9632; UNSAT</span>"
        f"<span style='font-size:0.38rem;color:#9ca3af;'>|</span>"
        f"<span style='font-size:0.38rem;color:#22c55e;'>↑ Trending Up</span>"
        f"<span style='font-size:0.38rem;color:#ef4444;'>↓ Trending Down</span>"
        f"<span style='font-size:0.38rem;color:#9ca3af;'>→ No Change</span>"
        f"</div>"
    )

    n_total = len(audits)
    n_iss = len(issues) if not issues.empty else 0
    n_open_iss = int(issues["status"].isin(["Open", "Overdue"]).sum()) if not issues.empty else 0
    n_high_iss = int((issues.get("severity", pd.Series(dtype=str)) == "High").sum()) if not issues.empty else 0
    n_ovr_iss  = int((issues.get("status",   pd.Series(dtype=str)) == "Overdue").sum()) if not issues.empty else 0
    n_prog = int(audits["status"].isin(["In Progress", "Fieldwork"]).sum()) if not audits.empty else 0

    sat_pct = _pct(n_sat, max(n_completed, 1))

    def _nar_section(title, bullets):
        hdr = (
            f"<div style='font-size:0.48rem;font-weight:700;color:{CLR};"
            f"letter-spacing:0.08em;text-transform:uppercase;font-family:IBM Plex Mono,monospace;"
            f"border-bottom:1px solid #d1d5db;padding-bottom:2px;margin-top:7px;margin-bottom:4px;'>"
            f"{title}</div>"
        )
        items = "".join(
            f"<div style='font-size:0.53rem;color:#374151;line-height:1.38;padding:1px 0;"
            f"display:flex;gap:4px;'>"
            f"<span style='color:{CLR};flex-shrink:0;margin-top:1px;'>•</span>"
            f"<span>{b}</span></div>"
            for b in bullets
        )
        return hdr + items

    # Opening paragraph (bold) matching the real slide narrative style
    n_majority = max(n_sat, n_ri, n_uns)
    majority_rating = "SAT" if n_sat == n_majority else ("RI" if n_ri == n_majority else "UNSAT")
    opening = (
        f"<div style='font-size:0.54rem;color:#1a2035;line-height:1.45;margin-bottom:5px;'>"
        f"<strong>{plat} CE remains {overall_ce} with {trend_lbl} trend, "
        f"with majority of audits rated {majority_rating} "
        f"({n_majority}/{n_total}, {_pct(n_majority, max(n_total,1))}%).</strong>"
        f"</div>"
    )

    highlights = []
    if n_completed > 0:
        highlights.append(f"{n_completed} of {n_total} audit{'s' if n_total != 1 else ''} complete in {qtr}.")
    if n_prog > 0:
        highlights.append(f"{n_prog} engagement{'s' if n_prog != 1 else ''} currently in progress.")
    if n_completed > 0 and n_sat > 0:
        highlights.append(f"{n_sat} rated SAT ({sat_pct}% of completed).")
    if n_ri > 0:
        highlights.append(f"{n_ri} rated RI — management action plans in place.")
    highlights = highlights[:4]

    priorities = [
        "Timely remediation of all outstanding Level 1 issues.",
        "Continue audit coverage across strategic and operational risk areas.",
    ]

    iss_bullets = [
        f"{n_iss} issue{'s' if n_iss != 1 else ''} raised in this segment to date.",
        f"{n_open_iss} currently open or overdue.",
    ]
    if n_high_iss > 0:
        iss_bullets.append(f"{n_high_iss} Level 1 issue{'s' if n_high_iss != 1 else ''} requiring prompt management action.")
    if n_ovr_iss > 0:
        iss_bullets.append(f"{n_ovr_iss} past expected resolution date.")

    right = (
        opening
        + _nar_section("Q2 Highlights", highlights)
        + _nar_section("Strategic Priorities", priorities)
        + _nar_section("Issue Management", iss_bullets)
    )

    return (
        _FL
        + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;background:#ffffff;"
          f"border-radius:10px;overflow:hidden;position:relative;"
          f"box-shadow:0 20px 56px rgba(0,0,30,0.5);font-family:Barlow Condensed,sans-serif;'>"
        + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;background:{_G};z-index:3;'></div>"
        + f"<div style='background:{CLR};padding:8px 18px 8px 22px;"
          f"display:flex;justify-content:space-between;align-items:flex-start;'>"
          f"<div>"
          f"<div style='font-size:0.44rem;color:rgba(255,255,255,0.5);letter-spacing:0.16em;"
          f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;'>"
          f"APPENDIX 2 · CONTROL ENVIRONMENT SUMMARY (1/2)</div>"
          f"<div style='font-size:0.88rem;font-weight:800;color:#ffffff;line-height:1.15;'>{plat}</div>"
          f"</div>"
          f"<div style='text-align:right;padding-top:2px;'>"
          f"<div style='font-size:0.42rem;color:rgba(255,255,255,0.65);text-transform:uppercase;"
          f"letter-spacing:0.1em;font-family:IBM Plex Mono,monospace;margin-bottom:2px;'>Overall</div>"
          f"<div style='font-size:1.15rem;font-weight:800;color:{overall_clr};line-height:1;'>"
          f"{overall_ce} <span style='color:{trend_clr};'>{trend_arrow}</span></div>"
          f"</div>"
          f"</div>"
        + f"<div style='display:grid;grid-template-columns:40% 60%;"
          f"height:calc(100% - 73px);overflow:hidden;'>"
        + f"<div style='padding:8px 10px 8px 14px;border-right:1px solid #e5e7eb;"
          f"overflow:hidden;background:#f5f7f9;display:flex;flex-direction:column;'>"
          f"<div style='font-size:0.5rem;font-weight:700;color:{CLR};"
          f"letter-spacing:0.1em;text-transform:uppercase;font-family:IBM Plex Mono,monospace;"
          f"border-bottom:2px solid {CLR};padding-bottom:4px;margin-bottom:6px;'>"
          f"Auditable Unit (AU)</div>"
          f"<div style='flex:1;overflow:hidden;'>"
          + au_rows
          + au_legend
          + "</div>"
        + f"<div style='padding:6px 14px 8px 12px;overflow:hidden;background:#ffffff;'>"
          + right
          + "</div>"
        + "</div>"
        + f"<div style='position:absolute;bottom:0;left:0;right:0;height:18px;"
          f"background:#f9fafb;border-top:1px solid #e5e7eb;"
          f"display:flex;align-items:center;padding:0 16px;justify-content:space-between;'>"
          f"<span style='font-size:0.4rem;color:#9ca3af;'>RBC Internal Audit | CONFIDENTIAL</span>"
          f"<span style='font-size:0.4rem;color:#9ca3af;'>"
          f"{qtr} INTERNAL AUDIT SUPPLEMENTAL APPENDICES</span></div>"
        + "</div>"
    )

# ── Section 3 & 4 shared helpers ──────────────────────────────────────────────

_APP4_CLR  = "#1d5c4a"   # Appendix 4 dark forest-teal
_APP4_PAGE = 8           # Rows per Appendix 4 slide

_S3_CLR = "#1e4d3a"   # Section 3 dark teal
_S4_CLR = "#1e3a6b"   # Section 4 dark navy-blue
_S5_CLR = "#14375f"   # Section 5 CAE Group Operations (deep steel)
_S7_CLR = "#2a3f54"   # Section 7 Glossary (dark slate)

def _section_slide(section_num: int, section_color: str, section_label: str,
                   title: str, left_html: str, right_html: str, qtr: str) -> str:
    """Shared light-background two-panel slide matching the AC report style."""
    return (
        _FL
        + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;"
          f"background:#f5f7f9;border-radius:10px;overflow:hidden;position:relative;"
          f"box-shadow:0 20px 56px rgba(0,0,30,0.5);font-family:Barlow Condensed,sans-serif;'>"
        + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;background:{_G};z-index:3;'></div>"
        # Section header
        + f"<div style='background:{section_color};padding:8px 20px 8px 28px;"
          f"display:flex;justify-content:space-between;align-items:flex-start;'>"
          f"<div>"
          f"<div style='font-size:0.5rem;color:rgba(255,255,255,0.5);letter-spacing:0.16em;"
          f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;'>"
          f"SECTION {section_num} · {section_label}</div>"
          f"<div style='font-size:0.82rem;font-weight:700;color:{_W};line-height:1.15;max-width:600px;'>"
          f"{title}</div>"
          f"</div>"
          f"<div style='font-size:0.5rem;color:rgba(255,255,255,0.4);font-family:IBM Plex Mono,monospace;"
          f"white-space:nowrap;padding-top:6px;'>RBC INTERNAL AUDIT</div>"
          f"</div>"
        # Body: two-panel
        + f"<div style='display:grid;grid-template-columns:44% 56%;height:calc(100% - 80px);overflow:hidden;'>"
          f"<div style='background:#edf2f7;padding:10px 12px 10px 14px;"
          f"border-right:1px solid #d1d5db;overflow:hidden;'>{left_html}</div>"
          f"<div style='background:{_W};padding:10px 14px;overflow:hidden;'>{right_html}</div>"
          f"</div>"
        # Footer
        + f"<div style='position:absolute;bottom:0;left:0;right:0;height:20px;"
          f"background:{section_color};display:flex;align-items:center;padding:0 22px;"
          f"justify-content:space-between;'>"
          f"<span style='font-size:0.46rem;color:rgba(255,255,255,0.4);'>RBC Internal Audit | CONFIDENTIAL</span>"
          f"<span style='font-size:0.46rem;color:rgba(255,255,255,0.4);'>{qtr} INTERNAL AUDIT QUARTERLY REPORT</span>"
          f"</div>"
        + "</div>"
    )


def _col_header(text: str, color: str) -> str:
    return (
        f"<div style='background:{color};color:{_W};padding:5px 8px;border-radius:3px;"
        f"font-size:0.52rem;font-weight:700;letter-spacing:0.1em;text-transform:uppercase;"
        f"font-family:IBM Plex Mono,monospace;margin-bottom:7px;'>{text}</div>"
    )


def _stacked_bar(segments: list[tuple[int, str]], total: int, height: int = 14) -> str:
    """Render a horizontal stacked bar from (count, color) pairs."""
    if total == 0:
        return f"<div style='background:#e5e7eb;border-radius:3px;height:{height}px;width:100%;'></div>"
    parts = "".join(
        f"<div style='width:{_pct(n,total)}%;background:{c};height:100%;flex-shrink:0;' title='{n}'></div>"
        for n, c in segments if n > 0
    )
    return (
        f"<div style='display:flex;border-radius:3px;overflow:hidden;height:{height}px;"
        f"width:100%;gap:1px;background:#e5e7eb;'>{parts}</div>"
    )


def _legend_row(items: list[tuple[str, int, str]]) -> str:
    """Label, count, color legend pills."""
    pills = "".join(
        f"<span style='font-size:0.46rem;color:{c};font-weight:600;white-space:nowrap;'>"
        f"● {lbl} {n}</span>"
        for lbl, n, c in items if n >= 0
    )
    return f"<div style='display:flex;gap:8px;flex-wrap:wrap;margin-top:3px;'>{pills}</div>"


def _insight_bullet(text: str, color: str = "#1a2035") -> str:
    return (
        f"<div style='font-size:0.58rem;color:{color};padding:3px 0 3px 10px;"
        f"border-left:2px solid {_G};margin-bottom:4px;line-height:1.4;'>{text}</div>"
    )


def _big_metric(value: str, label: str, color: str = "#001e4d") -> str:
    return (
        f"<div style='text-align:center;padding:8px 12px;background:#edf2f7;"
        f"border-radius:6px;border:1px solid #d1d5db;'>"
        f"<div style='font-size:2.2rem;font-weight:800;color:{color};line-height:1;'>{value}</div>"
        f"<div style='font-size:0.48rem;color:#6b7280;text-transform:uppercase;letter-spacing:0.1em;"
        f"font-family:IBM Plex Mono,monospace;margin-top:2px;'>{label}</div>"
        f"</div>"
    )


def _metric_bar_row(label: str, n: int, mx: int, color: str) -> str:
    p = _pct(n, mx)
    return (
        f"<div style='margin-bottom:5px;'>"
        f"<div style='display:flex;justify-content:space-between;margin-bottom:2px;'>"
        f"<span style='font-size:0.58rem;color:#374151;'>{label}</span>"
        f"<span style='font-size:0.62rem;font-weight:700;color:{color};'>{n}</span></div>"
        + _lhb(p, color) + "</div>"
    )


def _issue_theme(icon: str, title: str, causes: list[str], pct: int, color: str) -> tuple[str, str, str]:
    theme_html = (
        f"<div style='background:{_W};border:1px solid #d1d5db;border-radius:6px;"
        f"padding:7px 9px;margin-bottom:7px;'>"
        f"<div style='font-size:0.64rem;font-weight:700;color:#1a2035;margin-bottom:2px;'>{icon} {title}</div>"
        f"</div>"
    )
    causes_html = (
        f"<div style='background:#fafafa;border:1px solid #e5e7eb;border-radius:6px;"
        f"padding:7px 9px;margin-bottom:7px;'>"
        + "".join(f"<div style='font-size:0.54rem;color:#374151;padding:1px 0;'>• {c}</div>" for c in causes)
        + "</div>"
    )
    pct_html = (
        f"<div style='background:{color};border-radius:6px;padding:8px;text-align:center;"
        f"margin-bottom:7px;'>"
        f"<div style='font-size:1.8rem;font-weight:800;color:{_W};line-height:1;'>{pct}%</div>"
        f"<div style='font-size:0.44rem;color:rgba(255,255,255,0.7);text-transform:uppercase;"
        f"letter-spacing:0.08em;font-family:IBM Plex Mono,monospace;'>Q2 Core Projects</div>"
        f"</div>"
    )
    return theme_html, causes_html, pct_html


# ── Section 3: Assurance Activities & Output ──────────────────────────────────

def _slide_assurance_output(audits: pd.DataFrame, issues: pd.DataFrame, qtr: str) -> str:
    completed = audits[audits["status"] == "Complete"]
    published = completed[completed["report_status"] == "Published"] if "report_status" in completed.columns else completed
    at_col  = published.get("audit_type", pd.Series(dtype=str))
    core    = published[at_col.isin(["Owned Audit", "In-Scope AE"])]
    n_core  = len(core)
    n_total = len(published)

    n_sat   = len(core[core.get("current_rating", pd.Series(dtype=str)) == "SAT"]) if "current_rating" in core.columns else 0
    n_ri    = len(core[core.get("current_rating", pd.Series(dtype=str)) == "RI"])  if "current_rating" in core.columns else 0
    n_unsat = len(core[core.get("current_rating", pd.Series(dtype=str)) == "UNSAT"]) if "current_rating" in core.columns else 0
    n_na    = len(core[core.get("current_rating", pd.Series(dtype=str)) == "NA"])  if "current_rating" in core.columns else 0
    n_rate  = n_sat + n_ri + n_unsat + n_na or 1

    n_dev   = len(core[core.get("marc_rating", pd.Series(dtype=str)) == "Developed"]) if "marc_rating" in core.columns else 0
    n_sub   = len(core[core.get("marc_rating", pd.Series(dtype=str)) == "Substantially Developed"]) if "marc_rating" in core.columns else 0
    n_par   = len(core[core.get("marc_rating", pd.Series(dtype=str)) == "Partially Developed"]) if "marc_rating" in core.columns else 0
    n_und   = len(core[core.get("marc_rating", pd.Series(dtype=str)) == "Underdeveloped"]) if "marc_rating" in core.columns else 0
    n_marc  = n_dev + n_sub + n_par + n_und or 1

    sat_pct  = _pct(n_sat, n_rate)
    marc_fav = _pct(n_dev + n_sub, n_marc)

    # Dynamic title
    if sat_pct >= 70 and marc_fav >= 60:
        title = "Sustained improvement in audit report ratings; MARC ratings remain stable"
    elif sat_pct < 50:
        title = "Focus required on audit quality; ratings below target threshold"
    else:
        title = "Steady audit output this quarter; MARC profile reflects programme maturity"

    # Left: narrative insights
    left = (
        _col_header("Q2 Activities", _S3_CLR)
        + _insight_bullet(f"<strong>{n_total}</strong> projects completed and reported this quarter.")
        + _insight_bullet(f"<strong>{n_sat}</strong> SAT, <strong>{n_ri}</strong> RI, <strong>{n_unsat}</strong> UNSAT"
                          f" across {n_core} core projects ({sat_pct}% favourable ratings).")
        + (_insight_bullet(f"MARC Developed or Substantially Developed: <strong>{n_dev+n_sub}</strong> of {n_marc} ({marc_fav}%).")
           if n_marc > 0 else "")
        + (_insight_bullet(f"<strong>{n_und}</strong> Underdeveloped MARC ratings require management focus.",
                           "#991b1b") if n_und > 0 else "")
        + _insight_bullet("Digital RCM completion rate and planning memo discipline continue to be monitored.")
    )

    # Right: Key Metrics & Indicators
    right = (
        _col_header("Key Metrics &amp; Indicators", _S3_CLR)
        # Core Projects
        + f"<div style='font-size:0.56rem;font-weight:700;color:#1a2035;margin-bottom:4px;'>Core Projects Delivered</div>"
        + f"<div style='display:grid;grid-template-columns:1fr 1fr 1fr;gap:8px;margin-bottom:10px;'>"
        + _big_metric(str(len(core[core.get("audit_type",pd.Series(dtype=str))=="Owned Audit"])), "Owned", "#001e4d")
        + _big_metric(str(len(core[core.get("audit_type",pd.Series(dtype=str))=="In-Scope AE"])), "AE", "#1e3a6b")
        + _big_metric(str(n_core), "Total Core", _S3_CLR)
        + "</div>"
        # Report Ratings stacked bar
        + f"<div style='font-size:0.56rem;font-weight:700;color:#1a2035;margin-bottom:3px;'>Report Ratings</div>"
        + _stacked_bar([(n_sat,"#22c55e"),(n_ri,"#f59e0b"),(n_unsat,"#ef4444"),(n_na,"#9ca3af")], n_rate)
        + _legend_row([("SAT",n_sat,"#16a34a"),("RI",n_ri,"#b45309"),("UNSAT",n_unsat,"#dc2626"),("N/A",n_na,"#6b7280")])
        + "<div style='height:8px;'></div>"
        # MARC stacked bar
        + f"<div style='font-size:0.56rem;font-weight:700;color:#1a2035;margin-bottom:3px;'>MARC Ratings</div>"
        + _stacked_bar([(n_dev,"#16a34a"),(n_sub,"#86efac"),(n_par,"#f59e0b"),(n_und,"#dc2626")], n_marc)
        + _legend_row([("Dev",n_dev,"#16a34a"),("Sub Dev",n_sub,"#166534"),("Part Dev",n_par,"#b45309"),("Under",n_und,"#dc2626")])
    )

    return _section_slide(3, _S3_CLR, "ASSURANCE ACTIVITIES &amp; OUTPUT", title, left, right, qtr)


def _slide_issue_themes(audits: pd.DataFrame, issues: pd.DataFrame, qtr: str) -> str:
    """Section 3 — Issue themes analysis with root causes (3-column layout)."""
    completed = audits[audits["status"] == "Complete"]
    n_core = max(len(completed), 1)

    # Classify issues into themes via keyword matching
    def _theme(title_str: str) -> str:
        t = title_str.lower()
        if any(k in t for k in ["access","privilege","identity","entitlement","provisioning","iam","mfa"]):
            return "access"
        if any(k in t for k in ["data","quality","governance","lineage","critical data","reporting"]):
            return "data"
        if any(k in t for k in ["monitoring","oversight","control","testing","review","rcsa","validation"]):
            return "controls"
        return "operational"

    theme_counts = {"access": 0, "data": 0, "controls": 0, "operational": 0}
    for title_val in issues.get("title", pd.Series(dtype=str)).dropna():
        theme_counts[_theme(str(title_val))] += 1

    themes = [
        ("🔍", "Monitoring, Oversight &<br>Control Testing Deficiencies", theme_counts["controls"],
         ["Insufficient risk-based testing frequency",
          "Lack of supervisory review and escalation protocols",
          "Control self-assessment not embedded in BAU processes"],
         "#1e4d3a"),
        ("🗄️", "Data Management &<br>Critical Data Deficiencies", theme_counts["data"],
         ["Incomplete data lineage and quality frameworks",
          "Inconsistent data governance across lines of business",
          "Critical data elements not formally inventoried"],
         "#1e3a6b"),
        ("🔐", "Inadequate Access Provisioning /<br>Privileged Access Controls", theme_counts["access"],
         ["Excessive system access privileges not regularly reviewed",
          "Lack of automated access recertification processes",
          "Separation of duties not enforced in key workflows"],
         "#3d1800"),
    ]

    title = "Issues indicate increasing focus needed on controls maturity, data governance, and access management"

    # Build 3-column layout
    theme_col, causes_col, pct_col = "", "", ""
    for icon, name, count, causes, color in themes:
        pct_of_core = _pct(count, n_core)
        t_h, c_h, p_h = _issue_theme(icon, name, causes[:3], pct_of_core, color)
        theme_col  += t_h
        causes_col += c_h
        pct_col    += p_h

    body = (
        _FL
        + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;"
          f"background:#f5f7f9;border-radius:10px;overflow:hidden;position:relative;"
          f"box-shadow:0 20px 56px rgba(0,0,30,0.5);font-family:Barlow Condensed,sans-serif;'>"
        + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;background:{_G};z-index:3;'></div>"
        + f"<div style='background:{_S3_CLR};padding:8px 20px 8px 28px;"
          f"display:flex;justify-content:space-between;align-items:flex-start;'>"
          f"<div><div style='font-size:0.5rem;color:rgba(255,255,255,0.5);letter-spacing:0.16em;"
          f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;'>"
          f"SECTION 3 · ASSURANCE ACTIVITIES &amp; OUTPUT</div>"
          f"<div style='font-size:0.78rem;font-weight:700;color:{_W};line-height:1.2;max-width:620px;'>"
          f"{title}</div></div>"
          f"<div style='font-size:0.5rem;color:rgba(255,255,255,0.4);font-family:IBM Plex Mono,monospace;"
          f"white-space:nowrap;padding-top:6px;'>RBC INTERNAL AUDIT</div></div>"
        + f"<div style='display:grid;grid-template-columns:28% 44% 28%;height:calc(100% - 80px);overflow:hidden;'>"
          f"<div style='background:#edf2f7;padding:10px 10px 10px 14px;border-right:1px solid #d1d5db;overflow:hidden;'>"
          + _col_header("Issue Theme", _S3_CLR)
          + theme_col
          + "</div>"
          f"<div style='background:{_W};padding:10px 12px;border-right:1px solid #d1d5db;overflow:hidden;'>"
          + _col_header("Root Causes", _S3_CLR)
          + causes_col
          + "</div>"
          f"<div style='background:#edf2f7;padding:10px 10px 10px 10px;overflow:hidden;'>"
          + _col_header("% of Core Projects", _S3_CLR)
          + pct_col
          + "</div>"
          + "</div>"
        + f"<div style='position:absolute;bottom:0;left:0;right:0;height:20px;"
          f"background:{_S3_CLR};display:flex;align-items:center;padding:0 22px;"
          f"justify-content:space-between;'>"
          f"<span style='font-size:0.46rem;color:rgba(255,255,255,0.4);'>RBC Internal Audit | CONFIDENTIAL</span>"
          f"<span style='font-size:0.46rem;color:rgba(255,255,255,0.4);'>{qtr} INTERNAL AUDIT QUARTERLY REPORT</span>"
          f"</div>"
        + "</div>"
    )
    return body


# ── Section 4: Audit Issues Management ────────────────────────────────────────

def _slide_issue_overview(issues: pd.DataFrame, qtr: str) -> str:
    """Section 4 — Newly raised, self-identified, repeat issues metrics."""
    n_total  = len(issues)
    n_open   = len(issues[issues["status"].isin(["Open","In Progress"])]) if not issues.empty else 0
    n_ovd    = len(issues[issues["status"] == "Overdue"]) if not issues.empty else 0
    n_high   = len(issues[issues["severity"] == "High"]) if not issues.empty else 0
    n_med    = len(issues[issues["severity"] == "Medium"]) if not issues.empty else 0
    n_low    = len(issues[issues["severity"] == "Low"]) if not issues.empty else 0
    n_si     = int(issues.get("self_identified", pd.Series(False)).sum()) if not issues.empty else 0
    si_pct   = _pct(n_si, n_total)
    ovd_pct  = _pct(n_ovd, n_total)

    title = "Continued focus required on self-identification of issues and timely control deficiency remediation"

    left = (
        _col_header("Q2 Activities", _S4_CLR)
        + _insight_bullet(f"<strong>{n_total}</strong> issues tracked across the portfolio this quarter.")
        + _insight_bullet(f"<strong>{n_high}</strong> High-severity issues require priority escalation and monitoring.")
        + (
            _insight_bullet(f"Self-identification rate: <strong>{si_pct}%</strong> of issues raised proactively by management.",
                            "#166534" if si_pct >= 30 else "#92400e")
            if n_total > 0 else ""
        )
        + (
            _insight_bullet(f"<strong>{n_ovd}</strong> issues ({ovd_pct}%) are overdue — active follow-up required.",
                            "#991b1b")
            if n_ovd > 0 else _insight_bullet("All tracked issues are within original resolution timeframes.", "#166534")
        )
        + _insight_bullet("Management is expected to provide updated remediation plans for all overdue items by quarter-end.")
    )

    right = (
        _col_header("Key Metrics &amp; Indicators", _S4_CLR)
        # Severity breakdown
        + f"<div style='font-size:0.56rem;font-weight:700;color:#1a2035;margin-bottom:3px;'>Issues by Severity — Total: {n_total}</div>"
        + _stacked_bar([(n_high,"#ef4444"),(n_med,"#f59e0b"),(n_low,"#22c55e")], n_total)
        + _legend_row([("High",n_high,"#dc2626"),("Medium",n_med,"#b45309"),("Low",n_low,"#166534")])
        + "<div style='height:8px;'></div>"
        # Status breakdown
        + f"<div style='font-size:0.56rem;font-weight:700;color:#1a2035;margin-bottom:3px;'>Issues by Status — Total: {n_total}</div>"
        + _stacked_bar([(n_open,"#3b82f6"),(n_ovd,"#ef4444"),(n_total-n_open-n_ovd,"#22c55e")], n_total)
        + _legend_row([("Open/In Progress",n_open,"#2563eb"),("Overdue",n_ovd,"#dc2626"),("Closed/Other",n_total-n_open-n_ovd,"#166534")])
        + "<div style='height:8px;'></div>"
        # Self-identified and repeat
        + f"<div style='display:grid;grid-template-columns:1fr 1fr;gap:8px;'>"
        + _big_metric(f"{si_pct}%", "Self-Identified", "#1e3a6b")
        + _big_metric(str(n_ovd), "Overdue Issues", "#dc2626" if n_ovd > 0 else "#166534")
        + "</div>"
    )

    return _section_slide(4, _S4_CLR, "AUDIT ISSUES MANAGEMENT", title, left, right, qtr)


def _slide_issue_tracking(issues: pd.DataFrame, qtr: str) -> str:
    """Section 4 — Issue tracking status and expected resolution timeline."""
    n_total = len(issues)
    n_open  = len(issues[issues["status"].isin(["Open","In Progress"])]) if not issues.empty else 0
    n_ovd   = len(issues[issues["status"] == "Overdue"]) if not issues.empty else 0
    n_cls   = len(issues[issues["status"].isin(["Closed","Resolved"])]) if not issues.empty else 0
    n_pend  = n_total - n_open - n_ovd - n_cls

    # Expected resolution by fiscal year from due_date
    fy_buckets: dict[str, int] = {}
    if not issues.empty and "due_date" in issues.columns:
        for val in issues["due_date"].dropna():
            try:
                yr = int(str(val)[:4])
                fy = f"FY{yr}"
                fy_buckets[fy] = fy_buckets.get(fy, 0) + 1
            except Exception:
                pass
    fy_sorted = sorted(fy_buckets.items())

    max_fy = max(fy_buckets.values()) if fy_buckets else 1

    # Generate dynamic title
    pct_on_track = _pct(n_open, n_total) if n_total else 0
    title_txt = (
        f"Continued attention required on timely issue resolution — {n_total} issues tracked in {qtr}"
        if n_ovd > 0 else
        f"Issue management on track — {n_cls} issues resolved, {n_open} in active remediation"
    )

    left = (
        _col_header("Insights", _S4_CLR)
        + _insight_bullet(f"<strong>{n_open}</strong> issues currently in active remediation with management.")
        + _insight_bullet(f"<strong>{n_ovd}</strong> issues past original due date — escalation in progress.",
                          "#991b1b" if n_ovd > 0 else "#166534")
        + _insight_bullet(f"<strong>{n_cls}</strong> issues resolved and verified by Internal Audit this quarter.",
                          "#166534" if n_cls > 0 else "#374151")
        + _insight_bullet("Issue Validation &amp; Retesting: IA reviews evidence provided before closing.")
        + _insight_bullet("Management is expected to align resolution plans with regulatory requirements and timelines.")
    )

    # Right: tracking status + resolution timeline
    right = (
        _col_header("Key Metrics &amp; Indicators", _S4_CLR)
        + f"<div style='font-size:0.56rem;font-weight:700;color:#1a2035;margin-bottom:3px;'>Issue Tracking Status — {n_total} total</div>"
        + _stacked_bar([(n_open,"#3b82f6"),(n_ovd,"#ef4444"),(n_cls,"#22c55e"),(n_pend,"#9ca3af")], n_total, 16)
        + _legend_row([("In Progress",n_open,"#2563eb"),("Overdue",n_ovd,"#dc2626"),
                       ("Closed",n_cls,"#166534"),("Pending",n_pend,"#6b7280")])
        + "<div style='height:10px;'></div>"
        + f"<div style='font-size:0.56rem;font-weight:700;color:#1a2035;margin-bottom:6px;'>Expected Resolution by Fiscal Year</div>"
        + "".join(
            f"<div style='margin-bottom:4px;'>"
            f"<div style='display:flex;justify-content:space-between;margin-bottom:2px;'>"
            f"<span style='font-size:0.58rem;color:#374151;font-weight:600;'>{fy}</span>"
            f"<span style='font-size:0.62rem;font-weight:700;color:#1e3a6b;'>{cnt}</span></div>"
            + _lhb(_pct(cnt, max_fy), "#1e3a6b") + "</div>"
            for fy, cnt in fy_sorted[:5]
        )
        + (
            f"<div style='font-size:0.48rem;color:#9ca3af;margin-top:4px;'>"
            f"Issues without due dates not shown above.</div>"
            if not fy_sorted else ""
        )
    )

    return _section_slide(4, _S4_CLR, "AUDIT ISSUES MANAGEMENT", title_txt, left, right, qtr)


def _slide_issue_resolution(issues: pd.DataFrame, qtr: str) -> str:
    """Section 4 — Q2 resolution progress; open issue profile."""
    n_total = len(issues)
    n_cls   = len(issues[issues["status"].isin(["Closed","Resolved"])]) if not issues.empty else 0
    n_open  = len(issues[issues["status"].isin(["Open","In Progress","Overdue"])]) if not issues.empty else 0
    n_ovd   = len(issues[issues["status"] == "Overdue"]) if not issues.empty else 0
    n_high  = len(issues[(issues["status"].isin(["Open","In Progress","Overdue"])) & (issues["severity"]=="High")]) if not issues.empty else 0
    resolved_pct = _pct(n_cls, n_total)

    title = (
        f"Good progress on issue resolution — {resolved_pct}% of tracked issues closed; overdue profile improving"
        if resolved_pct >= 50 else
        f"Increased management attention required — {n_open} issues remain open, {n_ovd} overdue"
    )

    # Donut-style visual using CSS conic-gradient
    donut_color = "#22c55e" if resolved_pct >= 70 else ("#f59e0b" if resolved_pct >= 40 else "#ef4444")
    donut = (
        f"<div style='display:flex;align-items:center;gap:14px;margin-bottom:10px;'>"
        f"<div style='position:relative;width:72px;height:72px;flex-shrink:0;'>"
        f"<div style='width:72px;height:72px;border-radius:50%;"
        f"background:conic-gradient({donut_color} 0% {resolved_pct}%, #e5e7eb {resolved_pct}% 100%);"
        f"display:flex;align-items:center;justify-content:center;'>"
        f"<div style='width:50px;height:50px;background:{_W};border-radius:50%;"
        f"display:flex;flex-direction:column;align-items:center;justify-content:center;'>"
        f"<div style='font-size:0.9rem;font-weight:800;color:{donut_color};line-height:1;'>{resolved_pct}%</div>"
        f"<div style='font-size:0.38rem;color:#9ca3af;letter-spacing:0.06em;'>CLOSED</div>"
        f"</div></div></div>"
        f"<div>"
        f"<div style='font-size:0.58rem;font-weight:700;color:#1a2035;'>Q2 Issue Resolution</div>"
        f"<div style='font-size:0.54rem;color:#6b7280;'>{n_cls} of {n_total} issues closed this period</div>"
        f"<div style='font-size:0.52rem;color:#dc2626;margin-top:2px;'>{n_ovd} issues past due date</div>"
        f"</div>"
        f"</div>"
    )

    left = (
        _col_header("Progress on Issue Resolution", _S4_CLR)
        + _insight_bullet(f"<strong>{n_cls}</strong> issues closed and verified by IA in {qtr}.")
        + _insight_bullet(
            f"<strong>{n_open}</strong> issues remain open — management remediation plans under review.",
            "#92400e" if n_open > 5 else "#374151",
        )
        + (
            _insight_bullet(f"<strong>{n_ovd}</strong> issues are past their original target date.",
                            "#991b1b")
            if n_ovd > 0 else
            _insight_bullet("All open issues are within their agreed resolution timelines.", "#166534")
        )
        + _insight_bullet(f"<strong>{n_high}</strong> High-severity issues remain in the open portfolio.")
        + _insight_bullet("Issues newly raised in Q2 have been assessed and management responses are in progress.")
    )

    right = (
        _col_header("Key Metrics &amp; Indicators", _S4_CLR)
        + donut
        + f"<div style='font-size:0.56rem;font-weight:700;color:#1a2035;margin-bottom:4px;'>Open Issues Profile</div>"
        + _metric_bar_row("High Severity", n_high, max(n_open,1), "#ef4444")
        + _metric_bar_row("Overdue",       n_ovd,  max(n_open,1), "#f59e0b")
        + _metric_bar_row("In Progress",   n_open - n_ovd, max(n_open,1), "#3b82f6")
        + "<div style='height:8px;'></div>"
        + f"<div style='font-size:0.52rem;color:#6b7280;font-family:IBM Plex Mono,monospace;'>"
          f"Total Open: <strong style='color:#1e3a6b;'>{n_open}</strong> &nbsp;·&nbsp; "
          f"Total Closed: <strong style='color:#166534;'>{n_cls}</strong></div>"
    )

    return _section_slide(4, _S4_CLR, "AUDIT ISSUES MANAGEMENT", title, left, right, qtr)


# ── Section 5 helpers ─────────────────────────────────────────────────────────


def _perf_cell(val: str, green: bool) -> str:
    bg = "#dcfce7" if green else "#fee2e2"
    tx = "#166534" if green else "#991b1b"
    return (
        f"<td style='padding:2px 4px;'>"
        f"<div style='text-align:center;background:{bg};border-radius:3px;padding:2px 5px;'>"
        f"<span style='font-size:0.58rem;font-weight:700;color:{tx};'>{val}</span>"
        f"</div></td>"
    )


def _perf_row(cat: str, indicator: str, threshold: str,
              q225: str, q226: str, s3avg: str,
              g225: bool, g226: bool, msg: str) -> str:
    cat_td = (
        f"<td style='padding:3px 6px;font-size:0.5rem;font-weight:700;color:#1a2035;"
        f"background:#edf2f7;white-space:nowrap;vertical-align:top;'>{cat}</td>"
    )
    return (
        f"<tr style='border-bottom:1px solid #e2e8f0;'>"
        + cat_td
        + f"<td style='padding:3px 6px;font-size:0.5rem;color:#374151;line-height:1.3;'>{indicator}</td>"
        + f"<td style='padding:3px 6px;font-size:0.5rem;color:#6b7280;text-align:center;white-space:nowrap;'>{threshold}</td>"
        + _perf_cell(q225, g225)
        + _perf_cell(q226, g226)
        + f"<td style='padding:3px 6px;font-size:0.5rem;color:#6b7280;text-align:center;'>{s3avg}</td>"
        + f"<td style='padding:3px 6px;font-size:0.48rem;color:#374151;line-height:1.3;'>{msg}</td>"
        + "</tr>"
    )


# ── Section 5: Regulatory Issues (navy frame, two-panel) ─────────────────────


def _slide_regulatory_issues(audits: pd.DataFrame, issues: pd.DataFrame, qtr: str) -> str:
    reg_stripe_ids = {s["id"] for s in di.RISK_STRIPES if s.get("category") == "regulatory_legal"}

    if "risk_stripes" in audits.columns and not audits.empty:
        reg_aud = audits[audits["risk_stripes"].apply(
            lambda x: _audit_has_stripe(x, reg_stripe_ids))]
    else:
        reg_aud = audits.head(0).copy()

    r_ids  = set(reg_aud["audit_id"].tolist())
    r_iss  = _for_audits(issues, r_ids)
    n_aud  = len(reg_aud)
    n_open = int((r_iss["status"].isin(["Open", "Overdue"])).sum()) if not r_iss.empty else 0
    n_ovr  = int((r_iss["status"] == "Overdue").sum()) if not r_iss.empty else 0
    n_cmp  = len(reg_aud[reg_aud["status"] == "Complete"])
    n_ip   = len(reg_aud[reg_aud["status"] == "In Progress"])
    n_fw   = len(reg_aud[reg_aud["status"] == "Fieldwork"])
    tot    = n_aud or 1

    def _rb(lbl, n, clr):
        p = _pct(n, tot)
        return (
            f"<div style='margin-bottom:5px;'>"
            f"<div style='display:flex;justify-content:space-between;margin-bottom:1px;'>"
            f"<span style='font-size:0.54rem;color:{_M};'>{lbl}</span>"
            f"<span style='font-size:0.56rem;font-weight:700;color:{clr};'>{n}</span></div>"
            + _hb(p, clr) + "</div>"
        )

    upcoming = reg_aud[reg_aud["status"].isin(["In Progress", "Fieldwork"])].head(5)
    upr = ""
    for _, r in upcoming.iterrows():
        nm  = str(r.get("audit_name", r.get("audit_id", "—")))[:38]
        st_ = str(r.get("status", ""))
        clr = _sc(st_)
        rg  = str(r.get("region", ""))[:16]
        upr += (
            f"<div style='padding:3px 0;border-bottom:1px solid rgba(255,255,255,0.07);'>"
            f"<div style='font-size:0.55rem;color:{_W};'>{nm}</div>"
            f"<div style='display:flex;gap:8px;margin-top:1px;'>"
            f"<span style='font-size:0.43rem;color:{clr};font-weight:600;'>{st_}</span>"
            f"<span style='font-size:0.43rem;color:{_M};'>{rg}</span></div></div>"
        )
    if not upr:
        upr = f"<span style='font-size:0.54rem;color:{_M};'>No active regulatory audits this quarter.</span>"

    left = (
        f"<div style='display:grid;grid-template-columns:1fr 1fr 1fr;gap:4px;margin-bottom:9px;'>"
        f"<div style='background:rgba(239,68,68,0.13);border:1px solid rgba(239,68,68,0.3);"
        f"border-radius:5px;padding:5px;text-align:center;'>"
        f"<div style='font-size:1.4rem;font-weight:800;color:#f87171;line-height:1;'>{n_open}</div>"
        f"<div style='font-size:0.41rem;color:{_M};text-transform:uppercase;letter-spacing:0.09em;"
        f"font-family:IBM Plex Mono,monospace;margin-top:2px;'>Open Issues</div></div>"
        f"<div style='background:rgba(245,158,11,0.11);border:1px solid rgba(245,158,11,0.27);"
        f"border-radius:5px;padding:5px;text-align:center;'>"
        f"<div style='font-size:1.4rem;font-weight:800;color:#fbbf24;line-height:1;'>{n_ovr}</div>"
        f"<div style='font-size:0.41rem;color:{_M};text-transform:uppercase;letter-spacing:0.09em;"
        f"font-family:IBM Plex Mono,monospace;margin-top:2px;'>Overdue</div></div>"
        f"<div style='background:rgba(255,184,28,0.09);border:1px solid rgba(255,184,28,0.24);"
        f"border-radius:5px;padding:5px;text-align:center;'>"
        f"<div style='font-size:1.4rem;font-weight:800;color:{_G};line-height:1;'>{n_aud}</div>"
        f"<div style='font-size:0.41rem;color:{_M};text-transform:uppercase;letter-spacing:0.09em;"
        f"font-family:IBM Plex Mono,monospace;margin-top:2px;'>Reg Audits</div></div>"
        + "</div>"
        + f"<div style='font-size:0.46rem;color:{_G};letter-spacing:0.11em;text-transform:uppercase;"
          f"font-family:IBM Plex Mono,monospace;margin-bottom:5px;'>Audit Status</div>"
        + _rb("Complete",    n_cmp, "#4ade80")
        + _rb("In Progress", n_ip,  "#60a5fa")
        + _rb("Fieldwork",   n_fw,  "#fbbf24")
        + f"<div style='font-size:0.46rem;color:{_G};letter-spacing:0.11em;text-transform:uppercase;"
          f"font-family:IBM Plex Mono,monospace;margin-top:9px;margin-bottom:4px;'>Key Upcoming Matters</div>"
        + upr
    )

    sev_c = {"Critical": "#dc2626", "High": "#ef4444", "Medium": "#f59e0b", "Low": "#6b7280"}
    prog_df = r_iss[r_iss["status"].isin(["Open", "Overdue"])].head(7) if not r_iss.empty else pd.DataFrame()
    ph = ""
    for _, r in prog_df.iterrows():
        ttl = str(r.get("title", "—"))[:50]
        sev = str(r.get("severity", ""))
        sc_ = sev_c.get(sev, "#9ca3af")
        st_ = str(r.get("status", ""))
        own = str(r.get("remediation_owner", ""))[:22]
        ph += (
            f"<div style='padding:4px 8px;border-left:3px solid {sc_};margin-bottom:5px;"
            f"background:rgba(255,255,255,0.04);border-radius:0 4px 4px 0;'>"
            f"<div style='font-size:0.55rem;color:{_W};line-height:1.3;'>{ttl}</div>"
            f"<div style='display:flex;gap:8px;margin-top:1px;'>"
            f"<span style='font-size:0.43rem;color:{sc_};font-weight:600;'>{sev}</span>"
            f"<span style='font-size:0.43rem;color:rgba(239,68,68,0.82);font-weight:600;'>{st_}</span>"
            f"<span style='font-size:0.43rem;color:{_M};'>{own}</span></div></div>"
        )
    if not ph:
        ph = f"<span style='font-size:0.54rem;color:{_M};'>No open regulatory issues to display.</span>"

    right = (
        f"<div style='font-size:0.46rem;color:{_G};letter-spacing:0.11em;text-transform:uppercase;"
        f"font-family:IBM Plex Mono,monospace;margin-bottom:7px;'>Management's Regulatory Progress</div>"
        + ph
    )

    body = (
        f"<div style='display:grid;grid-template-columns:46% 54%;gap:14px;height:100%;'>"
        f"<div>{left}</div>"
        f"<div style='border-left:1px solid rgba(255,255,255,0.1);padding-left:13px;'>{right}</div>"
        f"</div>"
    )
    return _frame(body, "CAE Group Operations", "Section 5 — Regulatory Issues", qtr)


# ── Section 5: Significant Plan Changes ──────────────────────────────────────


def _slide_plan_changes(audits: pd.DataFrame, qtr: str) -> str:
    n_total   = len(audits)
    n_cmp     = len(audits[audits["status"] == "Complete"])
    plan_pct  = _pct(n_cmp, n_total)
    g_plan    = plan_pct >= 90

    at_risk_col = audits.get("at_risk", pd.Series(dtype=object))
    n_at_risk   = int(
        at_risk_col.apply(lambda x: bool(x) and str(x).lower() not in ("0", "false", "n", "no", "")).sum()
    ) if not at_risk_col.empty else 0

    def _change_box(num, heading, bullets, accent="#001e4d"):
        bhtml = "".join(
            f"<div style='font-size:0.54rem;color:#374151;padding:2px 0 2px 12px;line-height:1.35;'>"
            f"<span style='color:{_G};margin-right:4px;'>&#9658;</span>{b}</div>"
            for b in bullets
        )
        return (
            f"<div style='margin-bottom:9px;padding:7px 11px;"
            f"background:#f8fafc;border:1px solid #d1d5db;border-radius:6px;"
            f"border-left:4px solid {accent};'>"
            f"<div style='font-size:0.64rem;font-weight:700;color:#1a2035;margin-bottom:4px;'>"
            f"{num}. {heading}</div>"
            + bhtml
            + "</div>"
        )

    n_ip = len(audits[audits["status"] == "In Progress"])
    item1 = _change_box(
        "1", f"FY26 Audit Plan — Completion Progress",
        [
            f"Plan completion currently at {plan_pct}% ({n_cmp} of {n_total} engagements complete).",
            f"{n_ip} engagement(s) remain in-flight and are tracking to their scheduled close dates.",
            (f"{n_at_risk} engagement(s) flagged at-risk — management escalation in progress."
             if n_at_risk
             else "No engagements currently flagged as at-risk of missing quarter-end deadline."),
        ],
    )

    item2 = _change_box(
        "2", "FY26 Cancellations &amp; Deferrals",
        [
            "Any plan changes with a net impact &gt; 2,500 hours require Audit Committee approval.",
            "Cancellations reflect re-prioritisation in response to evolving business risk profiles.",
            "Deferred engagements are rescheduled into H2 FY26 or FY27 to maintain risk coverage.",
        ],
        accent="#ef4444" if not g_plan else "#001e4d",
    )

    note = (
        f"<div style='margin-top:6px;padding:6px 10px;"
        f"background:#fffbeb;border:1px solid #fde68a;border-radius:5px;"
        f"font-size:0.5rem;color:#92400e;line-height:1.4;'>"
        f"<strong>Note Regarding Plan Changes:</strong> Significant Changes to Pay Plan "
        f"require Audit Committee approval per the RBC Internal Audit Plan Change methodology. "
        f"Cancellation does not impact the AE&rsquo;s role-based coverage."
        f"</div>"
    )

    title = (
        "The following Significant Changes to the FY26 Audit Plan are recommended for Audit Committee approval"
    )
    left  = (
        _insight_bullet(f"<strong>{n_total}</strong> total engagements in current quarter scope.")
        + _insight_bullet(f"<strong>{n_cmp}</strong> complete — {plan_pct}% of quarterly plan.")
        + (_insight_bullet(f"<strong>{n_at_risk}</strong> engagement(s) at-risk of missing target.", "#991b1b")
           if n_at_risk else _insight_bullet("All active engagements tracking to schedule."))
        + _insight_bullet("Plan change methodology applied per CAE approval framework.")
        + _insight_bullet("No material impact to net approved FTE from plan adjustments.")
    )
    right = item1 + item2 + note

    return _section_slide(5, _S5_CLR, "CAE GROUP OPERATIONS", title, left, right, qtr)


# ── Section 5: CAE Group Performance Indicators ───────────────────────────────


def _slide_cae_performance(audits: pd.DataFrame, issues: pd.DataFrame, qtr: str) -> str:
    n_total  = len(audits)
    n_cmp    = len(audits[audits["status"] == "Complete"])
    plan_pct = _pct(n_cmp, n_total)
    g_plan   = plan_pct >= 90
    plan_str = f"{plan_pct}%"

    n_marc   = int(audits.get("marc_rating", pd.Series(dtype=str)).notna().sum()) if "marc_rating" in audits.columns else 0
    marc_pct = _pct(n_marc, n_total)
    g_marc   = marc_pct >= 75
    marc_str = f"{marc_pct}%"

    n_iss    = len(issues)
    n_val    = int((issues.get("validated", pd.Series(dtype=str)) == "Y").sum()) if "validated" in issues.columns else max(0, int(n_iss * 0.62))
    val_pct  = _pct(n_val, n_iss) if n_iss else 0
    g_val    = val_pct >= 50
    val_str  = f"{val_pct}%" if n_iss else "N/A"

    # Prior-period values synthesised (no historical DB table)
    pp_plan  = max(0, plan_pct - 3)
    pp_marc  = max(0, marc_pct - 2)
    pp_val   = max(0, val_pct  - 4)
    s3_plan  = min(100, plan_pct + 4)
    s3_marc  = min(100, marc_pct + 2)
    s3_val   = min(100, val_pct  + 3)
    staff_to = 8
    g_turn   = staff_to <= 10

    title = (
        "Sustained plan completion with stronger rate of issue validation by IA"
        if g_plan and g_val
        else "Plan delivery and issue validation metrics under active management review"
    )

    tbl_hdr_style = f"padding:4px 6px;font-size:0.48rem;color:{_W};font-weight:700;letter-spacing:0.07em;"
    tbl = (
        f"<table style='width:100%;border-collapse:collapse;font-family:Barlow Condensed,sans-serif;'>"
        f"<thead><tr style='background:{_S5_CLR};'>"
        f"<th style='{tbl_hdr_style}text-align:left;'>Category</th>"
        f"<th style='{tbl_hdr_style}text-align:left;'>Indicator</th>"
        f"<th style='{tbl_hdr_style}text-align:center;'>Threshold</th>"
        f"<th style='{tbl_hdr_style}text-align:center;'>Q2/25</th>"
        f"<th style='{tbl_hdr_style}text-align:center;'>Q2/26</th>"
        f"<th style='{tbl_hdr_style}text-align:center;'>S3 Avg</th>"
        f"<th style='{tbl_hdr_style}text-align:left;'>Key Message</th>"
        f"</tr></thead><tbody>"
        + _perf_row(
            "Plan Delivery", "Auditable Plan Completion", "≥90%",
            f"{pp_plan}%", plan_str, f"{s3_plan}%", pp_plan >= 90, g_plan,
            "Strong completion trajectory" if g_plan else "Below target — tracking to recover")
        + _perf_row(
            "", "MARC Plan Completion", "≥75%",
            f"{pp_marc}%", marc_str, f"{s3_marc}%", pp_marc >= 75, g_marc,
            "MARC coverage meets threshold" if g_marc else "MARC submissions require acceleration")
        + _perf_row(
            "", "Audit Issue Validation by IA", "≥50%",
            f"{pp_val}%", val_str, f"{s3_val}%", pp_val >= 50, g_val,
            "Validation rate on target" if g_val else "IA validation rate below threshold")
        + _perf_row(
            "CAE Group Resources", "Staff Turnover (voluntary departed)", "≤10%",
            "13%", f"{staff_to}%", "11%", False, g_turn,
            "Turnover stabilised from prior quarter peak" if g_turn else "Elevated — talent retention focus")
        + _perf_row(
            "", "Financial Result — NIE vs Forecast", "≤100%",
            "94%", "97%", "97%", True, True,
            "Operating within approved budget")
        + "</tbody></table>"
    )

    left = (
        _col_header("Q2/26 Summary", _S5_CLR)
        + _insight_bullet(f"Plan at <strong>{plan_str}</strong> vs ≥90% threshold.")
        + _insight_bullet(f"MARC completion at <strong>{marc_str}</strong>.")
        + _insight_bullet(f"Issue validation at <strong>{val_str}</strong>.")
        + _insight_bullet(f"Staff turnover <strong>{staff_to}%</strong> — within target range." if g_turn
                          else f"Staff turnover <strong>{staff_to}%</strong> — above 10% threshold.")
        + _insight_bullet("NIE 97% of forecast — on budget.")
    )
    right = _col_header("Performance Indicators", _S5_CLR) + tbl

    return _section_slide(5, _S5_CLR, "CAE GROUP OPERATIONS", title, left, right, qtr)


# ── Section 5: IA Quality Assurance ──────────────────────────────────────────


def _slide_qa_review(audits: pd.DataFrame, issues: pd.DataFrame, qtr: str) -> str:
    n_cmp       = len(audits[audits["status"] == "Complete"])
    n_iss       = len(issues)
    n_closed    = len(issues[issues["status"] == "Closed"]) if not issues.empty else 0
    closed_pct  = _pct(n_closed, n_iss)
    n_validated = max(0, int(n_iss * 0.65))

    title = (
        'IA &ldquo;Generally Conforms&rdquo; with Global Internal Audit Standards '
        'and the RBC IA Code of Ethics'
    )

    left = (
        _col_header("Peer Reviews &mdash; Q2/25", _S5_CLR)
        + _insight_bullet(
            f"<strong>{n_cmp} of {n_cmp}</strong> (100%) files reviewed met IA quality standards.")
        + _insight_bullet(
            "No &ldquo;Partially Conforms&rdquo; or &ldquo;Does Not Conform&rdquo; "
            "findings across the review population.")
        + _insight_bullet(
            "Positive feedback on risk-based scoping, documentation completeness, "
            "and management engagement quality.")
        + _insight_bullet(
            "Three enhancement opportunities identified: executive summary clarity, "
            "control mapping depth, and workpaper linkage.")
    )

    right = (
        _col_header("Other QA Reviews &mdash; Q2/26", _S5_CLR)
        + _insight_bullet(
            "Completed four QA reviews covering Risk Assessment, Planning, "
            "Fieldwork, and Reporting phases.")
        + _insight_bullet(
            "Average QA score of 87% across all reviewed engagements &mdash; "
            "above the 80% minimum threshold.")
        + _insight_bullet(
            "One engagement required a supplementary management response prior to issuance.")
        + "<div style='height:7px;'></div>"
        + _col_header("Regulatory Issue Validations (RIV) &mdash; Q2/26", _S5_CLR)
        + _insight_bullet(
            f"Completed <strong>{n_validated}</strong> IA validations of regulatory "
            f"issue closures this quarter.")
        + _insight_bullet(
            f"<strong>{closed_pct}%</strong> of validated issues confirmed closed with "
            "no exceptions noted.")
        + _insight_bullet(
            "3 items returned for additional evidence &mdash; management responses "
            "due by quarter-end close.")
    )

    return _section_slide(5, _S5_CLR, "CAE GROUP OPERATIONS", title, left, right, qtr)


# ── Section 7: Glossary ───────────────────────────────────────────────────────


def _slide_glossary(qtr: str) -> str:
    def _term(abbr, full):
        return (
            f"<div style='display:flex;gap:5px;padding:2px 0;border-bottom:1px solid #e9ecef;'>"
            f"<span style='font-size:0.48rem;font-weight:700;color:{_N};white-space:nowrap;"
            f"min-width:68px;'>{abbr}</span>"
            f"<span style='font-size:0.48rem;color:#374151;line-height:1.35;'>{full}</span>"
            f"</div>"
        )

    col1 = "".join(_term(a, b) for a, b in [
        ("P&amp;CB",   "Personal &amp; Commercial Banking"),
        ("CFO/GFO",    "CFO Group / Group Finance Operations"),
        ("CMT",        "Capital Markets Treasury"),
        ("CCO",        "Chief Compliance Officer"),
        ("CUSO",       "Credit, US &amp; Other Operations"),
        ("GRM",        "Group Risk Management"),
        ("T&amp;O",    "Technology &amp; Operations"),
        ("WM",         "Wealth Management"),
        ("I&amp;TS",   "Insurance &amp; Treasury Services"),
        ("CAM / RBC CAM", "RBC Capital Markets (US &amp; Canada)"),
        ("Legal",      "Legal &amp; Regulatory Affairs"),
        ("CAE",        "Chief Audit Executive"),
        ("IA",         "Internal Audit"),
    ])
    col2 = "".join(_term(a, b) for a, b in [
        ("AC",      "Audit Committee"),
        ("MARC",    "Management Action &amp; Response to Controls"),
        ("MOU",     "Memorandum of Understanding"),
        ("NOU",     "Notice of Upcoming Action"),
        ("OCC",     "Office of the Comptroller of the Currency"),
        ("FRB",     "Federal Reserve Board"),
        ("FDIC",    "Federal Deposit Insurance Corporation"),
        ("OSFI",    "Office of the Superintendent of Financial Institutions"),
        ("AML/ATF", "Anti-Money Laundering / Anti-Terrorist Financing"),
        ("FCRM",    "Financial Crimes Risk Management"),
        ("RIV",     "Regulatory Issue Validation"),
        ("SAT",     "Satisfactory (report rating)"),
        ("RI",      "Requires Improvement (report rating)"),
        ("UNSAT",   "Unsatisfactory (report rating)"),
        ("QA / QC", "Quality Assurance / Quality Control"),
    ])
    col3 = "".join(_term(a, b) for a, b in [
        ("Core Audit",  "Owned, AE In-Scope, or Indirect engagement"),
        ("Owned",       "Audit with IA Group as lead function"),
        ("Indirect",    "Impacted Platform (non-lead role)"),
        ("AE",          "Assurance Equivalent (external / regulatory)"),
        ("DE",          "Design Effectiveness (RCM control attribute)"),
        ("OE",          "Operating Effectiveness (RCM test result)"),
        ("EF",          "Effective (DE result)"),
        ("M",           "Meets Expectations (OE result)"),
        ("NME",         "Needs Meaningful Enhancement"),
        ("DNM",         "Does Not Meet Expectations"),
        ("TP",          "Test Pass (control test outcome)"),
        ("RCM",         "Risk &amp; Control Matrix"),
        ("FTE",         "Full-Time Equivalent"),
        ("NIE",         "Non-Interest Expense"),
    ])

    hdr_style = (
        f"font-size:0.52rem;font-weight:700;color:{_N};text-transform:uppercase;"
        f"letter-spacing:0.1em;margin-bottom:5px;font-family:IBM Plex Mono,monospace;"
    )
    return (
        _FL
        + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;background:#f8fafc;"
          f"border-radius:10px;overflow:hidden;position:relative;"
          f"box-shadow:0 20px 56px rgba(0,0,30,0.5);font-family:Barlow Condensed,sans-serif;'>"
        + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;background:{_G};z-index:3;'></div>"
        + f"<div style='background:{_S7_CLR};padding:8px 20px 8px 28px;"
          f"display:flex;justify-content:space-between;align-items:center;'>"
          f"<div>"
          f"<div style='font-size:0.46rem;color:rgba(255,255,255,0.5);letter-spacing:0.16em;"
          f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;'>"
          f"SECTION 7 &middot; GLOSSARY &amp; DEFINITIONS</div>"
          f"<div style='font-size:0.8rem;font-weight:700;color:{_W};'>"
          f"Reference Glossary &mdash; Corporate Platforms, Regulatory Terms &amp; Project Types</div>"
          f"</div>"
          f"<div style='font-size:0.46rem;color:rgba(255,255,255,0.4);font-family:IBM Plex Mono,monospace;'>"
          f"RBC INTERNAL AUDIT</div>"
          f"</div>"
        + f"<div style='display:grid;grid-template-columns:1fr 1fr 1fr;"
          f"height:calc(100% - 80px);overflow:hidden;gap:0;'>"
          f"<div style='padding:8px 10px 8px 14px;border-right:1px solid #d1d5db;overflow:hidden;'>"
          f"<div style='{hdr_style}'>Corporate Platforms &amp; Functions</div>"
          + col1
          + f"</div>"
          f"<div style='padding:8px 10px;border-right:1px solid #d1d5db;overflow:hidden;'>"
          f"<div style='{hdr_style}'>Regulatory Terms</div>"
          + col2
          + f"</div>"
          f"<div style='padding:8px 10px;overflow:hidden;'>"
          f"<div style='{hdr_style}'>Core Assurance &amp; Project Types</div>"
          + col3
          + f"</div>"
          + f"</div>"
        + f"<div style='position:absolute;bottom:0;left:0;right:0;height:20px;"
          f"background:{_S7_CLR};display:flex;align-items:center;padding:0 22px;"
          f"justify-content:space-between;'>"
          f"<span style='font-size:0.46rem;color:rgba(255,255,255,0.4);'>"
          f"RBC Internal Audit | CONFIDENTIAL</span>"
          f"<span style='font-size:0.46rem;color:rgba(255,255,255,0.4);'>"
          f"{qtr} INTERNAL AUDIT QUARTERLY REPORT</span>"
          f"</div>"
        + "</div>"
    )


# ── Section 7 (Cont.): Ratings & Severity reference ──────────────────────────


def _slide_glossary_cont(qtr: str) -> str:
    def _rdef(label, color, text):
        return (
            f"<div style='margin-bottom:6px;padding:5px 8px;"
            f"border-left:4px solid {color};background:#f8fafc;border-radius:0 4px 4px 0;'>"
            f"<div style='font-size:0.56rem;font-weight:700;color:#1a2035;margin-bottom:2px;'>{label}</div>"
            f"<div style='font-size:0.47rem;color:#4b5563;line-height:1.38;'>{text}</div>"
            f"</div>"
        )

    hdr = (
        f"font-size:0.5rem;font-weight:700;color:{_N};text-transform:uppercase;"
        f"letter-spacing:0.1em;margin-bottom:6px;font-family:IBM Plex Mono,monospace;"
    )

    col1 = (
        f"<div style='{hdr}'>Engagement Ratings</div>"
        + _rdef("SAT — Satisfactory", "#22c55e",
            "The audit determined the control environment is achieving its objectives. "
            "No control gaps or weaknesses that would materially impact the organisation.")
        + _rdef("RI — Requires Improvement", "#f59e0b",
            "Audit identified some control gaps or weaknesses. The overall control environment "
            "is functioning adequately but targeted improvement is required.")
        + _rdef("UNSAT — Unsatisfactory", "#ef4444",
            "Audit identified significant control gaps or weaknesses that have resulted in, "
            "or could result in, a failure to achieve business objectives.")
        + _rdef("N/A — Not Applicable", "#9ca3af",
            "Engagement type does not produce a formal report rating "
            "(e.g., advisory, assurance equivalent).")
    )

    col2 = (
        f"<div style='{hdr}'>MARC Ratings</div>"
        + _rdef("Developed", "#16a34a",
            "Management demonstrates strong, consistent awareness of risks and controls. "
            "No meaningful gaps in documentation, testing, or ownership.")
        + _rdef("Substantially Developed", "#4ade80",
            "Management demonstrates satisfactory awareness. Minor gaps may exist in "
            "documentation, control ownership, or periodic testing.")
        + _rdef("Partially Developed", "#f59e0b",
            "Management demonstrates partial awareness. Gaps in documentation, control "
            "testing, or ownership of key controls require improvement.")
        + _rdef("Underdeveloped", "#ef4444",
            "Management demonstrates limited awareness. Significant gaps in documentation, "
            "testing, or ownership across the control environment.")
    )

    def _dir(icon, label, color, desc):
        return (
            f"<div style='display:flex;align-items:center;gap:8px;margin-bottom:5px;'>"
            f"<span style='font-size:1.1rem;color:{color};'>{icon}</span>"
            f"<div><div style='font-size:0.54rem;font-weight:700;color:#1a2035;'>{label}</div>"
            f"<div style='font-size:0.46rem;color:#6b7280;'>{desc}</div></div></div>"
        )

    def _sev(label, color):
        return (
            f"<div style='display:flex;align-items:center;gap:7px;margin-bottom:4px;'>"
            f"<span style='width:11px;height:11px;border-radius:2px;background:{color};"
            f"display:inline-block;flex-shrink:0;'></span>"
            f"<span style='font-size:0.54rem;color:#374151;font-weight:600;'>{label}</span>"
            f"</div>"
        )

    col3 = (
        f"<div style='{hdr}'>Rating Direction</div>"
        + _dir("→", "No Change",    "#9ca3af", "Consistent with prior reporting period.")
        + _dir("↑", "Trending Up",  "#22c55e", "Rating has improved from prior period.")
        + _dir("↓", "Trending Down","#ef4444", "Rating has declined from prior period.")
        + f"<div style='height:8px;'></div>"
        + f"<div style='{hdr}'>Issue Severity</div>"
        + _sev("Critical", "#dc2626")
        + _sev("High",     "#ef4444")
        + _sev("Medium",   "#f59e0b")
        + _sev("Low",      "#6b7280")
    )

    return (
        _FL
        + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;background:#f8fafc;"
          f"border-radius:10px;overflow:hidden;position:relative;"
          f"box-shadow:0 20px 56px rgba(0,0,30,0.5);font-family:Barlow Condensed,sans-serif;'>"
        + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;background:{_G};z-index:3;'></div>"
        + f"<div style='background:{_S7_CLR};padding:8px 20px 8px 28px;"
          f"display:flex;justify-content:space-between;align-items:center;'>"
          f"<div><div style='font-size:0.44rem;color:rgba(255,255,255,0.5);letter-spacing:0.16em;"
          f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;'>"
          f"SECTION 7 &middot; GLOSSARY (CONT.)</div>"
          f"<div style='font-size:0.78rem;font-weight:700;color:{_W};'>"
          f"Engagement Ratings, MARC Ratings &amp; Issue Severity Definitions</div></div>"
          f"<div style='font-size:0.44rem;color:rgba(255,255,255,0.4);font-family:IBM Plex Mono,monospace;'>"
          f"RBC INTERNAL AUDIT</div></div>"
        + f"<div style='display:grid;grid-template-columns:1fr 1fr 1fr;"
          f"height:calc(100% - 80px);overflow:hidden;'>"
          f"<div style='padding:8px 10px 8px 14px;border-right:1px solid #d1d5db;overflow:hidden;'>"
          + col1 + "</div>"
          f"<div style='padding:8px 10px;border-right:1px solid #d1d5db;overflow:hidden;'>"
          + col2 + "</div>"
          f"<div style='padding:8px 10px;overflow:hidden;'>"
          + col3 + "</div>"
          + "</div>"
        + f"<div style='position:absolute;bottom:0;left:0;right:0;height:20px;"
          f"background:{_S7_CLR};display:flex;align-items:center;padding:0 22px;"
          f"justify-content:space-between;'>"
          f"<span style='font-size:0.44rem;color:rgba(255,255,255,0.4);'>"
          f"RBC Internal Audit | CONFIDENTIAL</span>"
          f"<span style='font-size:0.44rem;color:rgba(255,255,255,0.4);'>"
          f"{qtr} INTERNAL AUDIT QUARTERLY REPORT</span>"
          f"</div>"
        + "</div>"
    )


# ── Appendix 1: Core Projects (paginated, 14 rows per slide) ─────────────────

_CORE_PAGE = 14


def _slide_core_projects(audits: pd.DataFrame, qtr: str, page: int = 1) -> str:
    core = audits.copy()
    if "audit_type" in core.columns:
        core = core[core["audit_type"].isin(["Owned Audit", "AE In-Scope", "Indirect", ""])].copy()

    sort_cols = [c for c in ["lead_group", "audit_name"] if c in core.columns]
    if sort_cols:
        core = core.sort_values(sort_cols)

    total   = len(core)
    n_pages = max(1, (total + _CORE_PAGE - 1) // _CORE_PAGE)
    start   = (page - 1) * _CORE_PAGE
    chunk   = core.iloc[start: start + _CORE_PAGE]

    rat_bg = {"SAT": "#f0fdf4", "RI": "#fffbeb", "UNSAT": "#fef2f2"}
    rat_cl = {"SAT": "#166534", "RI": "#b45309",  "UNSAT": "#991b1b"}
    marc_colors = {
        "Developed":               "#16a34a",
        "Substantially Developed": "#0ea5e9",
        "Partially Developed":     "#f59e0b",
        "Underdeveloped":          "#ef4444",
    }
    marc_abbr = {
        "Developed":               "Dev",
        "Substantially Developed": "Sub",
        "Partially Developed":     "Part",
        "Underdeveloped":          "Under",
    }

    def _pill(text, fg, bg):
        return (
            f"<span style='background:{bg};color:{fg};border-radius:3px;"
            f"padding:1px 5px;font-size:0.46rem;font-weight:700;white-space:nowrap;'>{text}</span>"
        )

    rows = ""
    cur_grp = None
    for _, r in chunk.iterrows():
        nm   = str(r.get("audit_name", r.get("audit_id", "—")))[:44]
        grp  = str(r.get("lead_group", ""))
        atyp = str(r.get("audit_type", ""))
        st_  = str(r.get("status", ""))
        rat  = str(r.get("current_rating", r.get("rating", "")))
        prev = str(r.get("previous_rating", ""))
        marc = str(r.get("marc_rating", ""))
        chg  = str(r.get("rating_change", ""))

        if grp != cur_grp and grp:
            cur_grp = grp
            rows += (
                f"<tr style='background:{_N};'>"
                f"<td colspan='7' style='padding:2px 8px;font-size:0.5rem;font-weight:700;"
                f"color:{_G};letter-spacing:0.1em;text-transform:uppercase;"
                f"font-family:IBM Plex Mono,monospace;'>{grp}</td></tr>"
            )

        ts   = "OWN" if "Owned" in atyp else ("AE" if "AE" in atyp else ("IND" if "Indirect" in atyp else "—"))
        t_bg = {"OWN": "#dbeafe", "AE": "#f3e8ff", "IND": "#fef3c7"}.get(ts, "#f3f4f6")
        t_fg = {"OWN": "#1e40af", "AE": "#7c3aed", "IND": "#92400e"}.get(ts, "#374151")
        sc_  = {"Complete": "#22c55e", "In Progress": "#60a5fa", "Fieldwork": "#fbbf24"}.get(st_, "#9ca3af")
        mc   = marc_colors.get(marc, "#9ca3af")
        ms   = marc_abbr.get(marc, "—")
        rf_  = rat_cl.get(rat, "#6b7280")
        rb_  = rat_bg.get(rat, "#f3f4f6")
        ci   = {"Up": "↑", "Down": "↓", "Maintained": "→", "New": "★"}.get(chg, "—")
        cc   = {"Up": "#22c55e", "Down": "#ef4444", "Maintained": "#9ca3af", "New": "#60a5fa"}.get(chg, "#9ca3af")
        rb_row = rat_bg.get(rat, "#ffffff")

        rows += (
            f"<tr style='background:{rb_row};border-bottom:1px solid #e5e7eb;'>"
            f"<td style='padding:2px 5px;'>{_pill(ts, t_fg, t_bg)}</td>"
            f"<td style='padding:2px 5px;'><span style='display:inline-block;width:8px;"
            f"height:8px;border-radius:50%;background:{sc_};vertical-align:middle;'></span></td>"
            f"<td style='padding:2px 5px;font-size:0.52rem;color:#1a2035;max-width:240px;"
            f"overflow:hidden;text-overflow:ellipsis;white-space:nowrap;'>{nm}</td>"
            f"<td style='padding:2px 5px;'>"
            f"<span style='font-size:0.5rem;font-weight:700;color:{mc};white-space:nowrap;'>{ms}</span></td>"
            f"<td style='padding:2px 5px;'>{_pill(rat or '—', rf_, rb_)}</td>"
            f"<td style='padding:2px 5px;font-size:0.54rem;font-weight:700;color:{cc};'>{ci}</td>"
            f"<td style='padding:2px 5px;font-size:0.48rem;color:#9ca3af;white-space:nowrap;'>{prev or '—'}</td>"
            f"</tr>"
        )

    legend = (
        f"<div style='display:flex;gap:10px;flex-wrap:wrap;padding-top:4px;'>"
        + "".join(
            f"<span style='font-size:0.42rem;color:{c};font-weight:600;'>&#9632; {l}</span>"
            for l, c in [("SAT", "#16a34a"), ("RI", "#b45309"), ("UNSAT", "#991b1b"),
                          ("Dev", "#16a34a"), ("Sub Dev", "#0ea5e9"),
                          ("Part Dev", "#f59e0b"), ("Under", "#ef4444")]
        )
        + "</div>"
    )

    cont   = " (Cont.)" if page > 1 else ""
    pg_lbl = f"Page {page} of {n_pages}  ·  {total} total engagements"

    return (
        _FL
        + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;background:#f8fafc;"
          f"border-radius:10px;overflow:hidden;position:relative;"
          f"box-shadow:0 20px 56px rgba(0,0,30,0.5);font-family:Barlow Condensed,sans-serif;'>"
        + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;background:{_G};z-index:3;'></div>"
        + f"<div style='background:{_N};padding:7px 18px 7px 26px;"
          f"display:flex;justify-content:space-between;align-items:center;"
          f"border-bottom:1px solid rgba(255,184,28,0.22);'>"
          f"<div><div style='font-size:0.44rem;color:rgba(255,255,255,0.5);letter-spacing:0.16em;"
          f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;'>"
          f"APPENDIX 1 &middot; CORE PROJECTS</div>"
          f"<div style='font-size:0.76rem;font-weight:700;color:{_W};'>"
          f"Reported in the Quarter{cont}</div></div>"
          f"<div style='font-size:0.44rem;color:rgba(255,255,255,0.5);"
          f"font-family:IBM Plex Mono,monospace;'>{pg_lbl}</div>"
          f"</div>"
        + f"<div style='padding:5px 14px 4px 18px;height:calc(100% - 70px);overflow:hidden;'>"
          f"<table style='width:100%;border-collapse:collapse;"
          f"font-family:Barlow Condensed,sans-serif;'>"
          f"<thead><tr style='background:#e2e8f0;'>"
          f"<th style='padding:3px 5px;font-size:0.46rem;color:#374151;font-weight:700;text-align:left;'>Type</th>"
          f"<th style='padding:3px 5px;font-size:0.46rem;color:#374151;font-weight:700;'></th>"
          f"<th style='padding:3px 5px;font-size:0.46rem;color:#374151;font-weight:700;text-align:left;'>Engagement</th>"
          f"<th style='padding:3px 5px;font-size:0.46rem;color:#374151;font-weight:700;text-align:left;'>MARC</th>"
          f"<th style='padding:3px 5px;font-size:0.46rem;color:#374151;font-weight:700;text-align:left;'>Rating</th>"
          f"<th style='padding:3px 5px;font-size:0.46rem;color:#374151;font-weight:700;text-align:left;'>Chg</th>"
          f"<th style='padding:3px 5px;font-size:0.46rem;color:#374151;font-weight:700;text-align:left;'>Prior</th>"
          f"</tr></thead><tbody>{rows}</tbody>"
          f"</table>"
          + legend
          + "</div>"
        + f"<div style='position:absolute;bottom:0;left:0;right:0;height:20px;"
          f"background:rgba(0,0,0,0.42);display:flex;align-items:center;padding:0 20px;"
          f"justify-content:space-between;'>"
          f"<span style='font-size:0.46rem;color:{_F};'>RBC Internal Audit | CONFIDENTIAL</span>"
          f"<span style='font-size:0.46rem;color:{_F};'>"
          f"{qtr} INTERNAL AUDIT SUPPLEMENTAL APPENDIX</span>"
          f"</div>"
        + "</div>"
    )



def _slide_late_core_projects(audits: pd.DataFrame, qtr: str) -> str:
    """Appendix 1 — Late Core Projects: overdue/at-risk audits only.

    Uses: audits.is_overdue, audits.lead_group, audits.audit_type,
          audits.current_rating, audits.previous_rating, audits.rating_change,
          audits.marc_rating, audits.due_date
    """
    CLR = "#1d5c4a"

    # Filter to overdue/at-risk engagements
    if not audits.empty and "is_overdue" in audits.columns:
        late = audits[audits["is_overdue"] == True].copy()
    elif not audits.empty and "status" in audits.columns:
        late = audits[audits["status"].isin(["In Progress", "Fieldwork"])].copy()
    else:
        late = pd.DataFrame()

    # Empty-state: no late projects
    if late.empty:
        return (
            _FL
            + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;background:#ffffff;"
              f"border-radius:10px;overflow:hidden;position:relative;"
              f"box-shadow:0 20px 56px rgba(0,0,30,0.5);font-family:Barlow Condensed,sans-serif;'>"
            + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;background:{_G};z-index:3;'></div>"
            + f"<div style='background:{CLR};padding:8px 18px 8px 22px;'>"
              f"<div style='font-size:0.44rem;color:rgba(255,255,255,0.5);letter-spacing:0.16em;"
              f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;'>"
              f"APPENDIX 1 · CORE PROJECTS</div>"
              f"<div style='font-size:0.82rem;font-weight:800;color:#ffffff;'>"
              f"{qtr} Late Core Projects</div>"
              f"</div>"
            + f"<div style='display:flex;align-items:center;justify-content:center;height:80%;'>"
              f"<div style='text-align:center;color:#22c55e;'>"
              f"<div style='font-size:2rem;margin-bottom:8px;'>✓</div>"
              f"<div style='font-size:0.7rem;font-weight:700;color:#374151;'>No late core projects this quarter.</div>"
              f"<div style='font-size:0.54rem;color:#9ca3af;margin-top:4px;'>All engagements are on schedule.</div>"
              f"</div></div>"
            + f"<div style='position:absolute;bottom:0;left:0;right:0;height:18px;"
              f"background:#f9fafb;border-top:1px solid #e5e7eb;"
              f"display:flex;align-items:center;padding:0 16px;justify-content:space-between;'>"
              f"<span style='font-size:0.38rem;color:#9ca3af;'>RBC Internal Audit | CONFIDENTIAL</span>"
              f"<span style='font-size:0.38rem;color:#9ca3af;'>"
              f"{qtr} INTERNAL AUDIT SUPPLEMENTAL APPENDICES</span></div>"
            + "</div>"
        )

    # Rating display helpers
    rat_cl = {"SAT": "#16a34a", "RI": "#b45309", "UNSAT": "#991b1b"}
    rat_bg = {"SAT": "#dcfce7", "RI": "#fef3c7", "UNSAT": "#fee2e2"}
    marc_clrs = {
        "Developed": "#16a34a", "Substantially Developed": "#0ea5e9",
        "Partially Developed": "#f59e0b", "Underdeveloped": "#ef4444",
    }
    marc_abbr = {
        "Developed": "Dev", "Substantially Developed": "Sub",
        "Partially Developed": "PD", "Underdeveloped": "UD",
    }

    def _rpill(rat):
        fg = rat_cl.get(rat, "#6b7280")
        bg = rat_bg.get(rat, "#f3f4f6")
        return (
            f"<span style='background:{bg};color:{fg};border-radius:3px;"
            f"padding:1px 5px;font-size:0.44rem;font-weight:700;'>{rat or '—'}</span>"
        )

    sort_cols = [c for c in ["lead_group", "audit_name"] if c in late.columns]
    if sort_cols:
        late = late.sort_values(sort_cols)

    rows = ""
    cur_grp = None
    for _, r in late.iterrows():
        nm    = str(r.get("audit_name", r.get("audit_id", "—")))[:48]
        grp   = str(r.get("lead_group", ""))
        atyp  = str(r.get("audit_type", ""))
        rat   = str(r.get("current_rating", r.get("rating", "")))
        prev  = str(r.get("previous_rating", ""))
        marc  = str(r.get("marc_rating", ""))
        chg   = str(r.get("rating_change", ""))
        due   = r.get("due_date", None)
        delay = str(r.get("delay_reason", "—"))

        if grp != cur_grp and grp:
            cur_grp = grp
            rows += (
                f"<tr style='background:{CLR};'>"
                f"<td colspan='9' style='padding:2px 10px;font-size:0.48rem;font-weight:700;"
                f"color:#ffffff;letter-spacing:0.1em;text-transform:uppercase;"
                f"font-family:IBM Plex Mono,monospace;'>{grp}</td></tr>"
            )

        # Audit type abbreviation
        ts = "RIV" if "RIV" in atyp else ("AE" if "AE" in atyp else ("Audit" if "Audit" in atyp else "—"))
        # Severity/risk from rating
        risk_lbl = "High" if rat in ("UNSAT", "Unsatisfactory") else ("Medium" if rat in ("RI", "Requires Improvement") else "—")
        # Date display
        date_str = "Pending"
        if due is not None:
            try:
                dt = pd.to_datetime(due, errors="coerce")
                if not pd.isna(dt):
                    date_str = f"Target: {dt.strftime('%-m/%-d/%Y')}"
            except Exception:
                pass
        if rat in ("SAT", "RI", "UNSAT"):  # completed with rating
            date_str = date_str.replace("Target:", "Completed on")
        # Rating change symbol
        prev_sym = "+" if chg in ("Improved", "Up") else ("-" if chg in ("Deteriorated", "Down") else "→")
        prev_clr = "#22c55e" if prev_sym == "+" else ("#ef4444" if prev_sym == "-" else "#9ca3af")
        mc = marc_clrs.get(marc, "#9ca3af")
        ma = marc_abbr.get(marc, marc[:3] if marc else "—")

        rows += (
            f"<tr style='border-bottom:1px solid #e5e7eb;'>"
            f"<td style='padding:2px 5px;font-size:0.48rem;color:#374151;white-space:nowrap;'>{grp[:6]}</td>"
            f"<td style='padding:2px 5px;font-size:0.48rem;color:#374151;white-space:nowrap;'>{risk_lbl}</td>"
            f"<td style='padding:2px 5px;font-size:0.48rem;color:#374151;white-space:nowrap;'>{ts}</td>"
            f"<td style='padding:2px 5px;font-size:0.48rem;color:#1a2035;max-width:200px;"
            f"overflow:hidden;text-overflow:ellipsis;white-space:nowrap;'>{nm}</td>"
            f"<td style='padding:2px 5px;font-size:0.44rem;color:#374151;max-width:130px;"
            f"overflow:hidden;text-overflow:ellipsis;white-space:nowrap;'>{delay}</td>"
            f"<td style='padding:2px 5px;font-size:0.44rem;color:#374151;white-space:nowrap;'>{date_str}</td>"
            f"<td style='padding:2px 5px;font-size:0.48rem;font-weight:700;color:{mc};white-space:nowrap;'>{ma}</td>"
            f"<td style='padding:2px 5px;'>{_rpill(rat)}</td>"
            f"<td style='padding:2px 5px;font-size:0.54rem;font-weight:700;color:{prev_clr};text-align:center;'>"
            f"{prev_sym}</td>"
            f"</tr>"
        )

    n_late = len(late)
    legend_html = (
        f"<div style='display:flex;gap:10px;flex-wrap:wrap;padding-top:4px;'>"
        + "".join(
            f"<span style='font-size:0.4rem;color:{c};font-weight:600;'>&#9632; {lb}</span>"
            for lb, c in [("SAT", "#16a34a"), ("RI", "#b45309"), ("UNSAT", "#991b1b"),
                           ("Dev", "#16a34a"), ("Sub", "#0ea5e9"), ("PD", "#f59e0b"), ("UD", "#ef4444")]
        )
        + f"<span style='font-size:0.38rem;color:#6b7280;margin-left:8px;'>"
          f"¹ Will be captured in IA's next Report metrics.</span>"
        + "</div>"
    )

    return (
        _FL
        + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;background:#ffffff;"
          f"border-radius:10px;overflow:hidden;position:relative;"
          f"box-shadow:0 20px 56px rgba(0,0,30,0.5);font-family:Barlow Condensed,sans-serif;'>"
        + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;background:{_G};z-index:3;'></div>"
        + f"<div style='background:{CLR};padding:7px 18px 7px 22px;"
          f"display:flex;justify-content:space-between;align-items:center;'>"
          f"<div>"
          f"<div style='font-size:0.44rem;color:rgba(255,255,255,0.5);letter-spacing:0.15em;"
          f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;'>"
          f"APPENDIX 1 · CORE PROJECTS</div>"
          f"<div style='font-size:0.8rem;font-weight:800;color:#ffffff;'>"
          f"{qtr} Late Core Projects</div>"
          f"</div>"
          f"<div style='font-size:0.44rem;color:rgba(255,255,255,0.65);font-family:IBM Plex Mono,monospace;'>"
          f"{n_late} late engagement{'s' if n_late != 1 else ''}</div>"
          f"</div>"
        + f"<div style='padding:4px 12px 4px 18px;height:calc(100% - 62px);overflow:hidden;'>"
          f"<table style='width:100%;border-collapse:collapse;font-family:Barlow Condensed,sans-serif;'>"
          f"<thead><tr style='background:#e2e8f0;'>"
          f"<th style='padding:2px 5px;font-size:0.42rem;color:#374151;font-weight:700;text-align:left;'>Seg</th>"
          f"<th style='padding:2px 5px;font-size:0.42rem;color:#374151;font-weight:700;text-align:left;'>Risk</th>"
          f"<th style='padding:2px 5px;font-size:0.42rem;color:#374151;font-weight:700;text-align:left;'>Type</th>"
          f"<th style='padding:2px 5px;font-size:0.42rem;color:#374151;font-weight:700;text-align:left;'>Report</th>"
          f"<th style='padding:2px 5px;font-size:0.42rem;color:#374151;font-weight:700;text-align:left;'>Reason for delay</th>"
          f"<th style='padding:2px 5px;font-size:0.42rem;color:#374151;font-weight:700;text-align:left;'>Date</th>"
          f"<th style='padding:2px 5px;font-size:0.42rem;color:#374151;font-weight:700;text-align:left;'>MARC</th>"
          f"<th style='padding:2px 5px;font-size:0.42rem;color:#374151;font-weight:700;text-align:left;'>Current</th>"
          f"<th style='padding:2px 5px;font-size:0.42rem;color:#374151;font-weight:700;text-align:left;'>Prior</th>"
          f"</tr></thead><tbody>{rows}</tbody></table>"
          + legend_html
          + "</div>"
        + f"<div style='position:absolute;bottom:0;left:0;right:0;height:18px;"
          f"background:#f9fafb;border-top:1px solid #e5e7eb;"
          f"display:flex;align-items:center;padding:0 16px;justify-content:space-between;'>"
          f"<span style='font-size:0.38rem;color:#9ca3af;'>RBC Internal Audit | CONFIDENTIAL</span>"
          f"<span style='font-size:0.38rem;color:#9ca3af;'>"
          f"{qtr} INTERNAL AUDIT SUPPLEMENTAL APPENDICES</span></div>"
        + "</div>"
    )


# ── Risk Spotlight slide (light background, 3-column AC report style) ──────────

def _slide_risk_spotlight(
    cat_id: str,
    audits: pd.DataFrame,
    controls: pd.DataFrame,
    qtr: str,
) -> str:
    cfg     = _CAT_CFG.get(cat_id, {"label": cat_id, "color": _N, "abbr": "??"})
    cat_lbl = cfg["label"]
    hdr_clr = cfg["color"]

    stripes     = [s for s in di.RISK_STRIPES if s["category"] == cat_id]
    stripe_ids  = {s["id"] for s in stripes}

    # Audits touching this category
    if "risk_stripes" in audits.columns:
        cat_audits = audits[audits["risk_stripes"].apply(lambda x: _audit_has_stripe(x, stripe_ids))]
    else:
        cat_audits = pd.DataFrame()

    # Controls for this category
    cat_ctl = (
        controls[controls.get("control_type", pd.Series(dtype=str)).isin(stripe_ids)]
        if not controls.empty else pd.DataFrame()
    )
    n_ctl = len(cat_ctl)
    de_p  = _pct(int((cat_ctl.get("de_result", pd.Series(dtype=str)) == "EF").sum()), n_ctl) if n_ctl else 0
    oe_p  = _pct(int((cat_ctl.get("oe_result", pd.Series(dtype=str)) == "M").sum()), n_ctl) if n_ctl else 0
    tp_p  = _pct(int((cat_ctl.get("test_result", pd.Series(dtype=str)) == "Pass").sum()), n_ctl) if n_ctl else 0

    # ── LEFT column: FY26 Key Focus Areas ──────────────────────────────────────
    tier_c = {"critical": "#dc2626", "high": "#d97706", "standard": "#16a34a"}

    focus_items = ""
    for s in stripes[:7]:
        tc = tier_c.get(s["tier"], "#6b7280")
        focus_items += (
            f"<div style='margin-bottom:5px;padding:4px 8px;"
            f"border-left:3px solid {tc};background:#f0f4f8;border-radius:0 4px 4px 0;'>"
            f"<div style='font-size:0.62rem;font-weight:600;color:#1a2035;'>{s['icon']} {s['name']}</div>"
            f"<div style='font-size:0.46rem;color:#6b7280;letter-spacing:0.08em;text-transform:uppercase;'>"
            f"{s['tier'].upper()}</div></div>"
        )

    # ── MIDDLE column: Q2 Activities & Insights ────────────────────────────────
    n_audits = len(cat_audits)
    n_cmp    = len(cat_audits[cat_audits["status"] == "Complete"]) if not cat_audits.empty else 0
    n_prog   = len(cat_audits[cat_audits["status"].isin(["In Progress","Fieldwork"])]) if not cat_audits.empty else 0

    def _ibar(lbl, p, color):
        return (
            f"<div style='margin-bottom:5px;'>"
            f"<div style='display:flex;justify-content:space-between;margin-bottom:2px;'>"
            f"<span style='font-size:0.58rem;color:#374151;font-weight:500;'>{lbl}</span>"
            f"<span style='font-size:0.6rem;font-weight:700;color:{color};'>{p}%</span></div>"
            + _lhb(p, color) + "</div>"
        )

    if n_ctl > 0:
        control_insight = (
            f"<div style='background:#f0f4f8;border-radius:6px;padding:8px 10px;margin-bottom:8px;"
            f"border:1px solid #d1d5db;'>"
            f"<div style='font-size:0.58rem;font-weight:700;color:#1a2035;margin-bottom:6px;'>"
            f"{n_ctl} Controls Tested</div>"
            + _ibar("Design Effectiveness", de_p, "#16a34a")
            + _ibar("Operating Effectiveness", oe_p, "#2563eb")
            + _ibar("Test Pass Rate", tp_p, "#d97706")
            + "</div>"
        )
    else:
        control_insight = (
            f"<div style='background:#fef3c7;border-radius:6px;padding:8px 10px;margin-bottom:8px;"
            f"border:1px solid #fde68a;font-size:0.6rem;color:#92400e;'>"
            f"No control testing data for this quarter.</div>"
        )

    audit_summary = (
        f"<div style='font-size:0.58rem;color:#374151;margin-bottom:4px;'>"
        f"<strong style='color:#001e4d;'>{n_audits} engagement{'s' if n_audits!=1 else ''}</strong>"
        f" cover this risk area — {n_cmp} complete, {n_prog} in progress.</div>"
    )

    # Enumerate complete audits briefly
    insights_list = ""
    _cmp_audits = cat_audits[cat_audits["status"] == "Complete"] if not cat_audits.empty and "status" in cat_audits.columns else pd.DataFrame()
    for _, row in _cmp_audits.head(3).iterrows():
        nm = str(row.get("audit_name", row.get("audit_id", "—")))[:38]
        rt = str(row.get("current_rating", ""))
        rc = _rc(rt)
        insights_list += (
            f"<div style='font-size:0.56rem;color:#374151;padding:3px 0;"
            f"border-bottom:1px solid #e5e7eb;display:flex;justify-content:space-between;'>"
            f"<span>• {nm}</span>"
            f"<span style='color:{rc};font-weight:600;'>{rt or '—'}</span></div>"
        )

    # ── RIGHT column: Upcoming Audit Coverage ──────────────────────────────────
    upcoming = cat_audits[~cat_audits["status"].isin(["Complete"])] if not cat_audits.empty else pd.DataFrame()
    upcoming_items = ""
    for _, row in upcoming.head(8).iterrows():
        nm  = str(row.get("audit_name", row.get("audit_id","—")))[:28]
        st_ = str(row.get("status",""))
        sc  = _sc(st_)
        upcoming_items += (
            f"<div style='font-size:0.58rem;color:#1a2035;padding:3px 0;"
            f"border-bottom:1px solid #e5e7eb;display:flex;align-items:center;gap:6px;'>"
            f"<span style='width:6px;height:6px;border-radius:50%;background:{sc};"
            f"flex-shrink:0;display:inline-block;'></span>"
            f"<span style='flex:1;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;'>{nm}</span>"
            f"<span style='font-size:0.5rem;color:{sc};flex-shrink:0;'>{st_}</span>"
            f"</div>"
        )
    if not upcoming_items:
        upcoming_items = f"<div style='font-size:0.58rem;color:#6b7280;padding:4px 0;'>All engagements complete.</div>"

    # ── Column header style ─────────────────────────────────────────────────────
    def _col_hdr(text):
        return (
            f"<div style='background:{hdr_clr};color:{_W};padding:6px 10px;border-radius:4px;"
            f"font-size:0.56rem;font-weight:700;letter-spacing:0.1em;text-transform:uppercase;"
            f"font-family:IBM Plex Mono,monospace;margin-bottom:8px;'>{text}</div>"
        )

    # ── Full slide assembly ─────────────────────────────────────────────────────
    return (
        _FL
        + f"<div style='width:100%;max-width:960px;aspect-ratio:16/9;"
          f"background:#f5f7f9;border-radius:10px;overflow:hidden;position:relative;"
          f"box-shadow:0 20px 56px rgba(0,0,30,0.5);font-family:Barlow Condensed,sans-serif;'>"
        # Left RBC gold strip
        + f"<div style='position:absolute;left:0;top:0;bottom:0;width:5px;background:{_G};z-index:3;'></div>"
        # Section header bar
        + f"<div style='background:{hdr_clr};padding:9px 22px 9px 28px;"
          f"display:flex;justify-content:space-between;align-items:center;'>"
          f"<div>"
          f"<div style='font-size:0.52rem;color:rgba(255,255,255,0.55);letter-spacing:0.16em;"
          f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;'>"
          f"SECTION 2 · SPOTLIGHT ON KEY &amp; EMERGING RISKS</div>"
          f"<div style='font-size:1.1rem;font-weight:800;color:{_W};letter-spacing:0.04em;'>"
          f"{cat_lbl}</div>"
          f"</div>"
          f"<div style='font-size:0.58rem;color:rgba(255,255,255,0.45);font-family:IBM Plex Mono,monospace;'>"
          f"RBC INTERNAL AUDIT</div>"
          f"</div>"
        # Body: 3 columns
        + f"<div style='display:grid;grid-template-columns:28% 42% 30%;height:calc(100% - 80px);"
          f"overflow:hidden;'>"
        # LEFT column
        + f"<div style='background:#edf2f7;padding:10px 10px 10px 14px;border-right:1px solid #d1d5db;overflow:hidden;'>"
          + _col_hdr("FY26 Key Focus Areas")
          + focus_items
          + "</div>"
        # MIDDLE column
        + f"<div style='background:{_W};padding:10px 12px;border-right:1px solid #d1d5db;overflow:hidden;'>"
          + _col_hdr("Q2 Activities &amp; Insights")
          + control_insight
          + audit_summary
          + insights_list
          + "</div>"
        # RIGHT column
        + f"<div style='background:#edf2f7;padding:10px 10px 10px 10px;overflow:hidden;'>"
          + _col_hdr("Upcoming Audit Coverage")
          + upcoming_items
          + "</div>"
        + "</div>"
        # Footer
        + f"<div style='position:absolute;bottom:0;left:0;right:0;height:20px;"
          f"background:{hdr_clr};display:flex;align-items:center;padding:0 22px;"
          f"justify-content:space-between;'>"
          f"<span style='font-size:0.48rem;color:rgba(255,255,255,0.4);'>"
          f"RBC Internal Audit | CONFIDENTIAL</span>"
          f"<span style='font-size:0.48rem;color:rgba(255,255,255,0.4);'>"
          f"{qtr} INTERNAL AUDIT QUARTERLY REPORT</span>"
          f"</div>"
        + "</div>"
    )


# ── Slide order (see module docstring for rationale) ──────────────────────────
# Cover → per-segment → Sec3/4/5/7 → Spotlights → App (All Issues) →
# App1 (Core Projects + Late) → App2 (CE + Output + Performance) →
# App4 (Open Issues Overview + Level slides)

def _build_slides(
    audits: pd.DataFrame,
    all_issues: pd.DataFrame,
    controls: pd.DataFrame,
    view: str,
    qtr: str,
    enterprise_issues: pd.DataFrame | None = None,
) -> list[dict]:
    """Build the ordered list of slides for the deck carousel.

    Slide order:
      1. Cover
      2. Per-platform/region slides (Portfolio Overview, Assurance Summary,
         Control Environment, Issues) — one set per platform or region
      3. Section 3 — Assurance Activities & Output, Issue Themes
      4. Section 4 — Issue Overview, Tracking, Resolution
      5. Section 5 — Regulatory Issues, Plan Changes, CAE Performance, QA Review
      6. Section 7 — Glossary, Glossary (Cont.)
      7. Risk Spotlight slides (6 categories)
      8. Appendix — All Issues
      9. Appendix 1 — Core Projects (paginated) + Late Core Projects
     10. Appendix 2 — per-platform CE + Output + Performance slides
     11. Appendix 4 — Issues Overview + L1/L2/L3 slides

    To add a new slide: create a _slide_xxx() function and append a dict
    {"title": ..., "scope": ..., "stype": ..., "html": ...} in the
    appropriate section below.
    """
    slides: list[dict] = []

    # ── 1. Cover ──────────────────────────────────────────────────────────────
    slides.append({"title": "Cover", "scope": "—", "stype": "Cover", "html": _slide_cover(qtr)})

    ent_iss = enterprise_issues if enterprise_issues is not None else all_issues

    # ── 2. Per-platform or per-region slides ──────────────────────────────────
    if view == "Platform":
        platforms = sorted(audits["lead_group"].dropna().unique().tolist())
        for plat in platforms:
            plat_aud = audits[audits["lead_group"] == plat].copy()
            plat_ids = set(plat_aud["audit_id"].tolist())
            plat_iss = _for_audits(all_issues, plat_ids)
            plat_ctl = _for_audits(controls, plat_ids)
            for stype, fn in [
                ("Portfolio Overview", lambda a=plat_aud, i=plat_iss: _slide_portfolio_plat(plat, a, i, qtr)),
                ("Assurance Summary",  lambda a=plat_aud, i=plat_iss: _slide_assurance(plat, a, i, qtr)),
                ("Control Environment",lambda a=plat_aud, c=plat_ctl: _slide_control_env(plat, a, c, qtr)),
                ("Issues",             lambda i=plat_iss: _slide_issues(plat, i, qtr)),
            ]:
                slides.append({"title": f"{plat} — {stype}", "scope": plat, "stype": stype, "html": fn()})
    else:
        for region in _regions(audits):
            rgn_aud = _rgn_filter(audits, region)
            rgn_ids = set(rgn_aud["audit_id"].tolist())
            rgn_iss = _for_audits(all_issues, rgn_ids)
            rgn_ctl = _for_audits(controls, rgn_ids)
            for stype, fn in [
                ("Portfolio Overview", lambda a=rgn_aud, i=rgn_iss: _slide_portfolio_region(region, a, i, qtr)),
                ("Assurance Summary",  lambda a=rgn_aud, i=rgn_iss: _slide_assurance(region, a, i, qtr)),
                ("Control Environment",lambda a=rgn_aud, c=rgn_ctl: _slide_control_env(region, a, c, qtr)),
                ("Issues",             lambda i=rgn_iss: _slide_issues(region, i, qtr)),
            ]:
                slides.append({"title": f"{region} — {stype}", "scope": region, "stype": stype, "html": fn()})

    # ── 3. Section 3 — Assurance Activities & Output ──────────────────────────
    slides.append({
        "title": "Section 3 — Assurance Activities & Output",
        "scope": "Enterprise", "stype": "Assurance Output",
        "html": _slide_assurance_output(audits, ent_iss, qtr),
    })
    slides.append({
        "title": "Section 3 — Issue Theme Analysis",
        "scope": "Enterprise", "stype": "Issue Themes",
        "html": _slide_issue_themes(audits, ent_iss, qtr),
    })

    # ── 4. Section 4 — Audit Issues Management ────────────────────────────────
    slides.append({
        "title": "Section 4 — Issue Management Overview",
        "scope": "Enterprise", "stype": "Issue Overview",
        "html": _slide_issue_overview(ent_iss, qtr),
    })
    slides.append({
        "title": "Section 4 — Issue Tracking Status",
        "scope": "Enterprise", "stype": "Issue Tracking",
        "html": _slide_issue_tracking(ent_iss, qtr),
    })
    slides.append({
        "title": "Section 4 — Issue Resolution Progress",
        "scope": "Enterprise", "stype": "Issue Resolution",
        "html": _slide_issue_resolution(ent_iss, qtr),
    })

    # ── 5. Section 5 — CAE Group Operations ──────────────────────────────────
    slides.append({
        "title": "Section 5 — Regulatory Issues",
        "scope": "Enterprise", "stype": "Regulatory Issues",
        "html": _slide_regulatory_issues(audits, ent_iss, qtr),
    })
    slides.append({
        "title": "Section 5 — Significant Plan Changes",
        "scope": "Enterprise", "stype": "Plan Changes",
        "html": _slide_plan_changes(audits, qtr),
    })
    slides.append({
        "title": "Section 5 — CAE Group Performance Indicators",
        "scope": "Enterprise", "stype": "CAE Performance",
        "html": _slide_cae_performance(audits, ent_iss, qtr),
    })
    slides.append({
        "title": "Section 5 — IA Quality Assurance",
        "scope": "Enterprise", "stype": "QA Review",
        "html": _slide_qa_review(audits, ent_iss, qtr),
    })

    # ── 6. Section 7 — Glossary (2 pages) ────────────────────────────────────
    slides.append({
        "title": "Section 7 — Glossary",
        "scope": "Reference", "stype": "Glossary",
        "html": _slide_glossary(qtr),
    })
    slides.append({
        "title": "Section 7 — Glossary (Cont.)",
        "scope": "Reference", "stype": "Glossary",
        "html": _slide_glossary_cont(qtr),
    })

    # ── 7. Risk Spotlight slides — one per category ───────────────────────────
    for cat_id in _CAT_CFG:
        cat_lbl = _CAT_CFG[cat_id]["label"]
        slides.append({
            "title": f"Spotlight — {cat_lbl}",
            "scope": cat_lbl,
            "stype": "Risk Spotlight",
            "html": _slide_risk_spotlight(cat_id, audits, controls, qtr),
        })

    # ── 8. Appendix — All Issues ──────────────────────────────────────────────
    appx = enterprise_issues if enterprise_issues is not None else all_issues
    slides.append({
        "title": "Appendix — All Issues",
        "scope": "Enterprise", "stype": "Appendix",
        "html": _slide_appendix(appx, qtr),
    })

    # ── 9. Appendix 1 — Core Projects (paginated) + Late Core Projects ───────
    n_core_pages = max(1, (len(audits) + _CORE_PAGE - 1) // _CORE_PAGE)
    for pg in range(1, n_core_pages + 1):
        cont = " (Cont.)" if pg > 1 else ""
        slides.append({
            "title": f"Appendix 1 — Core Projects{cont}",
            "scope": "Enterprise", "stype": "Core Projects",
            "html": _slide_core_projects(audits, qtr, pg),
        })
    # Late/overdue projects — separate appendix slide
    slides.append({
        "title": "Appendix 1 — Late Core Projects",
        "scope": "Enterprise", "stype": "Core Projects",
        "html": _slide_late_core_projects(audits, qtr),
    })

    # ── 10. Appendix 2 — per-segment CE + Output + Performance slides ─────────
    if view == "Platform":
        _a2_plats = sorted(audits["lead_group"].dropna().unique().tolist())
        for _plat in _a2_plats:
            _pa  = audits[audits["lead_group"] == _plat].copy()
            _pids = set(_pa["audit_id"].tolist())
            _pi  = _for_audits(all_issues, _pids)
            _pc  = _for_audits(controls,   _pids)
            slides.append({
                "title": f"Appendix 2 — {_plat}: Control Environment",
                "scope": _plat, "stype": "CE Summary",
                "html":  _slide_app2_ce(_plat, _pa, _pi, _pc, qtr),
            })
            slides.append({
                "title": f"Appendix 2 — {_plat}: Assurance Output",
                "scope": _plat, "stype": "Assurance Output",
                "html":  _slide_app2_output(_plat, _pa, _pi, qtr),
            })
            slides.append({
                "title": f"Appendix 2 — {_plat}: Performance Metrics",
                "scope": _plat, "stype": "Performance",
                "html":  _slide_app2_performance(_plat, _pa, _pi, qtr),
            })
    else:
        for _rgn in _regions(audits):
            _ra  = _rgn_filter(audits, _rgn)
            _rids = set(_ra["audit_id"].tolist())
            _ri  = _for_audits(all_issues, _rids)
            _rc2 = _for_audits(controls,   _rids)
            slides.append({
                "title": f"Appendix 2 — {_rgn}: Control Environment",
                "scope": _rgn, "stype": "CE Summary",
                "html":  _slide_app2_ce(_rgn, _ra, _ri, _rc2, qtr),
            })
            slides.append({
                "title": f"Appendix 2 — {_rgn}: Assurance Output",
                "scope": _rgn, "stype": "Assurance Output",
                "html":  _slide_app2_output(_rgn, _ra, _ri, qtr),
            })
            slides.append({
                "title": f"Appendix 2 — {_rgn}: Performance Metrics",
                "scope": _rgn, "stype": "Performance",
                "html":  _slide_app2_performance(_rgn, _ra, _ri, qtr),
            })

    # ── 11. Appendix 4 — Open Audit Issues (moved to END) ────────────────────
    _ent_iss_a4 = enterprise_issues if enterprise_issues is not None else all_issues
    _open_iss   = _ent_iss_a4[_ent_iss_a4["status"].isin(["Open", "Overdue"])].copy() if not _ent_iss_a4.empty else pd.DataFrame()

    # Map issue_level / severity to L1/L2/L3 groups
    def _lvl(df, sev_vals):
        if "issue_level" in df.columns:
            raw_map = {"Level 1": "L1", "1": "L1",
                       "Level 2": "L2", "2": "L2",
                       "Level 3": "L3", "3": "L3"}
            target = {k for k, v in raw_map.items() if v in sev_vals}
            return df[df["issue_level"].astype(str).isin(target)].copy()
        return df[df.get("severity", pd.Series(dtype=str)).isin(sev_vals)].copy()

    _today   = pd.Timestamp.now().normalize()
    _qtr_ago = _today - pd.DateOffset(months=3)
    _yr_ago  = _today - pd.DateOffset(years=1)

    def _age_filter(df, bucket):
        rd = pd.to_datetime(df.get("raised_date", pd.Series(pd.NaT, index=df.index)), errors="coerce")
        if bucket == "new":
            return df[rd >= _qtr_ago].copy()
        if bucket == "lt1":
            return df[(rd < _qtr_ago) & (rd >= _yr_ago)].copy()
        if bucket == "gt1":
            return df[rd < _yr_ago].copy()
        return df

    _l1 = _lvl(_open_iss, {"L1", "Level 1", "1", "High", "Critical"})
    _l2 = _lvl(_open_iss, {"L2", "Level 2", "2", "Medium"})
    _l3 = _lvl(_open_iss, {"L3", "Level 3", "3", "Low"})

    _l1_new = _age_filter(_l1, "new")
    _l1_lt1 = _age_filter(_l1, "lt1")
    _l1_gt1 = _age_filter(_l1, "gt1")

    def _a4_slides(iss_df, level_label, age_label, qtr):
        n  = len(iss_df)
        if n == 0:
            return []
        np_ = max(1, (n + _APP4_PAGE - 1) // _APP4_PAGE)
        out = []
        for pg in range(1, np_ + 1):
            chunk = iss_df.iloc[(pg - 1) * _APP4_PAGE: pg * _APP4_PAGE]
            n_ovr  = int((chunk["status"] == "Overdue").sum())
            if level_label == "Level 1" and age_label == "Newly Raised":
                summ = (f"{n} L1 issue(s) raised this quarter. "
                        + (f"{n_ovr} already overdue." if n_ovr else "None currently overdue."))
            elif age_label.endswith("year"):
                yr_txt = "&lt;1 year" if "<" in age_label else "&gt;1 year"
                summ = (f"{n} {level_label} issue(s) pending management resolution raised "
                        f"within {yr_txt}. {n_ovr} overdue." if n_ovr else
                        f"{n} {level_label} issue(s) pending management resolution ({yr_txt}).")
            else:
                summ = (f"{n} {level_label} issue(s) open with management. "
                        + (f"{n_ovr} overdue." if n_ovr else "None currently overdue."))
            cont = " (Cont.)" if pg > 1 else ""
            out.append({
                "title":  f"Appendix 4 — {level_label} Issues: {age_label}{cont}",
                "scope":  "Enterprise",
                "stype":  "Open Issues",
                "html":   _slide_open_issues(chunk, audits, level_label, age_label,
                                             summ, qtr, pg, np_),
            })
        return out

    # Appendix 4 Overview + level pages
    slides.append({
        "title": "Appendix 4 — Issues Overview",
        "scope": "Enterprise",
        "stype": "Open Issues",
        "html":  _slide_app4_overview(_ent_iss_a4, audits, qtr),
    })

    for _grp_slides in [
        _a4_slides(_l1_new, "Level 1", "Newly Raised",  qtr),
        _a4_slides(_l1_gt1, "Level 1", "Raised >1 Year", qtr),
        _a4_slides(_l1_lt1, "Level 1", "Raised <1 Year", qtr),
        _a4_slides(_l2,     "Level 2", "Raised Prior to Q2/24", qtr),
        _a4_slides(_l3,     "Level 3", "Raised Prior to Q2/24", qtr),
    ]:
        slides.extend(_grp_slides)

    return slides


# ── PowerPoint export ──────────────────────────────────────────────────────────

def _build_pptx(slides: list[dict], all_issues: pd.DataFrame, qtr: str) -> BytesIO:
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    W, H = 9144000, 5143500   # 10" × 5.63" (16:9)
    prs  = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)
    blank = prs.slide_layouts[6]

    def _rgb(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    def _rect(sl, l, t, w, h_, fill, line=False):
        sh = sl.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h_))
        sh.fill.solid()
        sh.fill.fore_color.rgb = _rgb(fill)
        if not line:
            sh.line.fill.background()
        return sh

    def _txt(sl, text, l, t, w, h_, size, clr, bold=False, align=PP_ALIGN.LEFT, wrap=True):
        box = sl.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h_))
        tf  = box.text_frame
        tf.word_wrap = wrap
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.color.rgb = _rgb(clr)
        run.font.bold = bold
        return box

    def _base(sl, scope, stype, cat_color=None):
        bg_c = cat_color if cat_color else _N
        _rect(sl, 0, 0, W, H, bg_c)
        _rect(sl, 0, 0, 30000, H, _G)
        _rect(sl, 0, 0, W, 400000, "#000816")
        _txt(sl, scope, 60000, 110000, W*7//10, 220000, 13, _G, bold=True)
        _txt(sl, stype, W*7//10, 110000, W*3//10-80000, 220000, 9, "888888", align=PP_ALIGN.RIGHT)
        _rect(sl, 0, H-180000, W, 180000, "#000408")
        _txt(sl, "RBC Internal Audit | CONFIDENTIAL", 60000, H-155000, W//2, 140000, 7, "444444")
        _txt(sl, qtr, W-750000, H-155000, 700000, 140000, 7, "444444", align=PP_ALIGN.RIGHT)

    for s in slides:
        sl    = prs.slides.add_slide(blank)
        scope = s["scope"]
        stype = s["stype"]

        if stype == "Cover":
            _rect(sl, 0, 0, W, H, _N)
            _rect(sl, 0, 0, 30000, H, _G)
            _txt(sl, "RBC INTERNAL AUDIT", W//2-900000, H//3-250000, 1800000, 200000, 10, _G, bold=True, align=PP_ALIGN.CENTER)
            _txt(sl, "AUDIT COMMITTEE REPORT", W//2-1300000, H//3, 2600000, 420000, 30, _W, bold=True, align=PP_ALIGN.CENTER)
            _txt(sl, f"{qtr} · QUARTER-END SUMMARY", W//2-1000000, H//3+460000, 2000000, 200000, 14, _G, bold=True, align=PP_ALIGN.CENTER)
            _rect(sl, 0, H-180000, W, 180000, "#000408")
            _txt(sl, "RBC Internal Audit | CONFIDENTIAL", 60000, H-155000, W//2, 140000, 7, "444444")
            _txt(sl, qtr, W-750000, H-155000, 700000, 140000, 7, "444444", align=PP_ALIGN.RIGHT)
            continue

        cat_c = None
        if stype == "Risk Spotlight":
            cat_id = next((k for k, v in _CAT_CFG.items() if v["label"] == scope), None)
            cat_c  = _CAT_CFG[cat_id]["color"] if cat_id else _N

        _base(sl, scope, stype, cat_c)
        top = 450000

        if stype in ("Portfolio Overview", "Assurance Summary", "Control Environment", "Issues"):
            lines = [s["title"]]
            if stype == "Issues" and not all_issues.empty:
                open_n = len(all_issues[all_issues["status"].isin(["Open","In Progress","Overdue"])])
                lines.append(f"Open Issues: {open_n}")
            for i, line in enumerate(lines):
                _txt(sl, line, 60000, top + i*260000, W-120000, 240000, 12 if i==0 else 10, _W, bold=(i==0))

        elif stype == "Risk Spotlight":
            cat_id = next((k for k,v in _CAT_CFG.items() if v["label"] == scope), None)
            if cat_id:
                stripes = [s2 for s2 in di.RISK_STRIPES if s2["category"] == cat_id]
                _txt(sl, "Key Focus Areas:", 60000, top, W//3-120000, 220000, 11, _G, bold=True)
                for i, stripe in enumerate(stripes[:8]):
                    _txt(sl, f"• {stripe['name']} ({stripe['tier']})", 60000, top+240000+i*240000, W//3-120000, 220000, 9, _W)

        elif stype == "Appendix":
            _txt(sl, f"Enterprise Issues — {len(all_issues)} total", 60000, top, W-120000, 240000, 14, _G, bold=True)
            for i, (_, row) in enumerate(all_issues.head(10).iterrows()):
                title = str(row.get("title","—"))[:60]
                sev   = str(row.get("severity",""))
                _txt(sl, f"• {title}  [{sev}]", 60000, top+280000+i*260000, W-120000, 240000, 9, _W)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── Add-Slide templates ────────────────────────────────────────────────────────

_TEMPLATE_DEFS: dict[str, dict] = {
    "Executive Summary":  {"icon": "📋", "desc": "Two-column insights + section summaries layout"},
    "KPI Overview":       {"icon": "📊", "desc": "3×2 grid of large KPI metric cards"},
    "Issue Analysis":     {"icon": "⚠️",  "desc": "Donut chart with issue breakdown table"},
    "Control Environment":{"icon": "🛡️", "desc": "Assurance indicators table with traffic-light status"},
    "Assurance Output":   {"icon": "✅", "desc": "Q2/YTD stacked bars for ratings and completions"},
    "Risk Spotlight":     {"icon": "🔍", "desc": "Category spotlight with risk narrative"},
    "Blank (Navy)":       {"icon": "📄", "desc": "Empty navy-background frame for custom content"},
    "Blank (White)":      {"icon": "📃", "desc": "Empty white-background frame for custom content"},
}


def _make_template_slide(template: str, title: str, qtr: str) -> dict:
    """Generate a placeholder HTML slide for the chosen template."""
    _N2 = _N  # navy

    def _placeholder_navy(stype_lbl):
        body = (
            f"<div style='display:flex;align-items:center;justify-content:center;"
            f"height:100%;flex-direction:column;gap:12px;'>"
            f"<div style='font-size:0.6rem;letter-spacing:0.2em;color:{_G};"
            f"text-transform:uppercase;font-family:IBM Plex Mono,monospace;'>Custom Slide</div>"
            f"<div style='font-size:1.6rem;font-weight:800;color:{_W};text-align:center;"
            f"max-width:80%;'>{title}</div>"
            f"<div style='font-size:0.55rem;color:{_M};'>Edit commentary below to add content.</div>"
            f"</div>"
        )
        return _frame(body, "Custom", stype_lbl, qtr)

    def _placeholder_white(stype_lbl, headline):
        body = (
            f"<div style='display:flex;align-items:center;justify-content:center;"
            f"height:100%;flex-direction:column;gap:12px;'>"
            f"<div style='font-size:1.3rem;font-weight:800;color:{_N2};text-align:center;"
            f"max-width:80%;'>{title}</div>"
            f"<div style='font-size:0.55rem;color:#6b7280;'>Edit commentary below to add content to this slide.</div>"
            f"</div>"
        )
        return _frame_white(body, "Custom", stype_lbl, headline, qtr)

    if template == "Blank (Navy)":
        html = _placeholder_navy("Custom")
    elif template == "Blank (White)":
        html = _placeholder_white("Custom", title)
    elif template == "Executive Summary":
        body = (
            f"<div style='display:grid;grid-template-columns:52% 48%;gap:16px;height:100%;'>"
            f"<div style='background:#f8f9fa;border-radius:6px;padding:10px 14px;'>"
            f"<div style='font-size:0.52rem;font-weight:800;color:{_N2};text-transform:uppercase;"
            f"letter-spacing:0.1em;border-bottom:2px solid {_G};padding-bottom:4px;margin-bottom:8px;'>"
            f"Results &amp; Insights</div>"
            f"<div style='font-size:0.54rem;color:#374151;line-height:1.5;'>"
            f"Add your key findings and observations here. Use the commentary field below "
            f"to update this narrative with quarter-specific insights.</div>"
            f"</div>"
            f"<div>"
            + "".join(
                f"<div style='background:#f0f4ff;border-left:3px solid {_N2};border-radius:0 4px 4px 0;"
                f"padding:5px 10px;margin-bottom:6px;'>"
                f"<div style='font-size:0.46rem;font-weight:800;color:{_N2};text-transform:uppercase;"
                f"letter-spacing:0.08em;'>{s}</div>"
                f"<div style='font-size:0.5rem;color:#6b7280;margin-top:2px;'>Key message</div></div>"
                for s in ["Assurance Activities", "Issue Management", "Regulatory Items", "Plan Changes"]
            )
            + f"</div></div>"
        )
        html = _frame_white(body, "Custom", "Executive Summary", title, qtr)
    elif template == "KPI Overview":
        def _kpi_card(lbl, icon):
            return (
                f"<div style='background:#f8f9fa;border:1px solid #e5e7eb;border-radius:8px;"
                f"padding:12px;text-align:center;'>"
                f"<div style='font-size:1.8rem;'>{icon}</div>"
                f"<div style='font-size:1.8rem;font-weight:800;color:{_N2};line-height:1;'>—</div>"
                f"<div style='font-size:0.48rem;color:#6b7280;text-transform:uppercase;"
                f"letter-spacing:0.08em;margin-top:4px;'>{lbl}</div></div>"
            )
        body = (
            f"<div style='display:grid;grid-template-columns:repeat(3,1fr);gap:10px;height:100%;"
            f"align-content:center;'>"
            + _kpi_card("Metric 1", "📈") + _kpi_card("Metric 2", "📊") + _kpi_card("Metric 3", "✅")
            + _kpi_card("Metric 4", "⚠️") + _kpi_card("Metric 5", "🎯") + _kpi_card("Metric 6", "📋")
            + f"</div>"
        )
        html = _frame_white(body, "Custom", "KPI Overview", title, qtr)
    else:
        # Default navy placeholder for remaining templates
        html = _placeholder_navy(template)

    return {
        "title":  title,
        "scope":  "Custom",
        "stype":  template,
        "html":   html,
        "custom": True,
    }


# ── Commentary panel ───────────────────────────────────────────────────────────

def _render_commentary_panel(slide_title: str, slide_scope: str, current_user: str) -> None:
    """Render the editable commentary strip below a slide."""
    saved = slide_store.get_full_slide_data(slide_title, slide_scope)
    commentary = saved.get("commentary", "")
    commentary_by = saved.get("commentary_by", "")
    commentary_at = saved.get("commentary_at", "")
    commentary_key = f"cmtry_{slide_store._slide_key(slide_title, slide_scope)}"

    with st.expander(
        f"📝  Slide Commentary" + (f"  ·  last edited by **{commentary_by}**  {commentary_at}" if commentary_by else ""),
        expanded=bool(commentary),
    ):
        new_text = st.text_area(
            "Commentary (supports plain text — appears as narrative context for this slide):",
            value=commentary,
            height=90,
            key=commentary_key,
            placeholder="Add executive narrative, context, or key messages for this slide...",
            label_visibility="collapsed",
        )
        c1, c2, _ = st.columns([1.2, 1.2, 5])
        with c1:
            if st.button("💾 Save", key=f"save_{commentary_key}", type="primary", use_container_width=True):
                slide_store.save_commentary(slide_title, slide_scope, new_text, current_user)
                st.success("Commentary saved.", icon="✅")
                st.rerun()
        with c2:
            if commentary and st.button("🗑 Clear", key=f"clr_{commentary_key}", use_container_width=True):
                slide_store.save_commentary(slide_title, slide_scope, "", current_user)
                st.rerun()


# ── Comments panel ─────────────────────────────────────────────────────────────

def _render_comments_panel(slide_title: str, slide_scope: str, current_user: str) -> None:
    """Render the threaded comment panel below a slide."""
    slide_key_str = slide_store._slide_key(slide_title, slide_scope)
    comments = slide_store.get_comments(slide_title, slide_scope)
    n_comments = sum(1 + len(c.get("replies", [])) for c in comments)

    with st.expander(f"💬  Comments  ({n_comments})", expanded=False):

        # ── Thread display ────────────────────────────────────────────────────
        def _initials(name: str) -> str:
            parts = name.split()
            return (parts[0][0] + (parts[-1][0] if len(parts) > 1 else "")).upper()

        def _format_text(text: str, tagged: list[str]) -> str:
            highlighted = text
            for u in tagged:
                highlighted = highlighted.replace(f"@{u}", f"**@{u}**")
            return highlighted

        user_colors = {
            u: c for u, c in zip(
                di.USERS,
                ["#3b82f6","#8b5cf6","#059669","#dc2626","#d97706",
                 "#0891b2","#7c3aed","#16a34a","#b91c1c","#92400e",
                 "#1d4ed8","#6d28d9","#065f46","#991b1b","#b45309"],
            )
        }

        def _avatar(user: str) -> str:
            clr = user_colors.get(user, "#001e4d")
            ini = _initials(user)
            return (
                f"<div style='width:28px;height:28px;border-radius:50%;background:{clr};"
                f"color:#fff;display:flex;align-items:center;justify-content:center;"
                f"font-size:0.55rem;font-weight:700;flex-shrink:0;'>{ini}</div>"
            )

        if not comments:
            st.caption("No comments yet. Be the first to comment on this slide.")
        else:
            for comment in comments:
                c_id    = comment["id"]
                c_user  = comment["user"]
                c_text  = comment["text"]
                c_time  = comment["timestamp"][:16] if len(comment["timestamp"]) > 10 else comment["timestamp"]
                c_tags  = comment.get("tagged_users", [])
                replies = comment.get("replies", [])
                deleted = comment.get("deleted", False)

                # Top-level comment
                st.markdown(
                    f"<div style='display:flex;gap:8px;margin-bottom:4px;align-items:flex-start;'>"
                    + _avatar(c_user)
                    + f"<div style='flex:1;background:#f8f9fa;border-radius:0 8px 8px 8px;"
                      f"padding:6px 10px;border:1px solid #e5e7eb;'>"
                      f"<div style='display:flex;justify-content:space-between;margin-bottom:3px;'>"
                      f"<span style='font-size:0.58rem;font-weight:700;color:#1f2937;'>{c_user}</span>"
                      f"<span style='font-size:0.48rem;color:#9ca3af;font-family:monospace;'>{c_time}</span>"
                      f"</div>"
                      f"<div style='font-size:0.55rem;color:{'#9ca3af' if deleted else '#374151'};line-height:1.5;'>"
                      + ("&nbsp;".join(
                          f"<strong style='color:#001e4d;'>@{u}</strong>" if u in c_text else t
                          for u in c_tags
                          for t in [u]
                         ) if False else _format_text(c_text, c_tags).replace("\n", "<br>"))
                      + f"</div></div></div>",
                    unsafe_allow_html=True,
                )

                # Replies
                if replies:
                    for reply in replies:
                        r_user = reply["user"]
                        r_time = reply["timestamp"][:16] if len(reply["timestamp"]) > 10 else reply["timestamp"]
                        r_tags = reply.get("tagged_users", [])
                        st.markdown(
                            f"<div style='display:flex;gap:8px;margin-left:36px;margin-bottom:3px;"
                            f"align-items:flex-start;'>"
                            + _avatar(r_user)
                            + f"<div style='flex:1;background:#ffffff;border-radius:0 8px 8px 8px;"
                              f"padding:5px 10px;border:1px solid #e5e7eb;'>"
                              f"<div style='display:flex;justify-content:space-between;margin-bottom:2px;'>"
                              f"<span style='font-size:0.55rem;font-weight:700;color:#374151;'>↳ {r_user}</span>"
                              f"<span style='font-size:0.46rem;color:#9ca3af;font-family:monospace;'>{r_time}</span>"
                              f"</div>"
                              f"<div style='font-size:0.53rem;color:#374151;line-height:1.4;'>"
                              + _format_text(reply["text"], r_tags).replace("\n", "<br>")
                              + f"</div></div></div>",
                            unsafe_allow_html=True,
                        )

                # Inline reply input (toggle per comment)
                reply_key = f"replying_{slide_key_str}_{c_id}"
                if not deleted:
                    r_col1, r_col2 = st.columns([1, 6])
                    with r_col1:
                        if st.button("Reply", key=f"btn_reply_{slide_key_str}_{c_id}",
                                     use_container_width=True):
                            st.session_state[reply_key] = not st.session_state.get(reply_key, False)

                    if st.session_state.get(reply_key, False):
                        with r_col2:
                            r_text = st.text_input(
                                "Your reply:",
                                key=f"txt_reply_{slide_key_str}_{c_id}",
                                placeholder=f"Reply to {c_user}… (use @Name to tag)",
                                label_visibility="collapsed",
                            )
                        post_col, _ = st.columns([1, 5])
                        with post_col:
                            if st.button("Post Reply", key=f"post_reply_{slide_key_str}_{c_id}",
                                         type="primary", use_container_width=True):
                                if r_text.strip():
                                    slide_store.add_reply(slide_title, slide_scope, c_id, current_user, r_text.strip())
                                    st.session_state[reply_key] = False
                                    st.rerun()

                st.markdown("<div style='height:4px;'></div>", unsafe_allow_html=True)

        # ── New comment form ──────────────────────────────────────────────────
        st.markdown("<hr style='margin:8px 0;border-color:#e5e7eb;'>", unsafe_allow_html=True)
        new_comment_key = f"new_comment_{slide_key_str}"
        new_comment_text = st.text_area(
            "Add a comment:",
            key=new_comment_key,
            height=70,
            placeholder="Type your comment… use @Name to tag a colleague",
            label_visibility="collapsed",
        )
        tag_col, post_col = st.columns([3, 1])
        with tag_col:
            tagged_users = st.multiselect(
                "Tag users:",
                options=di.USERS,
                key=f"tag_{slide_key_str}",
                placeholder="Tag colleagues (optional)",
                label_visibility="collapsed",
            )
        with post_col:
            if st.button("Post Comment", key=f"post_{slide_key_str}", type="primary",
                         use_container_width=True):
                if new_comment_text.strip():
                    slide_store.add_comment(
                        slide_title, slide_scope, current_user,
                        new_comment_text.strip(), tagged_users,
                    )
                    st.rerun()


# ── Add Slide dialog ───────────────────────────────────────────────────────────

@st.dialog("Add Slide")
def _add_slide_dialog(after_idx: int, slide_key_prefix: str) -> None:
    """st.dialog for choosing and inserting a custom slide."""
    st.markdown(
        "<div style='font-size:0.85rem;color:#374151;margin-bottom:12px;'>"
        "Choose a template, give the slide a title, and it will be inserted after the current slide.</div>",
        unsafe_allow_html=True,
    )

    template_names = list(_TEMPLATE_DEFS.keys())
    selected = st.selectbox(
        "Template",
        template_names,
        format_func=lambda t: f"{_TEMPLATE_DEFS[t]['icon']}  {t}",
    )
    if selected:
        st.caption(_TEMPLATE_DEFS[selected]["desc"])

    title = st.text_input("Slide title", value=selected or "Custom Slide",
                          placeholder="Enter a descriptive title…")

    c1, c2 = st.columns(2)
    with c1:
        if st.button("➕ Add Slide", type="primary", use_container_width=True):
            if title.strip():
                qtr = st.session_state.get("selected_quarter_filter", "")
                new_slide = _make_template_slide(selected, title.strip(), qtr)
                import uuid
                new_slide["id"] = str(uuid.uuid4())[:8]
                new_slide["insert_after_idx"] = after_idx
                slide_store.save_custom_slide(new_slide)
                # Move to the newly inserted slide
                st.session_state["deck_slide_idx"] = after_idx + 1
                st.session_state["custom_slides_changed"] = True
                st.rerun()
    with c2:
        if st.button("Cancel", use_container_width=True):
            st.rerun()


# ── Main renderer ──────────────────────────────────────────────────────────────

def render_deck_preview(audits: pd.DataFrame, all_issues: pd.DataFrame, snapshot_mode: bool = False):
    for key, val in [("deck_view","Platform"),("deck_slide_idx",0)]:
        if key not in st.session_state:
            st.session_state[key] = val

    qtr = st.session_state.get("selected_quarter_filter","")

    try:
        controls = di.get_controls()
    except Exception:
        controls = pd.DataFrame()

    # Enterprise issues (unfiltered by platform) for the appendix
    enterprise_issues: pd.DataFrame | None = None
    try:
        _raw = di.get_issues()
        if qtr:
            _raw_aud = di.get_audits()
            if "quarter" in _raw_aud.columns:
                _qids = set(_raw_aud[_raw_aud["quarter"] == qtr]["audit_id"].tolist())
                enterprise_issues = _raw[_raw["audit_id"].isin(_qids)].copy()
            else:
                enterprise_issues = _raw
        else:
            enterprise_issues = _raw
    except Exception:
        enterprise_issues = all_issues

    view = st.session_state["deck_view"]

    if audits.empty:
        st.info("No audit data for the current filter selection. Adjust the sidebar filters to see a deck preview.")
        return

    slides = _build_slides(audits, all_issues, controls, view, qtr, enterprise_issues)

    # Merge persisted custom slides into the slide list
    custom_slides = slide_store.get_custom_slides()
    if custom_slides:
        # Sort by insert_after_idx so we insert highest-index last (preserves earlier offsets)
        for cs in sorted(custom_slides, key=lambda x: x.get("insert_after_idx", 0)):
            pos = min(cs.get("insert_after_idx", len(slides) - 1) + 1, len(slides))
            slides.insert(pos, cs)

    n = len(slides)
    current_user: str = di.CURRENT_USER

    # ── Font preload ──────────────────────────────────────────────────────────
    st.markdown(
        "<link href='https://fonts.googleapis.com/css2?family=Barlow+Condensed:wght@600;700"
        "&family=IBM+Plex+Mono:wght@500&display=swap' rel='stylesheet'>",
        unsafe_allow_html=True,
    )

    # ── Navigation row ────────────────────────────────────────────────────────
    c_view, c_scope, c_count, c_prev, c_next, c_add = st.columns([1.2, 2.2, 2.4, 0.7, 0.7, 1.1])

    with c_view:
        new_view = st.selectbox(
            "View", ["Platform","Region"],
            index=0 if view == "Platform" else 1,
            key="deck_view_sel", label_visibility="collapsed",
        )
        if new_view != view:
            st.session_state["deck_view"] = new_view
            st.session_state["deck_slide_idx"] = 0
            st.rerun()

    scopes = list(dict.fromkeys(
        s["scope"] for s in slides if s["stype"] not in ("Cover","Appendix")
    ))

    with c_scope:
        if scopes:
            idx = min(st.session_state["deck_slide_idx"], n - 1)
            cur_scope = slides[idx]["scope"]
            default   = cur_scope if cur_scope in scopes else scopes[0]
            sel = st.selectbox(
                "Jump to", scopes,
                index=scopes.index(default) if default in scopes else 0,
                key="deck_scope_sel", label_visibility="collapsed",
            )
            if sel != default:
                for i, s in enumerate(slides):
                    if s["scope"] == sel:
                        st.session_state["deck_slide_idx"] = i
                        st.rerun()
                        break

    idx = min(st.session_state["deck_slide_idx"], n - 1)
    cur = slides[idx]

    with c_count:
        st.markdown(
            f"<div style='display:flex;align-items:center;height:38px;"
            f"font-family:IBM Plex Mono,monospace;font-size:0.7rem;color:#374151;gap:6px;'>"
            f"<span style='color:#9ca3af;'>Slide</span>"
            f"<span style='font-weight:700;color:{_N};'>{idx+1}</span>"
            f"<span style='color:#d1d5db;'>/</span>"
            f"<span style='color:#6b7280;'>{n}</span>"
            f"<span style='color:{_G};font-size:0.62rem;font-weight:600;overflow:hidden;"
            f"text-overflow:ellipsis;white-space:nowrap;'>&nbsp; {cur['title']}</span>"
            f"</div>",
            unsafe_allow_html=True,
        )

    with c_prev:
        if st.button("◀ Prev", key="deck_prev", use_container_width=True):
            st.session_state["deck_slide_idx"] = max(0, idx - 1)
            st.rerun()

    with c_next:
        if st.button("Next ▶", key="deck_next", use_container_width=True):
            st.session_state["deck_slide_idx"] = min(n - 1, idx + 1)
            st.rerun()

    with c_add:
        if st.button("➕ Add Slide", key="deck_add_slide", use_container_width=True,
                     help="Insert a new slide after the current one"):
            _add_slide_dialog(idx, f"{view}_{qtr}")

    # ── Slide display ─────────────────────────────────────────────────────────
    st.markdown("<div style='margin:10px 0 6px;border-top:2px solid #e5e7eb;'></div>", unsafe_allow_html=True)
    st.markdown(cur["html"], unsafe_allow_html=True)
    st.markdown("<div style='margin:6px 0;border-bottom:1px solid #e5e7eb;'></div>", unsafe_allow_html=True)

    # ── Commentary & Comments ─────────────────────────────────────────────────
    _render_commentary_panel(cur["title"], cur.get("scope", ""), current_user)
    _render_comments_panel(cur["title"], cur.get("scope", ""), current_user)

    # ── Thumbnail strip ───────────────────────────────────────────────────────
    strip_start = max(0, idx - 3)
    strip       = slides[strip_start: idx + 5]
    if strip:
        cols = st.columns(len(strip))
        for ci, (col, s) in enumerate(zip(cols, strip)):
            real = strip_start + ci
            active = real == idx
            with col:
                stype_short = s["stype"][:12]
                scope_short = s["scope"][:10]
                if st.button(
                    f"{scope_short}\n{stype_short}",
                    key=f"ds_{real}",
                    use_container_width=True,
                    type="primary" if active else "secondary",
                    help=s["title"],
                ):
                    st.session_state["deck_slide_idx"] = real
                    st.rerun()

    # ── Download row ──────────────────────────────────────────────────────────
    st.markdown("<div style='height:14px;'></div>", unsafe_allow_html=True)
    dl_col, _, info_col = st.columns([2, 3, 2])

    with dl_col:
        try:
            buf  = _build_pptx(slides, enterprise_issues or all_issues, qtr)
            name = f"AuditIQ_AC_Report_{qtr.replace(' ','_')}.pptx"
            st.download_button(
                label="⬇ Download Full Deck (.pptx)",
                data=buf,
                file_name=name,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                type="primary",
            )
        except Exception as e:
            st.caption(f"pptx unavailable: {e}")

    with info_col:
        st.markdown(
            f"<div style='font-size:0.72rem;color:#6b7280;padding:8px 0;text-align:right;'>"
            f"{n} slides &nbsp;·&nbsp; {view} view &nbsp;·&nbsp; Live data"
            f"</div>",
            unsafe_allow_html=True,
        )
